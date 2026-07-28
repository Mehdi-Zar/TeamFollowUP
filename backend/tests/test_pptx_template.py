"""PPTX export template: storage, the admin CRUD endpoints, and the guarantee that
every PPTX export is built on the uploaded template."""
import io

from pptx import Presentation

from tests.conftest import login

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _marked_template(marker: str = "ORG-TEMPLATE") -> bytes:
    """A minimal valid .pptx carrying a recognizable marker in its core properties,
    so we can prove an exported deck was built on it."""
    prs = Presentation()
    prs.core_properties.category = marker
    prs.slides.add_slide(prs.slide_layouts[0])   # a sample slide the strip must drop
    buf = io.BytesIO(); prs.save(buf)
    return buf.getvalue()


# ---- unit: the render-time helpers -------------------------------------------

def test_new_presentation_builds_on_template_and_strips_sample_slides():
    from app import pptxtpl

    pptxtpl.use(_marked_template())
    try:
        prs = pptxtpl.new_presentation()
        assert len(list(prs.slides)) == 0                     # the template's sample slide is dropped
        prs.slides.add_slide(pptxtpl.blank_layout(prs))
        out = io.BytesIO(); prs.save(out); out.seek(0)        # must save cleanly (no duplicate-part clash)
        assert Presentation(out).core_properties.category == "ORG-TEMPLATE"   # export IS the template
    finally:
        pptxtpl.use(None)
    # with no template the default deck carries no marker
    assert pptxtpl.new_presentation().core_properties.category != "ORG-TEMPLATE"


def test_blank_layout_avoids_placeholder_prompts():
    from app import pptxtpl
    prs = pptxtpl.new_presentation()          # default deck
    lay = pptxtpl.blank_layout(prs)
    assert lay.name == "Blank"                # same choice the renderers made historically


def _add_master_branding(prs):
    """Give the master a non-placeholder shape (a logo/footer stand-in) by cloning a
    master placeholder and stripping its <p:ph>, so blank_layout treats the deck as
    branded."""
    from copy import deepcopy
    from pptx.oxml.ns import qn
    spTree = prs.slide_masters[0].element.find(qn("p:cSld")).find(qn("p:spTree"))
    clone = deepcopy(spTree.find(qn("p:sp")))
    nvPr = clone.find(qn("p:nvSpPr")).find(qn("p:nvPr"))
    ph = nvPr.find(qn("p:ph"))
    if ph is not None:
        nvPr.remove(ph)
    spTree.append(clone)


def test_blank_layout_never_hides_the_template_master_branding():
    """Regression: a template's 'blank' layout often sets showMasterSp="0", which hides
    the master's footer/logo. When the master carries branding, blank_layout must pick a
    layout that SHOWS the master, not the emptiest one that hides it."""
    from app import pptxtpl
    from pptx import Presentation

    prs = Presentation()
    _add_master_branding(prs)                 # the master now has a logo/footer to show
    hidden = next(l for l in prs.slide_layouts if (l.name or "").lower() == "blank")
    hidden.element.set("showMasterSp", "0")   # the real-world '1_Blank' trap

    chosen = pptxtpl.blank_layout(prs)
    assert pptxtpl._shows_master(chosen)      # must not be the branding-hiding layout
    assert chosen is not hidden


# ---- admin CRUD --------------------------------------------------------------

def test_admin_pptx_template_crud(client, seeded):
    login(client, "admin@test")
    assert client.get("/api/admin/pptx-template").json()["present"] is False

    tpl = _marked_template()
    r = client.post("/api/admin/pptx-template",
                    files={"file": ("brand.pptx", tpl, _PPTX_MIME)})
    assert r.status_code == 200 and r.json()["present"] is True and r.json()["filename"] == "brand.pptx"

    got = client.get("/api/admin/pptx-template").json()
    assert got["present"] and got["filename"] == "brand.pptx" and got["size"] == len(tpl)

    dl = client.get("/api/admin/pptx-template/download")
    assert dl.status_code == 200 and dl.content[:2] == b"PK"   # a real .pptx (zip)

    # a non-pptx is rejected
    assert client.post("/api/admin/pptx-template",
                       files={"file": ("x.pptx", b"not a pptx", _PPTX_MIME)}).status_code == 400

    assert client.delete("/api/admin/pptx-template").json()["present"] is False
    assert client.get("/api/admin/pptx-template/download").status_code == 404


def test_admin_pptx_template_requires_admin(client, seeded):
    login(client, seeded["sl_a"])            # a squad leader, not admin
    assert client.get("/api/admin/pptx-template").status_code in (403, 404)
    assert client.post("/api/admin/pptx-template",
                       files={"file": ("x.pptx", _marked_template(), _PPTX_MIME)}).status_code in (403, 404)


# ---- end-to-end: an export endpoint actually uses the template ---------------

def test_steerco_pptx_export_uses_the_uploaded_template(client, db, seeded):
    """Uploading a template must make the Steerco PPTX export inherit it (proves the
    export endpoint wires pptxtpl.use)."""
    from app.models import SteercoEntry

    login(client, "admin@test")
    client.put("/api/admin/modules-config", json={"steerco": {"enabled": True}})
    sid = seeded["squad_a"]
    login(client, seeded["sl_a"])
    client.put(f"/api/steerco/squad/{sid}/enabled", json={"enabled": True})
    db.add(SteercoEntry(squad_id=sid, period="2026-07",
                        data={"kpis": [{"label": "Cloud Users", "value": "10"}]}))
    db.commit()

    login(client, "admin@test")
    # baseline (no template): the deck carries no marker
    base = client.get("/api/steerco/document.pptx?period=2026-07")
    assert base.status_code == 200
    assert Presentation(io.BytesIO(base.content)).core_properties.category != "ORG-TEMPLATE"

    # upload the template, re-export: the deck is now built on it
    client.post("/api/admin/pptx-template", files={"file": ("brand.pptx", _marked_template(), _PPTX_MIME)})
    out = client.get("/api/steerco/document.pptx?period=2026-07")
    assert out.status_code == 200
    assert Presentation(io.BytesIO(out.content)).core_properties.category == "ORG-TEMPLATE"
