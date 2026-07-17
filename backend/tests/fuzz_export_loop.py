"""Loop-mode fuzz harness for the report/dashboard exports (manual runner).

Not auto-collected by pytest (filename isn't ``test_*``). Run it directly to
hammer the exporters across many randomized scenarios:

    cd backend
    DISABLE_SCHEDULER=1 FUZZ_ITERS=1000 python -m tests.fuzz_export_loop

Every iteration asserts the invariants that the 'silently dropped squads' bug
violated:
  1) no exception building or rendering (dashboard single+multi, roadmap, html), any lang/viewer;
  2) the PPTX is well-formed OOXML;
  3) the multi-squad deck has EXACTLY 1 summary + min(n, cap) detail slides + (1 notice iff n>cap),
     i.e. no squad is silently dropped and any overflow is announced;
  4) every rendered detail slide corresponds to a distinct selected squad (no squad lost/duplicated).
Content is deliberately hostile: XML metacharacters, emoji, em-dashes, control chars,
empty/whitespace/huge strings, zero/negative/None budgets, invalid status enums.

The durable pytest regressions for the same invariants live in test_report.py
(test_dashboard_pptx_never_silently_drops_squads / _marks_omitted_squads_when_cap_hit).
"""
import io, os, sys, random, zipfile, traceback

os.environ.setdefault("DISABLE_SCHEDULER", "1")

from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool
from lxml import etree

import app.database as database
engine = create_engine("sqlite+pysqlite:///:memory:",
                       connect_args={"check_same_thread": False}, poolclass=StaticPool, future=True)
SessionLocal = sessionmaker(bind=engine, autoflush=False, autocommit=False, future=True)
database.engine = engine
database.SessionLocal = SessionLocal
from app.database import Base
import app.models as M
Base.metadata.create_all(bind=engine)

from app import status as st
import app.report as report
from app.report import build_report_data, render_pptx, render_roadmap_pptx, render_html
from app.security import hash_password

YEAR = st.current_year_quarter()[0]
PW = hash_password("x")

NASTY = [
    "", "   ", "\t\n ", "A" * 2500,
    "R&D <team> \"quotes\" 'apos' & <b>bold</b>",   # XML metacharacters
    "Impact 90% ↑ - coûts ↓  €1 234,56",             # unicode, em-dash, euro
    "🚀 Livraison 🎯 clôturée ✅", "café naïve façade Œuvre",
    "line1\nline2\rline3", "vt\x0bff\x0cbel\x07nul\x00us\x1f",  # control chars
    "Squad - Data.", "＜full-width＞", "‮RTL override‬",
]
RAGS = ["green", "amber", "red", "grey", "weird", None]
STATUSES = ["on_track", "at_risk", "blocked", "done", "???", None]
KINDS = ["success", "alert", "risk", "bogus"]


def txt(rng):
    return rng.choice(NASTY)


def build_scenario(db, rng, n):
    db.rollback()
    Base.metadata.drop_all(bind=engine)
    Base.metadata.create_all(bind=engine)
    admin = M.User(email="a@x", display_name="Adm", role="admin", is_break_glass=True, password_hash=PW)
    db.add(admin); db.flush()
    tribes = []
    for ti in range(rng.randint(1, 4)):
        tr = M.Tribe(name=f"T{ti}", display_order=ti)
        db.add(tr); db.flush(); tribes.append(tr)
    for i in range(n):
        tr = rng.choice(tribes)
        lead = M.User(email=f"l{i}@x", display_name=txt(rng) or f"Lead{i}", role="squad_leader",
                      tribe_id=tr.id, password_hash=PW)
        db.add(lead); db.flush()
        budget_on = rng.random() < 0.5
        s = M.Squad(name=f"Sq{i:04d}", description=txt(rng), tribe_id=tr.id,
                    leader_user_id=lead.id, display_order=i, budget_enabled=budget_on,
                    kpis_enabled=True, squad_type=rng.choice(["product", "transverse"]),
                    products=[txt(rng)] if rng.random() < .5 else [], hardware=[])
        db.add(s); db.flush()
        for oi in range(rng.randint(0, 6)):
            db.add(M.Objective(squad_id=s.id, year=YEAR, title=txt(rng) or f"O{oi}",
                               rag_status=rng.choice(RAGS) or "green", is_active=rng.random() < .9,
                               target_date=None, weight=rng.randint(1, 3)))
        for ji in range(rng.randint(0, 30)):
            db.add(M.RoadmapItem(squad_id=s.id, year=YEAR, quarter=rng.randint(1, 4),
                                 title=txt(rng) or f"J{ji}", status=rng.choice(STATUSES) or "on_track",
                                 release_stage=rng.choice(["EA", "GA", ""]), theme=txt(rng),
                                 owner=txt(rng), display_order=ji, dependencies=txt(rng)))
        for ki in range(rng.randint(0, 5)):
            db.add(M.KeyMessage(squad_id=s.id, year=YEAR, kind=rng.choice(KINDS),
                                text=txt(rng) or "k", display_order=ki, created_by_user_id=lead.id))
        if budget_on:
            pick = lambda: rng.choice([None, 0, -5000, 1234.56, 1_000_000, 999_999_999])
            db.add(M.SquadBudget(squad_id=s.id, year=YEAR, total=pick(), spent=pick(),
                                 forecast=pick(), comment=txt(rng)))
    db.commit()
    return admin


def validate_ooxml(payload):
    zf = zipfile.ZipFile(io.BytesIO(payload))
    for name in zf.namelist():
        if name.endswith(".xml") or name.endswith(".rels"):
            etree.fromstring(zf.read(name))


def check_multi(db, admin, rng):
    lang = rng.choice(["fr", "en"])
    viewer = rng.choice([admin, None])
    cap = rng.randint(2, 25)
    report._MAX_DETAIL_SLIDES = cap
    try:
        data = build_report_data(db, None, YEAR, 7, lang=lang, viewer=viewer)
        flat = [r for blk in data["tribes"] for r in blk["squads"]]
        n = len(flat)
        payload = render_pptx(data)
        assert payload[:2] == b"PK"
        validate_ooxml(payload)
        from pptx import Presentation
        slides = list(Presentation(io.BytesIO(payload)).slides)
        rendered = min(n, cap)
        notice = 1 if n > cap else 0
        assert len(slides) == 1 + rendered + notice, \
            f"slidecount n={n} cap={cap} got={len(slides)} exp={1+rendered+notice}"
        first_lines = [slides[i].shapes[0].text_frame.text.split("\n")[0] for i in range(1, 1 + rendered)]
        expected = [r["name"] for r in flat[:cap]]
        assert sorted(first_lines) == sorted(expected), \
            f"detail mismatch n={n} cap={cap}: got={sorted(first_lines)[:3]} exp={sorted(expected)[:3]}"
        if notice:
            nt = " ".join(sh.text_frame.text for sh in slides[-1].shapes if sh.has_text_frame)
            assert str(n - cap) in nt, f"notice missing count {n-cap}: {nt!r}"
        return n, cap
    finally:
        report._MAX_DETAIL_SLIDES = 300


def check_single(db, admin, rng):
    s = db.query(M.Squad).order_by(M.Squad.id).first()
    if not s:
        return
    data = build_report_data(db, None, YEAR, 7, squad_id=s.id, lang=rng.choice(["fr", "en"]), viewer=admin)
    payload = render_pptx(data)
    validate_ooxml(payload)
    from pptx import Presentation
    assert len(list(Presentation(io.BytesIO(payload)).slides)) == 1  # exactly the one squad


def check_roadmap_and_html(db, admin, rng):
    data = build_report_data(db, None, YEAR, 7, lang=rng.choice(["fr", "en"]), viewer=admin)
    validate_ooxml(render_roadmap_pptx(data))
    render_html(data, standalone=True)


def main():
    iters = int(os.environ.get("FUZZ_ITERS", "250"))
    db = SessionLocal()
    ok = 0
    hist = {}
    for it in range(iters):
        rng = random.Random(it)  # reproducible per-iteration
        n = rng.choice([0, 1, 2, 3, rng.randint(0, 20), rng.randint(20, 45), rng.randint(0, 60)])
        try:
            admin = build_scenario(db, rng, n)
            res = check_multi(db, admin, rng)
            check_single(db, admin, rng)
            check_roadmap_and_html(db, admin, rng)
            ok += 1
            hist[res[0] > res[1]] = hist.get(res[0] > res[1], 0) + 1
        except Exception:
            print(f"\n!!! ITERATION {it} FAILED (n={n}) !!!")
            traceback.print_exc()
            print(f"Reproduce with: FUZZ_ITERS=1 and random.Random({it})")
            db.close()
            sys.exit(1)
        if (it + 1) % 25 == 0:
            print(f"[loop] {it+1}/{iters} scenarios passed "
                  f"(over-cap cases so far: {hist.get(True,0)})")
    db.close()
    print(f"\nALL GREEN: {ok}/{iters} randomized export scenarios upheld every invariant.")
    print(f"  overflow-notice path exercised in {hist.get(True,0)} scenarios; "
          f"under-cap in {hist.get(False,0)}.")


if __name__ == "__main__":
    main()
