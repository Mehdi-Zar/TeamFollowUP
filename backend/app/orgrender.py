"""Org-chart export: a tribe's editable org tree rendered to a printable HTML page
and a single-slide PPTX, both laid out to fit one page. Mirrors the brand palette.

A node is {id, title, person_name, squad_id, squad_status, children:[...]}. Status
colours (derived squad health) are applied when a node maps to a squad.
"""
from __future__ import annotations

import html
import io

from .report import _pptx_toolkit

_NAVY = "#1E2761"
_INK = "#1F2937"
_MUTED = "#6B7280"
_LINE = "#CBD5E1"
_CARD = "#FFFFFF"
_STATUS = {"on_track": "#027A48", "at_risk": "#B54708", "blocked": "#B42318"}
_STATUS_BG = {"on_track": "#ECFDF3", "at_risk": "#FFFAEB", "blocked": "#FEF3F2"}

_T = {
    "fr": {"org": "Organigramme", "no_org": "Aucun élément à afficher."},
    "en": {"org": "Org chart", "no_org": "Nothing to display."},
}


def _lang(lang: str | None) -> str:
    return "en" if lang == "en" else "fr"


# ----- HTML --------------------------------------------------------------------

_ORG_CSS = """<style>
.org-doc{font-family:-apple-system,Segoe UI,Roboto,Helvetica,Arial,sans-serif;color:#1F2937;
  padding:24px;background:#fff}
.org-doc h1{font-size:22px;margin:0 0 2px}
.org-doc .sub{color:#6B7280;font-size:13px;margin-bottom:18px}
.org-tree,.org-tree ul{list-style:none;margin:0;padding:0}
.org-tree,.org-tree ul{display:flex;justify-content:center}
.org-tree li{display:flex;flex-direction:column;align-items:center;position:relative;padding:18px 8px 0}
/* connector up from each child to the horizontal bar */
.org-tree li::before{content:"";position:absolute;top:0;left:50%;width:1px;height:18px;background:#CBD5E1}
/* horizontal bar between siblings */
.org-tree li::after{content:"";position:absolute;top:0;left:0;right:0;height:1px;background:#CBD5E1}
.org-tree li:first-child::after{left:50%}
.org-tree li:last-child::after{right:50%}
.org-tree li:only-child::after{display:none}
.org-tree > li::before{display:none}
.org-node{border:1px solid #CBD5E1;border-radius:10px;background:#fff;padding:8px 12px;min-width:120px;
  max-width:190px;text-align:center;box-shadow:0 1px 2px rgba(16,24,40,.06)}
.org-node .t{font-weight:700;font-size:12.5px;line-height:1.25}
.org-node .p{color:#6B7280;font-size:11px;margin-top:2px}
.org-node.has-status{border-top-width:3px}
.org-children{margin-top:18px}
</style>"""


def _node_html(node: dict, e) -> str:
    status = node.get("squad_status") if node.get("squad_id") else None
    cls = "org-node" + (" has-status" if status else "")
    style = ""
    if status in _STATUS:
        style = f' style="border-top-color:{_STATUS[status]};background:{_STATUS_BG[status]}"'
    person = f'<div class="p">{e(node["person_name"])}</div>' if node.get("person_name") else ""
    box = f'<div class="{cls}"{style}><div class="t">{e(node["title"])}</div>{person}</div>'
    kids = node.get("children") or []
    if not kids:
        return f"<li>{box}</li>"
    inner = "".join(_node_html(k, e) for k in kids)
    return f'<li>{box}<ul class="org-children">{inner}</ul></li>'


def render_org_html(roots: list[dict], scope_name: str, *, lang: str = "fr",
                    standalone: bool = True) -> str:
    e = html.escape
    lang = _lang(lang)
    title = _T[lang]["org"]
    body = [f'<div class="org-doc"><h1>{e(title)} - {e(scope_name)}</h1>',
            f'<div class="sub">{e(scope_name)}</div>']
    if not roots:
        body.append(f'<div class="sub">{e(_T[lang]["no_org"])}</div>')
    else:
        body.append('<ul class="org-tree">')
        body.append("".join(_node_html(r, e) for r in roots))
        body.append('</ul>')
    body.append('</div>')
    inner = _ORG_CSS + "".join(body)
    if not standalone:
        return inner
    return (f'<!doctype html><html lang="{e(lang)}"><head><meta charset="utf-8">'
            f'<title>{e(title)} - {e(scope_name)}</title></head><body>{inner}</body></html>')


# ----- PPTX --------------------------------------------------------------------

def _layout(roots: list[dict]) -> tuple[list[dict], int, int]:
    """Tidy top-down layout. Returns (placed_nodes, n_leaf_slots, max_depth).
    Each placed node gets x (leaf-slot centre, float) and depth (int)."""
    placed: list[dict] = []
    counter = {"leaf": 0}

    def walk(node: dict, depth: int) -> float:
        kids = node.get("children") or []
        if not kids:
            x = counter["leaf"] + 0.5
            counter["leaf"] += 1
        else:
            xs = [walk(k, depth + 1) for k in kids]
            x = sum(xs) / len(xs)
        placed.append({"node": node, "x": x, "depth": depth,
                       "parent_x": None})  # parent_x filled by caller
        node["_x"] = x
        node["_depth"] = depth
        return x

    max_depth = 0

    def depth_of(node, d=0):
        nonlocal max_depth
        max_depth = max(max_depth, d)
        for k in (node.get("children") or []):
            depth_of(k, d + 1)

    for r in roots:
        walk(r, 0)
        depth_of(r, 0)
    return placed, max(1, counter["leaf"]), max_depth


def render_org_pptx(roots: list[dict], scope_name: str, *, lang: str = "fr") -> bytes:
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE = _pptx_toolkit()
    from pptx.enum.shapes import MSO_CONNECTOR

    def rgb(h):
        return RGBColor.from_string(h.lstrip("#").upper())

    lang = _lang(lang)
    SLIDE_W, SLIDE_H = 13.333, 7.5
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # Header band
    head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.9))
    head.fill.solid(); head.fill.fore_color.rgb = rgb(_NAVY); head.line.fill.background(); head.shadow.inherit = False
    tf = head.text_frame; tf.margin_left = Inches(0.5); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; r = p.add_run(); r.text = f'{_T[lang]["org"]} - {scope_name}'
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = rgb("#FFFFFF")

    placed, n_leaf, max_depth = _layout(roots)
    if not placed:
        buf = io.BytesIO(); prs.save(buf); return buf.getvalue()

    top, bottom = 1.2, 7.1
    left_m, right_m = 0.4, 0.4
    avail_w = SLIDE_W - left_m - right_m
    col_w = avail_w / n_leaf                       # width per leaf slot
    row_h = (bottom - top) / (max_depth + 1)
    box_w = min(2.2, max(1.1, col_w * 0.92))
    box_h = min(0.8, max(0.5, row_h * 0.62))

    def cx(x):  # leaf-slot centre → inches
        return left_m + (x / n_leaf) * avail_w
    def cy(depth):
        return top + depth * row_h + (row_h - box_h) / 2

    # Connectors first (so boxes sit on top).
    def draw_connectors(node):
        for k in (node.get("children") or []):
            x1, y1 = cx(node["_x"]), cy(node["_depth"]) + box_h
            x2, y2 = cx(k["_x"]), cy(k["_depth"])
            cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
            cn.line.color.rgb = rgb(_LINE); cn.line.width = Pt(1)
            draw_connectors(k)
    for r0 in roots:
        draw_connectors(r0)

    for item in placed:
        node = item["node"]
        x, depth = cx(item["x"]) - box_w / 2, cy(item["depth"])
        status = node.get("squad_status") if node.get("squad_id") else None
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(depth),
                                 Inches(box_w), Inches(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(_STATUS_BG[status]) if status in _STATUS_BG else rgb(_CARD)
        box.line.color.rgb = rgb(_STATUS[status]) if status in _STATUS else rgb(_LINE)
        box.line.width = Pt(2 if status in _STATUS else 1)
        box.shadow.inherit = False
        try:
            box.adjustments[0] = 0.12
        except Exception:
            pass
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r1 = p.add_run(); r1.text = node["title"]
        r1.font.size = Pt(10.5 if max_depth <= 3 else 9); r1.font.bold = True; r1.font.color.rgb = rgb(_INK)
        if node.get("person_name"):
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run(); r2.text = node["person_name"]
            r2.font.size = Pt(8.5); r2.font.color.rgb = rgb(_MUTED)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
