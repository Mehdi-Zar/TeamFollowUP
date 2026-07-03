"""Weekly report: combined dashboard + progress-review, rendered to HTML and PPTX.

Used both for on-demand downloads/emails (routers/reports.py) and for the
automatic weekly send driven by the in-process scheduler (send_due_weekly_reports).
"""
from __future__ import annotations

import html
import io
import os
from datetime import datetime, timedelta, timezone

from sqlalchemy import select
from sqlalchemy.orm import Session, selectinload

from . import status as st
from .generalconfig import get_general
from .models import Squad, Tribe, utcnow
from .serializers import annual_progress, budget_out, dependency_label


def _budget_for_report(squad, year: int, viewer) -> dict | None:
    """Budget figures for a report, only when the viewer may see this squad's
    budget (admin / its tribe leader / its own squad leader) and it is enabled."""
    if viewer is None or not squad.budget_enabled:
        return None
    from .deps import is_squad_privileged
    if not is_squad_privileged(viewer, squad):
        return None
    b = budget_out(squad, year)
    return {
        "total": b.total, "spent": b.spent, "forecast": b.forecast,
        "status": b.status, "spent_pct": b.spent_pct, "forecast_pct": b.forecast_pct,
        "overrun": b.overrun, "overrun_pct": b.overrun_pct, "comment": b.comment,
    }


def _aware(dt: datetime | None) -> datetime | None:
    if dt is None:
        return None
    return dt if dt.tzinfo else dt.replace(tzinfo=timezone.utc)

# ----- RAG / status presentation -------------------------------------------------

RAG_COLOR = {"red": "#dc2626", "amber": "#d97706", "green": "#16a34a", "grey": "#6b7280"}

# Roadmap milestone palette (mirrors the reference "Global Roadmap" slide):
# navy ink for titles/themes, gold for EA, green for GA.
RM_INK = "#002060"
STAGE_COLOR = {"EA": "#FFC000", "GA": "#00B050"}


def group_by_theme(items: list[dict]) -> list[tuple[str, list[dict]]]:
    """Group roadmap items by their theme, preserving first-seen order.

    A blank/missing theme yields an empty key (rendered without a header)."""
    groups: list[tuple[str, list[dict]]] = []
    index: dict[str, int] = {}
    for it in items:
        key = (it.get("theme") or "").strip()
        if key not in index:
            index[key] = len(groups)
            groups.append((key, []))
        groups[index[key]][1].append(it)
    return groups

_STATUS_RAG = {"blocked": "red", "at_risk": "amber", "on_track": "green",
               "done": "green", "red": "red", "amber": "amber", "green": "green"}

_STATUS_LABELS = {
    "fr": {"blocked": "Bloqué", "at_risk": "À risque", "on_track": "En cours",
           "done": "Terminé", "red": "Rouge", "amber": "Orange", "green": "Vert"},
    "en": {"blocked": "Blocked", "at_risk": "At risk", "on_track": "On track",
           "done": "Done", "red": "Red", "amber": "Amber", "green": "Green"},
}

_CHANGE_LABELS = {
    "fr": {"jalon_added": "Nouveau jalon", "jalon_status": "Jalon",
           "quarter_pct": "Progression", "objective_rag": "Objectif", "kpi_trend": "KPI"},
    "en": {"jalon_added": "New milestone", "jalon_status": "Milestone",
           "quarter_pct": "Progress", "objective_rag": "Objective", "kpi_trend": "KPI"},
}

# All other report strings, by language.
_RT = {
    "fr": {
        "report": "Rapport hebdomadaire", "all_tribes": "Toutes les tribus",
        "year": "Année", "generated": "généré le", "window": "fenêtre {n} j",
        "generated_full": "Généré le {d}", "window_full": "Fenêtre {n} jours",
        "k_squads": "Squads", "k_progress": "Progression moy.", "k_blocked": "Jalons bloqués",
        "k_atrisk": "Jalons à risque", "k_obj_red": "Objectifs rouges", "k_stale": "Reporting périmé",
        "attention": "Points d'attention", "blocked_n": "bloqué(s)", "stale": "périmé",
        "h_squad": "Squad", "h_leader": "Responsable", "h_status": "Statut", "h_progress": "Progr.",
        "h_progress_long": "Progression", "h_delta": "Δ sem.", "h_blocked": "Bloqués",
        "h_atrisk": "À risque", "h_facts": "Faits de la semaine",
        "squad_scope": "Squad {name}", "more_squads": "… +{n} autres squads",
        "subject": "Rapport hebdomadaire — {scope} — semaine {w}", "synthesis": "Synthèse",
        "subject_personal": "Rapport — {scope} — {n} j",
        "detail_title": "Détail par squad", "h_objectives": "OTD",
        "h_roadmap": "Roadmap & jalons", "deadline": "échéance", "no_obj": "Aucun objectif",
        "no_jalon": "Aucun jalon", "dep": "Dép.", "roadmap_report": "Roadmap",
        "roadmap_subject": "Roadmap {year}",
        "stage_ea": "Accès anticipé", "stage_ga": "Disponibilité générale",
        "transverse_report": "Reporting transverse", "h_initiatives": "Initiatives",
        "no_initiative": "Aucune initiative", "h_jalons": "jalons", "h_objective": "Objectif",
        "h_jalon": "Jalon", "h_stage": "Phase", "h_otd": "OTD (engagements de livraison)",
        "no_otd": "Aucun OTD", "otd_commit": "engagé",
        "h_otd_section": "OTD", "h_freshness_ok": "Données à jour",
        "h_key_messages": "Messages clés", "no_key_message": "Aucun message clé",
        "km_success": "Succès", "km_alert": "Alerte", "km_risk": "Risque",
        "h_budget": "Budget", "no_budget": "Budget non renseigné",
        "b_total": "Total", "b_spent": "Consommé", "b_forecast": "Prévision",
        "b_on_track": "Sur les rails", "b_at_risk": "À risque", "b_over": "Dépassement",
        "leaves_upcoming": "Absences à venir (30 j)", "leaves_pending": "à valider", "days_short": "j",
        # --- "What's new since your last report" changelog ---
        "whatsnew": "Nouveautés depuis votre dernier rapport",
        "first_report": "Premier rapport — pas encore de comparaison.",
        "no_changes": "Aucun changement depuis le dernier rapport.",
        "subj_changes": "[{n} nouveauté(s)]", "subj_uptodate": "[à jour]",
        "chg_progress": "avancement {d} pts", "chg_status": "statut {frm} → {to}",
        "chg_ms_new": "nouveau jalon « {title} »", "chg_ms_status": "jalon « {title} » : {frm} → {to}",
        "chg_ms_removed": "jalon « {title} » retiré",
        "chg_obj_new": "nouvel OTD « {title} »", "chg_obj_status": "OTD « {title} » : {frm} → {to}",
        "chg_budget": "budget mis à jour", "chg_km": "{n} nouveau(x) message(s) clé(s)",
        "chg_stale": "reporting devenu périmé", "chg_unstale": "reporting de nouveau à jour",
        "chg_new_squad": "nouvelle squad « {name} »", "chg_squad_removed": "squad « {name} » retirée",
        "sum_moved": "{n} squad(s) ont bougé", "sum_delivered": "{n} jalon(s) livré(s)",
        "sum_blocked": "{n} nouveau(x) bloqueur(s)", "sum_stale": "{n} squad(s) périmée(s)",
        "ms_on_track": "En cours", "ms_at_risk": "À risque", "ms_blocked": "Bloqué", "ms_done": "Livré",
        "rag_green": "vert", "rag_amber": "orange", "rag_red": "rouge",
    },
    "en": {
        "report": "Weekly report", "all_tribes": "All tribes",
        "year": "Year", "generated": "generated on", "window": "window {n}d",
        "generated_full": "Generated on {d}", "window_full": "Window {n} days",
        "k_squads": "Squads", "k_progress": "Avg. progress", "k_blocked": "Blocked milestones",
        "k_atrisk": "At-risk milestones", "k_obj_red": "Red objectives", "k_stale": "Stale reporting",
        "attention": "Attention points", "blocked_n": "blocked", "stale": "stale",
        "h_squad": "Squad", "h_leader": "Leader", "h_status": "Status", "h_progress": "Progr.",
        "h_progress_long": "Progress", "h_delta": "Δ wk", "h_blocked": "Blocked",
        "h_atrisk": "At risk", "h_facts": "This week",
        "squad_scope": "Squad {name}", "more_squads": "… +{n} more squads",
        "subject": "Weekly report — {scope} — week {w}", "synthesis": "Summary",
        "subject_personal": "Report — {scope} — {n}d",
        "detail_title": "Detail by squad", "h_objectives": "OTD",
        "h_roadmap": "Roadmap & milestones", "deadline": "due", "no_obj": "No objective",
        "no_jalon": "No milestone", "dep": "Dep.", "roadmap_report": "Roadmap",
        "roadmap_subject": "Roadmap {year}",
        "stage_ea": "Early Access", "stage_ga": "General Availability",
        "transverse_report": "Transverse report", "h_initiatives": "Initiatives",
        "no_initiative": "No initiative", "h_jalons": "milestones", "h_objective": "Objective",
        "h_jalon": "Milestone", "h_stage": "Stage", "h_otd": "OTD (delivery commitments)",
        "no_otd": "No OTD", "otd_commit": "committed",
        "h_otd_section": "OTD", "h_freshness_ok": "Up to date",
        "h_key_messages": "Key messages", "no_key_message": "No key message",
        "km_success": "Success", "km_alert": "Alert", "km_risk": "Risk",
        "h_budget": "Budget", "no_budget": "Budget not set",
        "b_total": "Total", "b_spent": "Spent", "b_forecast": "Forecast",
        "b_on_track": "On track", "b_at_risk": "At risk", "b_over": "Over budget",
        "leaves_upcoming": "Upcoming absences (30 d)", "leaves_pending": "to approve", "days_short": "d",
        # --- "What's new since your last report" changelog ---
        "whatsnew": "What's new since your last report",
        "first_report": "First report — nothing to compare yet.",
        "no_changes": "No changes since the last report.",
        "subj_changes": "[{n} update(s)]", "subj_uptodate": "[up to date]",
        "chg_progress": "progress {d} pts", "chg_status": "status {frm} → {to}",
        "chg_ms_new": "new milestone “{title}”", "chg_ms_status": "milestone “{title}”: {frm} → {to}",
        "chg_ms_removed": "milestone “{title}” removed",
        "chg_obj_new": "new OTD “{title}”", "chg_obj_status": "OTD “{title}”: {frm} → {to}",
        "chg_budget": "budget updated", "chg_km": "{n} new key message(s)",
        "chg_stale": "reporting went stale", "chg_unstale": "reporting back up to date",
        "chg_new_squad": "new squad “{name}”", "chg_squad_removed": "squad “{name}” removed",
        "sum_moved": "{n} squad(s) moved", "sum_delivered": "{n} milestone(s) delivered",
        "sum_blocked": "{n} new blocker(s)", "sum_stale": "{n} squad(s) went stale",
        "ms_on_track": "On track", "ms_at_risk": "At risk", "ms_blocked": "Blocked", "ms_done": "Done",
        "rag_green": "green", "rag_amber": "amber", "rag_red": "red",
    },
}


def _lang(lang: str | None) -> str:
    return "en" if lang == "en" else "fr"


# Leave types are stored as French labels; in English show "English (French)".
_LEAVE_TYPE_EN = {
    "Congés payés": "Paid leave", "RTT": "RTT", "Maladie": "Sick leave",
    "Formation": "Training", "Autre": "Other",
}


def leave_type_label(label: str, lang: str) -> str:
    if _lang(lang) != "en":
        return label
    en = _LEAVE_TYPE_EN.get(label)
    return f"{en} ({label})" if en and en != label else label


def rt(lang: str, key: str, **kw) -> str:
    s = _RT[_lang(lang)].get(key, key)
    return s.format(**kw) if kw else s


def _status_rag(status: str | None) -> str:
    return _STATUS_RAG.get(status or "", "grey")


def _status_label(status: str | None, lang: str = "fr") -> str:
    return _STATUS_LABELS[_lang(lang)].get(status or "", status or "-")


def _change_text(ch: dict, lang: str = "fr") -> str:
    kind = ch.get("kind", "")
    label = ch.get("label", "")
    frm, to = ch.get("from"), ch.get("to")
    prefix = _CHANGE_LABELS[_lang(lang)].get(kind, kind)
    sep = " : " if _lang(lang) == "fr" else ": "
    if kind == "jalon_added":
        return f"{prefix}{sep}{label}"
    if to is not None and frm is not None:
        return f"{prefix} {label}{sep}{_status_label(str(frm), lang)} → {_status_label(str(to), lang)}"
    if to is not None:
        return f"{prefix} {label} → {_status_label(str(to), lang)}"
    return f"{prefix} {label}".strip()


# ----- Data assembly --------------------------------------------------------------

def build_report_data(db: Session, scope_tribe: int | None, year: int | None = None,
                      since_days: int = 7, now: datetime | None = None,
                      squad_id: int | None = None, lang: str | None = None,
                      squad_ids: list[int] | None = None, viewer=None) -> dict:
    """Assemble the combined dashboard + weekly-review data for the given scope.

    squad_id, when set, narrows the report to a single squad (ignoring scope_tribe).
    squad_ids, when set, restricts the report to that subset of squads (still within
    the caller's tribe scope) - used to pick which squads appear in a global roadmap.
    lang, when set, picks the report language; otherwise the general default_lang.
    """
    now = now or utcnow()
    cfg = get_general(db)
    threshold = cfg.get("staleness_threshold_days")
    year = year or st.current_year_quarter(now)[0]
    lang = _lang(lang or cfg.get("default_lang"))

    tribes = {t.id: t for t in db.scalars(select(Tribe)).all()}
    q = select(Squad).order_by(Squad.display_order, Squad.id).options(
        selectinload(Squad.objectives), selectinload(Squad.roadmap_items),
        selectinload(Squad.quarter_progress), selectinload(Squad.kpis),
        selectinload(Squad.snapshots), selectinload(Squad.leader),
        selectinload(Squad.budgets), selectinload(Squad.key_messages),
    )
    id_filter = set(squad_ids) if squad_ids else None
    squads = []
    for s in db.scalars(q).all():
        if squad_id is not None:
            if s.id == squad_id:
                squads.append(s)
            continue
        if scope_tribe is not None and s.tribe_id != scope_tribe:
            continue  # tribe scope is the security boundary
        if id_filter is not None and s.id not in id_filter:
            continue  # caller-chosen subset
        squads.append(s)

    # Initiatives assigned to each squad (shown in that squad's report/dashboard).
    from .models import Initiative
    init_by_squad: dict[int, list[dict]] = {}
    sq_ids = [s.id for s in squads]
    if sq_ids:
        irows = db.scalars(
            select(Initiative).where(Initiative.year == year, Initiative.squad_id.in_(sq_ids))
            .order_by(Initiative.display_order, Initiative.id)).all()
        for it in irows:
            init_by_squad.setdefault(it.squad_id, []).append({
                "title": it.title, "owner": it.owner,
                "deadline": it.deadline.date().isoformat() if it.deadline else None})

    by_tribe: dict[int | None, list[dict]] = {}
    totals = {"squads": 0, "blocked": 0, "at_risk": 0, "objectives_red": 0,
              "stale": 0, "progress_sum": 0}

    for s in squads:
        c = st.counts(s, year)
        f = st.freshness(s, threshold, now)
        prog = st.year_progress(s, year)
        comments = st.quarter_comments(s, year)
        ann = annual_progress(s, year)
        # Full per-squad content (objectives + roadmap by quarter + advancement),
        # so the report/PPTX can show everything, not just the dashboard summary.
        detail = {
            "initiatives": init_by_squad.get(s.id, []),
            "objectives": [
                {"title": o.title,
                 "rag": st.objective_status(o, s, now),
                 "target_date": o.target_date.date().isoformat() if o.target_date else None}
                for o in sorted(s.objectives, key=lambda x: x.id)
                if o.year == year and o.is_active
            ],
            "quarters": [
                {"q": q, "pct": prog[q], "comment": comments.get(q),
                 "items": [
                     {"title": r.title, "status": r.status, "owner": r.owner,
                      "stage": r.release_stage, "theme": r.theme,
                      "dependency": dependency_label(r)}
                     for r in sorted(s.roadmap_items, key=lambda x: (x.display_order, x.id))
                     if r.year == year and r.quarter == q
                 ]}
                for q in (1, 2, 3, 4)
            ],
            # Hand-curated key messages (success / alert / risk) for this squad/year.
            "key_messages": [
                {"kind": m.kind, "text": m.text,
                 "created_at": _aware(m.created_at).strftime("%Y-%m-%d %H:%M") if m.created_at else None}
                for m in sorted(s.key_messages, key=lambda x: (x.display_order, x.id))
                if m.year == year
            ],
            # Budget readout - only for a viewer allowed to see this squad's figures.
            "budget": _budget_for_report(s, year, viewer),
        }
        row = {
            "squad_id": s.id,
            "name": s.name,
            "leader": s.leader.display_name if s.leader else "",
            "status": st.squad_status(s, year),
            "status_rag": _status_rag(st.squad_status(s, year)),
            "quarters": {q: prog[q] for q in (1, 2, 3, 4)},
            "annual_pct": ann,
            "blocked": c["roadmap_blocked"],
            "at_risk": c["roadmap_at_risk"],
            "objectives_red": c["objectives_red"],
            "age_days": f.get("age_days"),
            "is_stale": bool(f.get("is_stale")),
            "delta": 0,
            "confidence": None,
            "note": None,
            "points_in_period": 0,
            "changes": [],
            "detail": detail,
        }
        by_tribe.setdefault(s.tribe_id, []).append(row)
        totals["squads"] += 1
        totals["blocked"] += row["blocked"]
        totals["at_risk"] += row["at_risk"]
        totals["objectives_red"] += row["objectives_red"]
        totals["stale"] += 1 if row["is_stale"] else 0
        totals["progress_sum"] += ann

    # Order squads within a tribe: most blocked / worst movers first.
    for rows in by_tribe.values():
        rows.sort(key=lambda r: (-r["blocked"], r["delta"], r["name"]))

    tribe_blocks = []
    for tid, rows in sorted(by_tribe.items(),
                            key=lambda kv: (tribes[kv[0]].display_order if kv[0] in tribes else 0,
                                            tribes[kv[0]].name if kv[0] in tribes else "")):
        tribe_blocks.append({
            "tribe_id": tid,
            "tribe_name": tribes[tid].name if tid in tribes else "-",
            "squads": rows,
        })

    # Attention list: blocked or regressing squads, across the whole scope.
    attention = [r for blk in tribe_blocks for r in blk["squads"]
                 if r["blocked"] > 0 or r["delta"] < 0 or r["is_stale"]]
    attention.sort(key=lambda r: (-r["blocked"], r["delta"]))

    avg = round(totals["progress_sum"] / totals["squads"]) if totals["squads"] else 0
    if squad_id is not None:
        sq = db.get(Squad, squad_id)
        scope_name = rt(lang, "squad_scope", name=sq.name) if sq else rt(lang, "h_squad")
    elif scope_tribe in tribes:
        scope_name = tribes[scope_tribe].name
    else:
        scope_name = rt(lang, "all_tribes")

    leaves_upcoming = _upcoming_leaves(db, scope_tribe, squad_id, sq_ids, now)

    return {
        "app_name": cfg.get("app_name") or "Tribe Cockpit",
        "subtitle": cfg.get("app_subtitle") or "",
        "scope_name": scope_name,
        "squad_scoped": squad_id is not None,
        "lang": lang,
        "year": year,
        "since_days": since_days,
        "generated_at": now,
        "summary": {
            "squads_total": totals["squads"],
            "blocked": totals["blocked"],
            "at_risk": totals["at_risk"],
            "objectives_red": totals["objectives_red"],
            "stale": totals["stale"],
            "avg_progress": avg,
        },
        "tribes": tribe_blocks,
        "attention": attention,
        "leaves_upcoming": leaves_upcoming,
    }


def _upcoming_leaves(db: Session, scope_tribe: int | None, squad_id: int | None,
                     sq_ids: list[int], now: datetime) -> list[dict]:
    """Approved/pending absences ending in the next 30 days, scoped like the report.
    Empty list when the leaves module is disabled."""
    from .modulesconfig import get_modules, is_active
    if not is_active(get_modules(db), "leaves"):
        return []
    from .leavesconfig import ACTIVE_STATUSES, leave_days
    from .models import Leave, Member, User

    today = now.date()
    horizon = today + timedelta(days=30)
    stmt = select(Leave).where(Leave.status.in_(ACTIVE_STATUSES),
                               Leave.end_date >= today, Leave.start_date <= horizon)
    if squad_id is not None:
        uids = list(db.scalars(select(Member.user_id).where(
            Member.squad_id == squad_id, Member.user_id.isnot(None))).all())
        stmt = stmt.where(Leave.user_id.in_(uids or [-1]))
    elif scope_tribe is not None:
        stmt = stmt.where(Leave.tribe_id == scope_tribe)
    elif sq_ids:
        uids = list(db.scalars(select(Member.user_id).where(
            Member.squad_id.in_(sq_ids), Member.user_id.isnot(None))).all())
        stmt = stmt.where(Leave.user_id.in_(uids or [-1]))

    out: list[dict] = []
    names: dict[int, str] = {}
    for lv in db.scalars(stmt.order_by(Leave.start_date, Leave.id)).all():
        if lv.user_id not in names:
            u = db.get(User, lv.user_id)
            names[lv.user_id] = u.display_name if u else f"#{lv.user_id}"
        out.append({
            "name": names[lv.user_id],
            "type_label": lv.type.label if lv.type else "",
            "detail": lv.detail or "",
            "type_color": lv.type.color if lv.type else "#6B7280",
            "start": lv.start_date.strftime("%d/%m"), "end": lv.end_date.strftime("%d/%m"),
            "days": leave_days(lv), "status": lv.status,
        })
    return out


# ----- HTML rendering -------------------------------------------------------------

def _bar(pct: int, rag: str = "green") -> str:
    pct = max(0, min(100, int(pct or 0)))
    color = RAG_COLOR.get(rag, RAG_COLOR["green"])
    return (
        f'<div class="bar"><div class="bar-fill" style="width:{pct}%;background:{color}"></div>'
        f'<span class="bar-label">{pct}%</span></div>'
    )


def _delta_html(delta: int) -> str:
    if delta > 0:
        return f'<span style="color:{RAG_COLOR["green"]}">▲ +{delta}</span>'
    if delta < 0:
        return f'<span style="color:{RAG_COLOR["red"]}">▼ {delta}</span>'
    return '<span style="color:#6b7280">→ 0</span>'


def _squad_detail_parts(r: dict, lang: str, e, *, with_title: bool = True) -> list[str]:
    """One squad's detail block, in the exact order of the squad page:
    Initiatives → OTD → Roadmap → Key messages → Budget."""
    det = r.get("detail") or {}
    parts: list[str] = ['<div class="sq-detail">']
    if with_title:
        parts.append(f'<h3>{e(r["name"])} <span class="muted">· {r["annual_pct"]}%</span></h3>')

    # Initiatives
    inits = det.get("initiatives") or []
    if inits:
        parts.append(f'<div class="d-sub">{e(rt(lang, "h_initiatives"))}</div><ul class="d-obj">')
        for ini in inits:
            meta = []
            if ini.get("owner"):
                meta.append(e(ini["owner"]))
            if ini.get("deadline"):
                meta.append(f'{e(rt(lang, "deadline"))} {e(ini["deadline"])}')
            tail = f' <span class="muted">({e(" · ".join(meta))})</span>' if meta else ""
            parts.append(f'<li>{e(ini["title"])}{tail}</li>')
        parts.append('</ul>')

    # OTD (annual objectives)
    parts.append(f'<div class="d-sub">{e(rt(lang, "h_otd_section"))}</div>')
    if det.get("objectives"):
        parts.append('<ul class="d-obj">')
        for o in det["objectives"]:
            rag = _status_rag(o["rag"])
            dl = f' · {e(rt(lang, "deadline"))} {e(o["target_date"])}' if o.get("target_date") else ""
            parts.append(f'<li><span class="dot" style="background:{RAG_COLOR[rag]}"></span>'
                         f'{e(o["title"])} <span class="muted">({e(_status_label(o["rag"], lang))}{dl})</span></li>')
        parts.append('</ul>')
    else:
        parts.append(f'<div class="muted small">{e(rt(lang, "no_obj"))}</div>')

    # Roadmap by quarter
    parts.append(f'<div class="d-sub">{e(rt(lang, "h_roadmap"))}</div><div class="d-quarters">')
    for qd in det.get("quarters", []):
        parts.append(f'<div class="d-q"><div class="d-q-head">Q{qd["q"]} '
                     f'<span class="muted">{qd["pct"]}%</span></div>')
        if qd["items"]:
            parts.append('<ul>')
            for it in qd["items"]:
                rag = _status_rag(it["status"])
                stage = f' <strong>({e(it["stage"])})</strong>' if it.get("stage") else ""
                dep = f' <span class="muted">· {e(rt(lang, "dep"))} {e(it["dependency"])}</span>' if it.get("dependency") else ""
                parts.append(f'<li><span class="dot" style="background:{RAG_COLOR[rag]}"></span>'
                             f'{e(it["title"])}{stage}{dep}</li>')
            parts.append('</ul>')
        else:
            parts.append(f'<div class="muted small">{e(rt(lang, "no_jalon"))}</div>')
        parts.append('</div>')
    parts.append('</div>')  # .d-quarters

    # Key messages
    kms = det.get("key_messages") or []
    parts.append(f'<div class="d-sub">{e(rt(lang, "h_key_messages"))}</div>')
    if kms:
        km_rag = {"success": "green", "alert": "amber", "risk": "red"}
        parts.append('<ul class="d-obj">')
        for m in kms:
            rag = km_rag.get(m["kind"], "grey")
            ts = f' <span class="muted">· {e(m["created_at"])}</span>' if m.get("created_at") else ""
            parts.append(f'<li><span class="dot" style="background:{RAG_COLOR[rag]}"></span>'
                         f'<strong>{e(rt(lang, "km_" + m["kind"]))}</strong> — {e(m["text"])}{ts}</li>')
        parts.append('</ul>')
    else:
        parts.append(f'<div class="muted small">{e(rt(lang, "no_key_message"))}</div>')

    # Budget (present only when the viewer may see this squad's figures)
    bud = det.get("budget")
    if bud is not None:
        fmtn = lambda v: "-" if v is None else f"{v:,.0f} €"
        st_rag = {"on_track": "green", "at_risk": "amber", "over": "red"}[bud["status"]]
        st_lbl = rt(lang, {"on_track": "b_on_track", "at_risk": "b_at_risk", "over": "b_over"}[bud["status"]])
        over = f' (+{fmtn(bud["overrun"])} · {bud["overrun_pct"]}%)' if bud["status"] == "over" else ""
        parts.append(f'<div class="d-sub">{e(rt(lang, "h_budget"))} '
                     f'<span class="dot" style="background:{RAG_COLOR[st_rag]}"></span> '
                     f'<span class="muted">{e(st_lbl)}{e(over)}</span></div>')
        if bud["total"] is None and bud["spent"] is None and bud["forecast"] is None:
            parts.append(f'<div class="muted small">{e(rt(lang, "no_budget"))}</div>')
        else:
            sp = f' <span class="muted">· {bud["spent_pct"]}%</span>' if bud.get("spent_pct") is not None else ""
            fp = f' <span class="muted">· {bud["forecast_pct"]}%</span>' if bud.get("forecast_pct") is not None else ""
            parts.append('<ul class="d-obj">')
            parts.append(f'<li>{e(rt(lang, "b_total"))} : <strong>{fmtn(bud["total"])}</strong></li>')
            parts.append(f'<li>{e(rt(lang, "b_spent"))} : <strong>{fmtn(bud["spent"])}</strong>{sp}</li>')
            parts.append(f'<li>{e(rt(lang, "b_forecast"))} : <strong>{fmtn(bud["forecast"])}</strong>{fp}</li>')
            if bud.get("comment"):
                parts.append(f'<li class="muted">{e(bud["comment"])}</li>')
            parts.append('</ul>')

    parts.append('</div>')  # .sq-detail
    return parts


_STATIC_ASSETS = os.path.join(os.path.dirname(__file__), "static", "assets")
_DOT_CLASS = {"green": "dot-green", "amber": "dot-orange", "red": "dot-red", "grey": "dot-grey"}
_KM_BADGE = {"success": "badge-green", "alert": "badge-orange", "risk": "badge-red"}
_BUD_BADGE = {"on_track": "badge-green", "at_risk": "badge-orange", "over": "badge-red"}


def _app_css() -> str:
    """The application's own built stylesheet (served under /assets), so a
    single-squad export renders with the exact look of the squad page. Falls back
    to the report stylesheet when no build is present (e.g. tests)."""
    try:
        for fn in sorted(os.listdir(_STATIC_ASSETS)):
            if fn.endswith(".css"):
                with open(os.path.join(_STATIC_ASSETS, fn), encoding="utf-8") as fh:
                    return f"<style>{fh.read()}</style>"
    except OSError:
        pass
    return _CSS


def _dot(rag: str) -> str:
    return f'<span class="dot {_DOT_CLASS.get(rag, "dot-green")}"></span>'


def _squad_app_cards(det: dict, lang: str, e, year: int) -> list[str]:
    """The squad page's cards, in page order: Initiatives → OTD → Roadmap →
    Key messages → Budget, using the application's own component classes."""
    fmtn = lambda v: "-" if v is None else f"{v:,.0f} €"
    C: list[str] = []

    # Initiatives — always shown (even empty), to mirror the squad page.
    inits = det.get("initiatives") or []
    C.append(f'<div class="card"><h2>{e(rt(lang, "h_initiatives"))}</h2>')
    if inits:
        C.append('<table class="init-tbl"><thead><tr>'
                 f'<th>{e(rt(lang, "h_initiatives"))}</th><th>{e(rt(lang, "h_leader"))}</th>'
                 f'<th>{e(rt(lang, "deadline"))}</th></tr></thead><tbody>')
        for ini in inits:
            C.append(f'<tr><td><strong>{e(ini["title"])}</strong></td>'
                     f'<td>{e(ini.get("owner") or "-")}</td><td>{e(ini.get("deadline") or "-")}</td></tr>')
        C.append('</tbody></table>')
    else:
        C.append(f'<div class="muted small">{e(rt(lang, "no_initiative"))}</div>')
    C.append('</div>')

    # OTD (annual objectives)
    C.append(f'<div class="card"><h2>{e(rt(lang, "h_otd_section"))} {year}</h2>')
    if det.get("objectives"):
        for o in det["objectives"]:
            dl = f' · {e(rt(lang, "deadline"))} {e(o["target_date"])}' if o.get("target_date") else ""
            C.append(f'<div class="item-row">{_dot(_status_rag(o["rag"]))}'
                     f'<div class="grow"><div>{e(o["title"])}</div></div>'
                     f'<span class="small muted">{e(_status_label(o["rag"], lang))}{dl}</span></div>')
    else:
        C.append(f'<div class="small muted">{e(rt(lang, "no_obj"))}</div>')
    C.append('</div>')

    # Roadmap by quarter
    C.append(f'<div class="card"><h2>{e(rt(lang, "h_roadmap"))} {year}</h2>'
             '<div class="grid" style="grid-template-columns:repeat(auto-fit,minmax(240px,1fr));gap:14px">')
    for qd in det.get("quarters", []):
        pct = max(0, min(100, int(qd["pct"] or 0)))
        C.append(f'<div class="quarter-block"><div class="between"><h4 style="margin:0">Q{qd["q"]}</h4>'
                 f'<span class="small muted">{qd["pct"]}%</span></div>'
                 f'<div class="progress"><div style="width:{pct}%"></div></div>')
        if qd.get("comment"):
            C.append(f'<div class="small muted" style="margin-top:6px">{e(qd["comment"])}</div>')
        C.append('<div style="margin-top:8px">')
        if not qd["items"]:
            C.append(f'<div class="small muted">{e(rt(lang, "no_jalon"))}</div>')
        for it in qd["items"]:
            stage = f'<span class="badge badge-navy" style="font-size:10px">{e(it["stage"])}</span>' if it.get("stage") else ""
            dep = f'<span class="small muted">· {e(rt(lang, "dep"))} {e(it["dependency"])}</span>' if it.get("dependency") else ""
            C.append(f'<div class="item-row">{_dot(_status_rag(it["status"]))}'
                     f'<span class="grow small">{e(it["title"])}</span>{stage}'
                     f'<span class="small muted">{e(_status_label(it["status"], lang))}</span>{dep}</div>')
        C.append('</div></div>')
    C.append('</div></div>')

    # Key messages
    C.append(f'<div class="card"><h2>{e(rt(lang, "h_key_messages"))}</h2>')
    kms = det.get("key_messages") or []
    if kms:
        for m in kms:
            ts = f'<div class="small muted">{e(m["created_at"])}</div>' if m.get("created_at") else ""
            C.append(f'<div class="item-row"><span class="badge {_KM_BADGE.get(m["kind"], "badge-grey")}">'
                     f'{e(rt(lang, "km_" + m["kind"]))}</span>'
                     f'<div class="grow"><div class="small">{e(m["text"])}</div>{ts}</div></div>')
    else:
        C.append(f'<div class="small muted">{e(rt(lang, "no_key_message"))}</div>')
    C.append('</div>')

    # Budget (present only when the viewer may see the figures)
    bud = det.get("budget")
    if bud is not None:
        st_lbl = rt(lang, {"on_track": "b_on_track", "at_risk": "b_at_risk", "over": "b_over"}[bud["status"]])
        over = f' (+{fmtn(bud["overrun"])} · {bud["overrun_pct"]}%)' if bud["status"] == "over" else ""
        C.append(f'<div class="card"><div class="between"><h2 style="margin:0">{e(rt(lang, "h_budget"))}</h2>'
                 f'<span class="badge {_BUD_BADGE[bud["status"]]}">{e(st_lbl)}{e(over)}</span></div>')
        if bud["total"] is None and bud["spent"] is None and bud["forecast"] is None:
            C.append(f'<div class="small muted">{e(rt(lang, "no_budget"))}</div>')
        else:
            sp = f' · {bud["spent_pct"]}%' if bud.get("spent_pct") is not None else ""
            fp = f' · {bud["forecast_pct"]}%' if bud.get("forecast_pct") is not None else ""
            C.append('<div class="stack" style="gap:6px;margin-top:6px">'
                     f'<div class="between"><span class="small muted">{e(rt(lang, "b_total"))}</span>'
                     f'<span class="strong">{fmtn(bud["total"])}</span></div>'
                     f'<div class="between"><span class="small muted">{e(rt(lang, "b_spent"))}</span>'
                     f'<span class="strong">{fmtn(bud["spent"])}{sp}</span></div>'
                     f'<div class="between"><span class="small muted">{e(rt(lang, "b_forecast"))}</span>'
                     f'<span class="strong">{fmtn(bud["forecast"])}{fp}</span></div>')
            if bud.get("comment"):
                C.append(f'<div class="small muted" style="margin-top:4px">{e(bud["comment"])}</div>')
            C.append('</div>')
        C.append('</div>')

    return C


def _render_squad_page(data: dict, standalone: bool, e, lang: str) -> str:
    """Single-squad export rendered with the application's own stylesheet and
    component markup, so it looks exactly like the squad page."""
    r = next((rr for blk in data["tribes"] for rr in blk["squads"] if rr.get("detail")), None)
    style = _app_css()
    if r is None:
        body = f'<div class="export-page"><h1>{e(data["scope_name"])}</h1></div>'
    else:
        year = data["year"]
        fresh_cls = "badge-grey" if r["is_stale"] else "badge-navy"
        fresh_lbl = rt(lang, "stale") if r["is_stale"] else rt(lang, "h_freshness_ok")
        badges = [f'<span class="badge badge-navy">{e(rt(lang, "h_progress_long"))} {r["annual_pct"]}%</span>']
        if r["blocked"]:
            badges.append(f'<span class="badge badge-red">{r["blocked"]} {e(rt(lang, "h_blocked"))}</span>')
        if r["at_risk"]:
            badges.append(f'<span class="badge badge-orange">{r["at_risk"]} {e(rt(lang, "h_atrisk"))}</span>')
        badges.append(f'<span class="badge {fresh_cls}">{e(fresh_lbl)}</span>')
        P = [f'<div class="export-page"><h1 style="color:var(--navy);margin:0 0 8px">{e(r["name"])}</h1>',
             f'<div class="inline" style="gap:10px;flex-wrap:wrap;margin-bottom:6px">{"".join(badges)}</div>',
             f'<div class="muted small" style="margin-bottom:16px">{e(rt(lang, "h_leader"))} : '
             f'<span class="strong">{e(r["leader"] or "-")}</span> · {e(rt(lang, "year"))} {year}</div>',
             '<div class="stack" style="gap:18px">']
        P.extend(_squad_app_cards(r["detail"], lang, e, year))
        P.append('</div></div>')
        body = "\n".join(P)
    page_css = ('<style>body{background:var(--bg,#F5F7FA);margin:0;padding:24px;color:var(--text,#1E293B);'
                'font-family:-apple-system,Segoe UI,Roboto,Helvetica,Arial,sans-serif}'
                '.export-page{max-width:1040px;margin:0 auto}'
                '.init-tbl{width:100%;border-collapse:collapse}'
                '.init-tbl th,.init-tbl td{text-align:left;padding:6px 8px;border-bottom:1px solid var(--line,#E2E8F0)}</style>')
    if not standalone:
        return f'{style}{page_css}{body}'
    title = e(r["name"]) if r else e(data["scope_name"])
    return (f'<!doctype html><html lang="{e(lang)}"><head><meta charset="utf-8">'
            f'<meta name="viewport" content="width=device-width, initial-scale=1">'
            f'<title>{title}</title>{style}{page_css}</head><body>{body}</body></html>')


# =============================================================================
# "What's new since your last report" — change detection against a stored
# per-scope baseline (see models.ReportBaseline).
# =============================================================================

def _ms_lbl(status: str, lang: str) -> str:
    return rt(lang, {"on_track": "ms_on_track", "at_risk": "ms_at_risk",
                     "blocked": "ms_blocked", "done": "ms_done"}.get(status, "h_status")) \
        if status in ("on_track", "at_risk", "blocked", "done") else status


def _rag_lbl(rag: str, lang: str) -> str:
    return rt(lang, {"green": "rag_green", "amber": "rag_amber", "red": "rag_red"}.get(rag, "rag_green")) \
        if rag in ("green", "amber", "red") else rag


def report_signature(data: dict) -> dict:
    """Compact, diff-friendly snapshot of a report's per-squad state."""
    sig: dict = {}
    for blk in data.get("tribes", []):
        for r in blk.get("squads", []):
            d = r.get("detail", {}) or {}
            b = d.get("budget") or None
            sig[str(r["squad_id"])] = {
                "name": r["name"],
                "status": r.get("status"),
                "annual_pct": r.get("annual_pct", 0),
                "is_stale": bool(r.get("is_stale")),
                "milestones": {it["title"]: it["status"]
                               for q in d.get("quarters", []) for it in q.get("items", [])},
                "objectives": {o["title"]: o["rag"] for o in d.get("objectives", [])},
                "budget": {k: b.get(k) for k in ("total", "spent", "forecast", "status")} if b else None,
                "km": len(d.get("key_messages", [])),
            }
    return sig


def diff_report(prev: dict | None, cur: dict, lang: str) -> dict:
    """Compare two report signatures. Returns
    {first, count, summary, by_squad:[{name, items:[str]}]}."""
    if not prev:
        return {"first": True, "count": 0, "summary": "", "by_squad": []}
    by_squad: list[dict] = []
    tally = {"moved": 0, "delivered": 0, "blocked": 0, "stale": 0}
    for sid, c in cur.items():
        p = prev.get(sid)
        items: list[str] = []
        if p is None:
            items.append(rt(lang, "chg_new_squad", name=c["name"]))
        else:
            if c["annual_pct"] != p["annual_pct"]:
                d = c["annual_pct"] - p["annual_pct"]
                items.append(rt(lang, "chg_progress", d=(f"+{d}" if d > 0 else str(d))))
                tally["moved"] += 1
            if c.get("status") != p.get("status") and p.get("status"):
                items.append(rt(lang, "chg_status", frm=p["status"], to=c["status"]))
            for title, stt in c["milestones"].items():
                if title not in p["milestones"]:
                    items.append(rt(lang, "chg_ms_new", title=title))
                elif p["milestones"][title] != stt:
                    items.append(rt(lang, "chg_ms_status", title=title,
                                    frm=_ms_lbl(p["milestones"][title], lang), to=_ms_lbl(stt, lang)))
                    if stt == "done":
                        tally["delivered"] += 1
                    if stt == "blocked" and p["milestones"][title] != "blocked":
                        tally["blocked"] += 1
            for title in p["milestones"]:
                if title not in c["milestones"]:
                    items.append(rt(lang, "chg_ms_removed", title=title))
            for title, rag in c["objectives"].items():
                if title not in p["objectives"]:
                    items.append(rt(lang, "chg_obj_new", title=title))
                elif p["objectives"][title] != rag:
                    items.append(rt(lang, "chg_obj_status", title=title,
                                    frm=_rag_lbl(p["objectives"][title], lang), to=_rag_lbl(rag, lang)))
            if (c.get("budget") or {}) != (p.get("budget") or {}):
                items.append(rt(lang, "chg_budget"))
            if c.get("km", 0) > p.get("km", 0):
                items.append(rt(lang, "chg_km", n=c["km"] - p["km"]))
            if c["is_stale"] and not p["is_stale"]:
                items.append(rt(lang, "chg_stale"))
                tally["stale"] += 1
            elif not c["is_stale"] and p["is_stale"]:
                items.append(rt(lang, "chg_unstale"))
        if items:
            by_squad.append({"name": c["name"], "items": items})
    for sid, p in prev.items():
        if sid not in cur:
            by_squad.append({"name": p["name"], "items": [rt(lang, "chg_squad_removed", name=p["name"])]})

    count = sum(len(s["items"]) for s in by_squad)
    parts = []
    if tally["moved"]:
        parts.append(rt(lang, "sum_moved", n=tally["moved"]))
    if tally["delivered"]:
        parts.append(rt(lang, "sum_delivered", n=tally["delivered"]))
    if tally["blocked"]:
        parts.append(rt(lang, "sum_blocked", n=tally["blocked"]))
    if tally["stale"]:
        parts.append(rt(lang, "sum_stale", n=tally["stale"]))
    return {"first": False, "count": count, "summary": " · ".join(parts), "by_squad": by_squad}


def subject_prefix(changes: dict | None, lang: str) -> str:
    """Subject tag: '[3 nouveautés] ' / '[à jour] ' / '' (first report)."""
    if not changes or changes.get("first"):
        return ""
    if changes["count"] == 0:
        return rt(lang, "subj_uptodate") + " "
    return rt(lang, "subj_changes", n=changes["count"]) + " "


def render_changes_html(changes: dict | None, lang: str) -> str:
    if not changes:
        return ""
    e = html.escape
    head = f'<div class="chg-h">{e(rt(lang, "whatsnew"))}</div>'
    if changes.get("first"):
        return f'<div class="changes-box"><div class="chg-h">{e(rt(lang, "whatsnew"))}</div>' \
               f'<div class="chg-empty">{e(rt(lang, "first_report"))}</div></div>'
    if changes["count"] == 0:
        return f'<div class="changes-box uptodate">{head}' \
               f'<div class="chg-empty">✓ {e(rt(lang, "no_changes"))}</div></div>'
    out = [f'<div class="changes-box">{head}']
    if changes["summary"]:
        out.append(f'<div class="chg-sum">{e(changes["summary"])}</div>')
    for sq in changes["by_squad"][:12]:
        out.append(f'<div class="chg-sq"><span class="chg-sqn">{e(sq["name"])}</span><ul>')
        for it in sq["items"][:8]:
            out.append(f'<li>{e(it)}</li>')
        out.append('</ul></div>')
    out.append('</div>')
    return "".join(out)


def get_baseline(db, scope_key: str) -> dict | None:
    from .models import ReportBaseline
    row = db.get(ReportBaseline, scope_key)
    return row.signature if row else None


def set_baseline(db, scope_key: str, signature: dict) -> None:
    from .models import ReportBaseline
    row = db.get(ReportBaseline, scope_key)
    if row is None:
        db.add(ReportBaseline(scope_key=scope_key, signature=signature, updated_at=utcnow()))
    else:
        row.signature = signature
        row.updated_at = utcnow()


def render_html(data: dict, *, standalone: bool = True, changes: dict | None = None) -> str:
    e = html.escape
    lang = data.get("lang", "fr")
    # A single-squad export mirrors the squad page, not the whole dashboard report.
    if data.get("squad_scoped"):
        return _render_squad_page(data, standalone, e, lang)
    s = data["summary"]
    gen = data["generated_at"]
    gen_str = gen.strftime("%d/%m/%Y %H:%M") if isinstance(gen, datetime) else str(gen)

    parts: list[str] = []
    parts.append(f'<div class="hdr"><h1>{e(data["app_name"])} - {e(rt(lang, "report"))}</h1>')
    parts.append(f'<div class="sub">{e(data["scope_name"])} · {e(rt(lang, "year"))} {data["year"]} · '
                 f'{e(rt(lang, "generated"))} {e(gen_str)} · {e(rt(lang, "window", n=data["since_days"]))}</div></div>')

    # "What's new since your last report" — right under the header.
    if changes is not None:
        parts.append(render_changes_html(changes, lang))

    # Summary cards
    cards = [
        (rt(lang, "k_squads"), s["squads_total"], "#111827"),
        (rt(lang, "k_progress"), f'{s["avg_progress"]}%', RAG_COLOR["green"]),
        (rt(lang, "k_blocked"), s["blocked"], RAG_COLOR["red"] if s["blocked"] else "#111827"),
        (rt(lang, "k_atrisk"), s["at_risk"], RAG_COLOR["amber"] if s["at_risk"] else "#111827"),
        (rt(lang, "k_obj_red"), s["objectives_red"], RAG_COLOR["red"] if s["objectives_red"] else "#111827"),
        (rt(lang, "k_stale"), s["stale"], RAG_COLOR["amber"] if s["stale"] else "#111827"),
    ]
    parts.append('<div class="cards">')
    for label, val, color in cards:
        parts.append(f'<div class="kpi"><div class="kpi-val" style="color:{color}">{e(str(val))}</div>'
                     f'<div class="kpi-lbl">{e(label)}</div></div>')
    parts.append('</div>')

    # Attention list
    if data["attention"]:
        parts.append(f'<h2>{e(rt(lang, "attention"))}</h2><ul class="attention">')
        for r in data["attention"][:12]:
            bits = []
            if r["blocked"]:
                bits.append(f'{r["blocked"]} {rt(lang, "blocked_n")}')
            if r["delta"] < 0:
                bits.append(f'{r["delta"]} pt')
            if r["is_stale"]:
                bits.append(rt(lang, "stale"))
            parts.append(f'<li><span class="dot" style="background:{RAG_COLOR["red"]}"></span>'
                         f'<strong>{e(r["name"])}</strong> - {e(", ".join(bits))}</li>')
        parts.append('</ul>')

    # Upcoming absences (next 30 days), when the leaves module is enabled.
    if data.get("leaves_upcoming"):
        parts.append(f'<h2>{e(rt(lang, "leaves_upcoming"))}</h2><ul class="attention">')
        for lv in data["leaves_upcoming"][:30]:
            pending = f' <span class="badge">{e(rt(lang, "leaves_pending"))}</span>' if lv["status"] == "pending" else ""
            parts.append(
                f'<li><span class="dot" style="background:{e(lv["type_color"])}"></span>'
                f'<strong>{e(lv["name"])}</strong> — {e(leave_type_label(lv["type_label"], lang))}'
                f'{e(" (" + lv["detail"] + ")") if lv.get("detail") else ""} '
                f'<span class="muted">({e(lv["start"])} → {e(lv["end"])}, {lv["days"]:g} {e(rt(lang, "days_short"))})</span>'
                f'{pending}</li>')
        parts.append('</ul>')

    # Per-tribe tables
    for blk in data["tribes"]:
        parts.append(f'<h2>{e(blk["tribe_name"])}</h2>')
        parts.append('<table><thead><tr>'
                     f'<th>{e(rt(lang, "h_squad"))}</th><th>{e(rt(lang, "h_leader"))}</th>'
                     f'<th>{e(rt(lang, "h_status"))}</th>'
                     f'<th>{e(rt(lang, "h_progress_long"))}</th><th>{e(rt(lang, "h_delta"))}</th>'
                     f'<th>{e(rt(lang, "h_blocked"))}</th>'
                     f'<th>{e(rt(lang, "h_atrisk"))}</th><th>{e(rt(lang, "h_facts"))}</th>'
                     '</tr></thead><tbody>')
        for r in blk["squads"]:
            changes = r["changes"][:4]
            ch_html = "<br>".join(e(_change_text(c, lang)) for c in changes) if changes else \
                ('<span class="muted">-</span>' if not r["note"] else "")
            if r["note"]:
                note = e(r["note"]).replace("\n", " ")
                if len(note) > 160:
                    note = note[:159] + "…"
                ch_html = (ch_html + "<br>" if ch_html else "") + f'<em class="note">« {note} »</em>'
            stale_badge = f' <span class="badge">{e(rt(lang, "stale"))}</span>' if r["is_stale"] else ""
            parts.append(
                f'<tr><td><strong>{e(r["name"])}</strong>{stale_badge}</td>'
                f'<td>{e(r["leader"])}</td>'
                f'<td><span class="pill" style="background:{RAG_COLOR[r["status_rag"]]}">'
                f'{e(_status_label(r["status"], lang))}</span></td>'
                f'<td>{_bar(r["annual_pct"], r["status_rag"])}</td>'
                f'<td>{_delta_html(r["delta"])}</td>'
                f'<td>{r["blocked"] or ""}</td><td>{r["at_risk"] or ""}</td>'
                f'<td class="changes">{ch_html}</td></tr>'
            )
        parts.append('</tbody></table>')

    # --- Full detail per squad: annual objectives + roadmap/milestones by quarter.
    all_squads = [r for blk in data["tribes"] for r in blk["squads"]]
    if any(r.get("detail") for r in all_squads):
        parts.append(f'<h2 class="detail-h">{e(rt(lang, "detail_title"))}</h2>')
    for r in all_squads:
        if r.get("detail"):
            parts.extend(_squad_detail_parts(r, lang, e))

    body = "\n".join(parts)
    if not standalone:
        return f'<div class="tc-report">{_CSS}{body}</div>'
    return (
        f'<!doctype html><html lang="{e(lang)}"><head><meta charset="utf-8">'
        f'<title>{e(data["app_name"])} - {e(rt(lang, "report"))}</title>'
        f'{_CSS}</head><body><div class="tc-report">{body}</div></body></html>'
    )


def render_roadmap_html(data: dict, *, standalone: bool = True) -> str:
    """Roadmap matrix web page: quarters in columns, squads (themes) in rows,
    milestone titles in the cells, colour-coded by status."""
    e = html.escape
    lang = data.get("lang", "fr")
    year = data["year"]
    gen = data["generated_at"]
    gen_str = gen.strftime("%d/%m/%Y %H:%M") if isinstance(gen, datetime) else str(gen)
    squads = [row for blk in data["tribes"] for row in blk["squads"]]

    parts: list[str] = []
    parts.append(f'<div class="hdr"><h1>{e(data["app_name"])} - {e(rt(lang, "roadmap_report"))}</h1>')
    parts.append(f'<div class="sub">{e(data["scope_name"])} · {e(rt(lang, "year"))} {year} · '
                 f'{e(rt(lang, "generated"))} {e(gen_str)}</div></div>')

    # EA/GA legend (status is no longer colour-coded in the roadmap view)
    parts.append('<div class="rm-legend">')
    parts.append(f'<span><b class="rm-ea">EA</b> {e(rt(lang, "stage_ea"))}</span>')
    parts.append(f'<span><b class="rm-ga">GA</b> {e(rt(lang, "stage_ga"))}</span>')
    parts.append('</div>')

    months = _MONTHS[lang]
    parts.append('<table class="rm"><thead>')
    parts.append('<tr><th class="rm-corner" rowspan="2"></th>')
    for q in (1, 2, 3, 4):
        parts.append(f'<th class="rm-q" colspan="3">Q{q} {year}</th>')
    parts.append('</tr><tr>')
    for mi in range(12):
        parts.append(f'<th class="rm-m">{e(months[mi])}</th>')
    parts.append('</tr></thead><tbody>')
    for sq in squads:
        parts.append(f'<tr><th class="rm-row">{e(sq["name"])}</th>')
        qmap = {qd["q"]: qd["items"] for qd in (sq.get("detail") or {}).get("quarters", [])}
        for q in (1, 2, 3, 4):
            items = qmap.get(q, [])
            parts.append('<td colspan="3"><div class="rm-card">' if items else '<td colspan="3">')
            for theme, group in group_by_theme(items):
                if theme:
                    parts.append(f'<div class="rm-theme">{e(theme)}</div>')
                for it in group:
                    stage = it.get("stage")
                    st_html = ""
                    if stage:
                        cls = "rm-ea" if stage == "EA" else "rm-ga"
                        st_html = f' (<span class="{cls}">{e(stage)}</span>)'
                    parts.append(f'<div class="rm-j">{e(it["title"])}{st_html}</div>')
            parts.append('</div></td>' if items else '</td>')
        parts.append('</tr>')
    parts.append('</tbody></table>')

    body = "\n".join(parts)
    if not standalone:
        return f'<div class="tc-report">{_CSS}{body}</div>'
    return (
        f'<!doctype html><html lang="{e(lang)}"><head><meta charset="utf-8">'
        f'<title>{e(data["app_name"])} - {e(rt(lang, "roadmap_report"))}</title>'
        f'{_CSS}</head><body><div class="tc-report">{body}</div></body></html>'
    )


_CSS = """<style>
.tc-report{font-family:-apple-system,Segoe UI,Roboto,Helvetica,Arial,sans-serif;color:#111827;
  max-width:980px;margin:0 auto;padding:24px;line-height:1.45;background:#fff}
.tc-report h1{font-size:22px;margin:0 0 4px}
.tc-report h2{font-size:16px;margin:26px 0 10px;border-bottom:2px solid #e5e7eb;padding-bottom:4px}
.tc-report .sub{color:#6b7280;font-size:13px}
.tc-report .cards{display:flex;flex-wrap:wrap;gap:10px;margin-top:16px}
.tc-report .kpi{flex:1;min-width:120px;border:1px solid #e5e7eb;border-radius:10px;padding:12px 14px;background:#f9fafb}
.tc-report .kpi-val{font-size:24px;font-weight:700}
.tc-report .kpi-lbl{font-size:12px;color:#6b7280;margin-top:2px}
.tc-report table{width:100%;border-collapse:collapse;font-size:13px}
.tc-report th{text-align:left;background:#f3f4f6;padding:7px 9px;border-bottom:2px solid #e5e7eb;font-size:12px;color:#374151}
.tc-report td{padding:7px 9px;border-bottom:1px solid #eef0f3;vertical-align:top}
.tc-report .pill{color:#fff;border-radius:999px;padding:2px 9px;font-size:11px;white-space:nowrap}
.tc-report .bar{position:relative;background:#eef0f3;border-radius:6px;height:16px;width:120px;overflow:hidden}
.tc-report .bar-fill{position:absolute;left:0;top:0;bottom:0;border-radius:6px}
.tc-report .bar-label{position:relative;font-size:11px;padding-left:6px;line-height:16px;color:#111827}
.tc-report .changes{color:#374151;font-size:12px}
.tc-report .changes-box{border:1px solid #dbe4ff;background:#f5f8ff;border-left:4px solid #175CD3;border-radius:10px;padding:12px 16px;margin:16px 0}
.tc-report .changes-box.uptodate{border-color:#d1fae5;background:#f0fdf6;border-left-color:#059669}
.tc-report .chg-h{font-size:13px;font-weight:700;color:#1E2761;text-transform:uppercase;letter-spacing:.03em;margin-bottom:6px}
.tc-report .chg-sum{font-size:14px;font-weight:600;color:#111827;margin-bottom:8px}
.tc-report .chg-empty{font-size:13px;color:#4b5563}
.tc-report .chg-sq{margin:6px 0}
.tc-report .chg-sqn{font-weight:700;font-size:13px;color:#175CD3}
.tc-report .chg-sq ul{margin:2px 0 0;padding-left:18px}
.tc-report .chg-sq li{font-size:12.5px;padding:1px 0}
.tc-report .note{color:#4b5563}
.tc-report .muted{color:#9ca3af}
.tc-report .badge{background:#fef3c7;color:#92400e;border-radius:4px;padding:1px 5px;font-size:10px}
.tc-report ul.attention{list-style:none;padding:0;margin:0}
.tc-report ul.attention li{padding:5px 0;font-size:13px}
.tc-report .dot{display:inline-block;width:9px;height:9px;border-radius:50%;margin-right:8px}
.tc-report h2.detail-h{margin-top:30px}
.tc-report .sq-detail{border:1px solid #e5e7eb;border-radius:10px;padding:12px 14px;margin:12px 0;background:#fff;break-inside:avoid}
.tc-report .sq-detail h3{margin:0 0 6px;font-size:15px}
.tc-report .d-sub{font-size:12px;font-weight:700;color:#374151;margin:10px 0 4px;text-transform:uppercase;letter-spacing:.03em}
.tc-report ul.d-obj{list-style:none;padding:0;margin:0}
.tc-report ul.d-obj li{padding:3px 0;font-size:13px}
.tc-report .d-quarters{display:flex;flex-wrap:wrap;gap:10px}
.tc-report .d-q{flex:1;min-width:170px;border:1px solid #eef0f3;border-radius:8px;padding:8px 10px;background:#f9fafb}
.tc-report .d-q-head{font-weight:700;font-size:13px;margin-bottom:4px}
.tc-report .d-q ul{list-style:none;padding:0;margin:0}
.tc-report .d-q li{padding:2px 0;font-size:12px}
.tc-report .rm-legend{display:flex;gap:16px;margin:14px 0 8px;font-size:12px;color:#374151}
.tc-report table.rm{table-layout:fixed;border-collapse:separate;border-spacing:3px}
.tc-report table.rm th.rm-q{background:#304957;color:#fff;text-align:center;font-size:15px;padding:7px}
.tc-report table.rm th.rm-m{background:#97A3AA;color:#fff;text-align:center;font-size:11px;font-weight:500;padding:3px}
.tc-report table.rm th.rm-corner{background:transparent;border:none}
.tc-report table.rm th.rm-row{background:#304957;color:#fff;text-align:center;width:58px;vertical-align:middle;font-size:12.5px;writing-mode:vertical-rl;transform:rotate(180deg);padding:7px 4px}
.tc-report table.rm td{vertical-align:top;padding:0}
.tc-report .rm-card{background:#F2F2F2;border-radius:8px;padding:8px 10px;height:100%;box-sizing:border-box}
.tc-report .rm-theme{font-size:13px;font-weight:800;color:#002060;margin:6px 0 1px;line-height:1.3}
.tc-report .rm-theme:first-child{margin-top:0}
.tc-report .rm-j{font-size:12.5px;padding:1px 0 1px 10px;line-height:1.32;color:#002060}
.tc-report .rm-legend b.rm-ea,.tc-report .rm-j .rm-ea{color:#FFC000;font-weight:800}
.tc-report .rm-legend b.rm-ga,.tc-report .rm-j .rm-ga{color:#00B050;font-weight:800}
</style>"""


# ----- PPTX rendering -------------------------------------------------------------

# Brand palette (mirrors the app theme): navy / accent + RAG.
_BRAND = {
    "navy": "#1E2761", "navy_deep": "#141B47", "accent": "#175CD3",
    "green": "#027A48", "orange": "#B54708", "red": "#B42318",
    "ink": "#1F2937", "muted": "#6B7280", "card": "#F1F5F9",
    "line": "#E2E8F0", "white": "#FFFFFF", "zebra": "#F8FAFC",
}
_RAG_BRAND = {"red": "#B42318", "amber": "#B54708", "green": "#027A48", "grey": "#6B7280"}


_MAX_DETAIL_SLIDES = 40  # safety cap on per-squad slides in a deck


def _pptx_toolkit():
    """Import python-pptx lazily and return the bits used to build a deck."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    return Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE


def render_pptx(data: dict) -> bytes:
    """Render the weekly report as a branded deck (requires python-pptx):
    a summary one-pager + one full detail slide per squad (annual objectives +
    roadmap/milestones by quarter + advancement). The roadmap-only swimlane deck
    is produced separately by render_roadmap_pptx."""
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE = _pptx_toolkit()

    def rgb(hexstr: str) -> RGBColor:
        return RGBColor.from_string(hexstr.lstrip("#").upper())

    B = {k: rgb(v) for k, v in _BRAND.items()}
    lang = data.get("lang", "fr")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def new_slide():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def textbox(s, left, top, width, height, text, size, *, bold=False, color=None,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        box = s.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color if color is not None else B["ink"]
        return box

    def rect(s, left, top, width, height, fill, line=None):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line
        sh.shadow.inherit = False
        return sh

    def place(sh, lines, *, anchor=MSO_ANCHOR.TOP, ml=0.1, mt=0.06, mr=0.1, mb=0.06):
        """Write paragraphs INTO a shape's own text frame, so the text is part of
        the shape (not a separate textbox floating on top). Each line is
        (text, size_pt, color[, bold, align, space_after_pt])."""
        tf = sh.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = Inches(ml); tf.margin_right = Inches(mr)
        tf.margin_top = Inches(mt); tf.margin_bottom = Inches(mb)
        for i, ln in enumerate(lines):
            txt, size, color = ln[0], ln[1], ln[2]
            bold = ln[3] if len(ln) > 3 else False
            align = ln[4] if len(ln) > 4 else PP_ALIGN.LEFT
            sa = ln[5] if len(ln) > 5 else 2
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(sa)
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        return sh

    def bullets(s, left, top, width, height, lines, size):
        """A text box with one paragraph per (text, color, bold) line."""
        box = s.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        for i, (txt, color, bold) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(2)
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        return box

    def style_cell(cell, text, size, color, *, bold=False, align=PP_ALIGN.LEFT, fill=None):
        cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
        cell.margin_top = Emu(0); cell.margin_bottom = Emu(0)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fill is not None:
            cell.fill.solid(); cell.fill.fore_color.rgb = fill
        else:
            cell.fill.background()
        cell.text = text or " "
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        if p.runs:
            r = p.runs[0]
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

    # ----- app-styled primitives (rounded cards, chips, progress bars) -----------
    def rrect(s, left, top, width, height, fill, *, line=None, radius=0.08):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line; sh.line.width = Pt(1)
        sh.shadow.inherit = False
        return sh

    def card(s, left, top, width, height, title=None):
        sh = rrect(s, left, top, width, height, B["white"], line=B["line"], radius=0.05)
        if title:
            # Title lives in the card's own text frame (top-anchored), not as an
            # overlay; body content is added by the caller as child shapes.
            place(sh, [(title, 12, B["navy"], True)], anchor=MSO_ANCHOR.TOP, ml=0.18, mt=0.12, mr=0.18)
        return sh

    def chip(s, left, top, text, fill, *, color=None, size=10):
        w = Inches(0.26 + 0.082 * len(text))
        sh = rrect(s, left, top, w, Inches(0.3), fill, radius=0.5)
        tf = sh.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = Inches(0.07)
        tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = color if color is not None else B["white"]
        return Emu(int(left) + int(w))

    def chips_row(s, left, top, items, *, gap=0.1):
        x = int(left)
        for text, fill in items:
            x = int(chip(s, Emu(x), top, text, fill))
            x += int(Inches(gap))

    def pbar(s, left, top, width, pct, color):
        rrect(s, left, top, width, Inches(0.12), rgb("#E2E8F0"), radius=0.5)
        p = max(0, min(100, int(pct or 0)))
        if p > 0:
            rrect(s, left, top, Emu(int(int(width) * p / 100)), Inches(0.12), color, radius=0.5)

    gen = data["generated_at"]
    gen_str = gen.strftime("%d/%m/%Y %H:%M") if isinstance(gen, datetime) else str(gen)
    sm = data["summary"]
    squads = [r for blk in data["tribes"] for r in blk["squads"]]
    margin = Inches(0.5)

    # ---------------- Summary one-pager (report kind) -----------------------------
    def summary_slide():
        s = new_slide()
        band = rect(s, Inches(0), Inches(0), prs.slide_width, Inches(1.12), B["navy"])
        # Title + subtitle are the band's own text; the right-corner meta stays a
        # corner label (left title and right meta can't share one text frame).
        place(band, [
            (f'{data["app_name"]} - {rt(lang, "report")}', 26, B["white"], True, PP_ALIGN.LEFT, 5),
            (f'{data["scope_name"]} · {rt(lang, "year")} {data["year"]}', 13, rgb("#C7D2FE"), False, PP_ALIGN.LEFT, 0),
        ], anchor=MSO_ANCHOR.TOP, ml=0.55, mt=0.18, mr=4.2)
        textbox(s, Inches(9.3), Inches(0.3), Inches(3.5), Inches(0.6),
                f'{rt(lang, "generated_full", d=gen_str)}\n{rt(lang, "window_full", n=data["since_days"])}', 11,
                color=rgb("#C7D2FE"), align=PP_ALIGN.RIGHT)

        kpis = [
            (rt(lang, "k_squads"), str(sm["squads_total"]), B["navy"]),
            (rt(lang, "k_progress"), f'{sm["avg_progress"]}%', B["accent"]),
            (rt(lang, "k_blocked"), str(sm["blocked"]), B["red"] if sm["blocked"] else B["ink"]),
            (rt(lang, "k_atrisk"), str(sm["at_risk"]), B["orange"] if sm["at_risk"] else B["ink"]),
            (rt(lang, "k_obj_red"), str(sm["objectives_red"]), B["red"] if sm["objectives_red"] else B["ink"]),
            (rt(lang, "k_stale"), str(sm["stale"]), B["orange"] if sm["stale"] else B["ink"]),
        ]
        gap = Inches(0.14)
        n = len(kpis)
        total_w = prs.slide_width - margin * 2
        card_w = Emu(int((total_w - gap * (n - 1)) / n))
        ky, kh = Inches(1.34), Inches(1.05)
        for i, (label, val, color) in enumerate(kpis):
            left = Emu(int(margin) + i * (int(card_w) + int(gap)))
            kp = rect(s, left, ky, card_w, kh, B["card"], line=B["line"])
            # Value + label are the card's own text, vertically centered.
            place(kp, [(val, 26, color, True, PP_ALIGN.CENTER, 4),
                       (label, 10, B["muted"], False, PP_ALIGN.CENTER, 0)],
                  anchor=MSO_ANCHOR.MIDDLE, ml=0.05, mr=0.05)

        headers = [rt(lang, "h_squad"), rt(lang, "h_leader"), rt(lang, "h_status"),
                   rt(lang, "h_progress"), rt(lang, "h_delta"), rt(lang, "h_blocked"), rt(lang, "h_atrisk")]
        wfrac = [0.275, 0.21, 0.12, 0.105, 0.082, 0.103, 0.105]
        table_w = int(prs.slide_width - margin * 2)
        widths = [Emu(int(table_w * f)) for f in wfrac]

        MAX = 18
        shown = squads[:MAX]
        overflow = len(squads) - len(shown)
        nrows = len(shown) + 1 + (1 if overflow > 0 else 0)
        top = Inches(2.62)
        tbl = s.shapes.add_table(max(nrows, 2), len(headers), margin, top, Emu(table_w),
                                 Inches(0.34) + Inches(0.255) * (nrows - 1)).table
        for ci, w in enumerate(widths):
            tbl.columns[ci].width = w
        for ci, h in enumerate(headers):
            align = PP_ALIGN.LEFT if ci < 2 else PP_ALIGN.CENTER
            style_cell(tbl.cell(0, ci), h, 10, B["white"], bold=True, align=align, fill=B["navy"])
        for ri, r in enumerate(shown, start=1):
            zebra = B["zebra"] if ri % 2 == 0 else B["white"]
            delta = r["delta"]
            cells = [
                (r["name"], B["ink"], True, PP_ALIGN.LEFT),
                (r["leader"] or "-", B["muted"], False, PP_ALIGN.LEFT),
                (_status_label(r["status"], lang), rgb(_RAG_BRAND[r["status_rag"]]), True, PP_ALIGN.CENTER),
                (f'{r["annual_pct"]}%', B["ink"], False, PP_ALIGN.CENTER),
                ((f'+{delta}' if delta > 0 else str(delta)),
                 (B["green"] if delta > 0 else B["red"]) if delta else B["muted"], False, PP_ALIGN.CENTER),
                (str(r["blocked"] or "-"), B["red"] if r["blocked"] else B["muted"], r["blocked"] > 0, PP_ALIGN.CENTER),
                (str(r["at_risk"] or "-"), B["orange"] if r["at_risk"] else B["muted"], r["at_risk"] > 0, PP_ALIGN.CENTER),
            ]
            for ci, (val, color, bold, align) in enumerate(cells):
                style_cell(tbl.cell(ri, ci), val, 9.5, color, bold=bold, align=align, fill=zebra)
        if overflow > 0:
            last = len(shown) + 1
            style_cell(tbl.cell(last, 0), rt(lang, "more_squads", n=overflow), 9, B["muted"], align=PP_ALIGN.LEFT)
            for ci in range(1, len(headers)):
                style_cell(tbl.cell(last, ci), " ", 9, B["muted"])

    # ---------------- Per-squad detail slide --------------------------------------
    def detail_slide(r, *, with_objectives):
        det = r.get("detail") or {}
        s = new_slide()
        band = rect(s, Inches(0), Inches(0), prs.slide_width, Inches(0.92), B["navy"])
        place(band, [
            (r["name"], 22, B["white"], True, PP_ALIGN.LEFT, 4),
            (f'{rt(lang, "year")} {data["year"]} · {rt(lang, "h_progress_long")} {r["annual_pct"]}%',
             12, rgb("#C7D2FE"), False, PP_ALIGN.LEFT, 0),
        ], anchor=MSO_ANCHOR.TOP, ml=0.5, mt=0.12, mr=0.5)

        # Budget line (status + figures), just under the header band.
        bud = det.get("budget")
        if bud is not None:
            f = lambda v: "-" if v is None else f"{v:,.0f} €"
            st_color = {"on_track": "green", "at_risk": "amber", "over": "red"}[bud["status"]]
            st_lbl = rt(lang, {"on_track": "b_on_track", "at_risk": "b_at_risk", "over": "b_over"}[bud["status"]])
            bline = f'{rt(lang, "h_budget")}: {st_lbl} · {rt(lang, "b_total")} {f(bud["total"])}'
            if bud["spent"] is not None:
                pc = f' ({bud["spent_pct"]}%)' if bud.get("spent_pct") is not None else ""
                bline += f' · {rt(lang, "b_spent")} {f(bud["spent"])}{pc}'
            if bud["forecast"] is not None:
                pc = f' ({bud["forecast_pct"]}%)' if bud.get("forecast_pct") is not None else ""
                bline += f' · {rt(lang, "b_forecast")} {f(bud["forecast"])}{pc}'
            if bud["status"] == "over":
                bline += f' · +{f(bud["overrun"])} ({bud["overrun_pct"]}%)'
            textbox(s, margin, Inches(0.96), Inches(12.3), Inches(0.22), bline, 10, bold=True,
                    color=rgb(_RAG_BRAND[st_color]))

        top = Inches(1.22) if bud is not None else Inches(1.15)
        if with_objectives:
            textbox(s, margin, top, Inches(12.3), Inches(0.3),
                    rt(lang, "h_otd_section"), 13, bold=True, color=B["navy"])
            objs = det.get("objectives", [])
            if objs:
                lines = []
                for o in objs[:6]:
                    rag = _status_rag(o["rag"])
                    dl = f' · {rt(lang, "deadline")} {o["target_date"]}' if o.get("target_date") else ""
                    lines.append((f'•  {o["title"]}   ({_status_label(o["rag"], lang)}{dl})',
                                  rgb(_RAG_BRAND[rag]), False))
                if len(objs) > 6:
                    lines.append((f'+{len(objs) - 6}…', B["muted"], False))
                bullets(s, margin, Inches(1.5), Inches(12.3), Inches(1.4), lines, 11)
            else:
                textbox(s, margin, Inches(1.5), Inches(12.3), Inches(0.3),
                        rt(lang, "no_obj"), 11, color=B["muted"])
            top = Inches(3.05)

        textbox(s, margin, top, Inches(12.3), Inches(0.3),
                rt(lang, "h_roadmap"), 13, bold=True, color=B["navy"])
        qtop = Emu(int(top) + int(Inches(0.42)))
        gap = Inches(0.2)
        n = 4
        total_w = prs.slide_width - margin * 2
        col_w = Emu(int((total_w - gap * (n - 1)) / n))
        # Reserve a band at the bottom for key messages when there are any.
        kms = det.get("key_messages") or []
        roadmap_bottom = int(Inches(6.35)) if kms else int(Inches(7.5))
        qh = Emu(roadmap_bottom - int(qtop) - int(Inches(0.3)))
        for i, qd in enumerate(det.get("quarters", [])):
            left = Emu(int(margin) + i * (int(col_w) + int(gap)))
            qc = rect(s, left, qtop, col_w, qh, B["card"], line=B["line"])
            items = qd.get("items", [])
            lines = []
            for it in items[:10]:
                rag = _status_rag(it["status"])
                dep = f'   · {rt(lang, "dep")} {it["dependency"]}' if it.get("dependency") else ""
                lines.append((f'•  {it["title"]}{dep}', rgb(_RAG_BRAND[rag]), False))
            if not items:
                lines.append((rt(lang, "no_jalon"), B["muted"], False))
            elif len(items) > 10:
                lines.append((f'+{len(items) - 10}…', B["muted"], False))
            # Header + milestones are the card's own text.
            paras = [(f'Q{qd["q"]} - {qd["pct"]}%', 12, B["ink"], True, PP_ALIGN.LEFT, 6)]
            paras += [(txt, 9.5, color, bold, PP_ALIGN.LEFT, 2) for (txt, color, bold) in lines]
            place(qc, paras, anchor=MSO_ANCHOR.TOP, ml=0.1, mt=0.08, mr=0.1)

        # Key messages band (success / alert / risk) along the bottom.
        if kms:
            ktop = Emu(roadmap_bottom + int(Inches(0.05)))
            textbox(s, margin, ktop, Inches(12.3), Inches(0.3),
                    rt(lang, "h_key_messages"), 12, bold=True, color=B["navy"])
            km_rag = {"success": "green", "alert": "amber", "risk": "red"}
            klines = []
            for m in kms[:4]:
                rag = km_rag.get(m["kind"], "grey")
                ts = f'   ({m["created_at"]})' if m.get("created_at") else ""
                klines.append((f'•  [{rt(lang, "km_" + m["kind"])}] {m["text"]}{ts}',
                               rgb(_RAG_BRAND[rag]), False))
            if len(kms) > 4:
                klines.append((f'+{len(kms) - 4}…', B["muted"], False))
            bullets(s, margin, Emu(int(ktop) + int(Inches(0.32))), Inches(12.3), Inches(0.7), klines, 9.5)

    # ---------------- Single-squad page slide (mirrors the squad page order) ------
    def squad_page_slide(r):
        det = r.get("detail") or {}
        s = new_slide()
        SW = prs.slide_width
        rect(s, Inches(0), Inches(0), SW, prs.slide_height, rgb("#F5F7FA"))  # app background
        L = Inches(0.4)
        FULLW = Emu(int(SW) - 2 * int(L))
        colw = Emu((int(FULLW) - int(Inches(0.2))) // 2)
        rcolx = Emu(int(L) + int(colw) + int(Inches(0.2)))

        # ----- header card (navy) -----
        hdr = rrect(s, L, Inches(0.32), FULLW, Inches(0.84), B["navy"], radius=0.08)
        place(hdr, [(r["name"], 20, B["white"], True)], anchor=MSO_ANCHOR.MIDDLE, ml=0.3, mr=4.6)
        textbox(s, Emu(int(SW) - int(Inches(4.6))), Inches(0.5), Inches(4.0), Inches(0.3),
                f'{rt(lang, "h_leader")} : {r["leader"] or "-"}', 11, color=rgb("#C7D2FE"), align=PP_ALIGN.RIGHT)

        # ----- badges row (mirrors the page header) -----
        items = [(f'{rt(lang, "h_progress_long")} {r["annual_pct"]}%', B["accent"])]
        if r["blocked"]:
            items.append((f'{r["blocked"]} {rt(lang, "h_blocked")}', B["red"]))
        if r["at_risk"]:
            items.append((f'{r["at_risk"]} {rt(lang, "h_atrisk")}', B["orange"]))
        items.append((rt(lang, "stale") if r["is_stale"] else rt(lang, "h_freshness_ok"),
                      B["orange"] if r["is_stale"] else B["green"]))
        chips_row(s, L, Inches(1.3), items)

        def list_card(x, y, w, h, title, lines, empty):
            # Title + lines are all the card's own text (one shape, no overlay).
            sh = rrect(s, x, y, w, h, B["white"], line=B["line"], radius=0.05)
            paras = [(title, 12, B["navy"], True, PP_ALIGN.LEFT, 5)]
            if lines:
                paras += [(txt, 9.5, color, bold, PP_ALIGN.LEFT, 2) for (txt, color, bold) in lines]
            else:
                paras.append((empty, 9.5, B["muted"], False, PP_ALIGN.LEFT, 0))
            place(sh, paras, anchor=MSO_ANCHOR.TOP, ml=0.18, mt=0.12, mr=0.18, mb=0.1)

        # ----- Row 1: Initiatives | OTD -----
        inits = det.get("initiatives") or []
        ilines = []
        for ini in inits[:3]:
            meta = [x for x in (ini.get("owner"),
                                (f'{rt(lang, "deadline")} {ini["deadline"]}' if ini.get("deadline") else None)) if x]
            tail = f'   ({" · ".join(meta)})' if meta else ""
            ilines.append((f'{ini["title"]}{tail}', B["ink"], False))
        if len(inits) > 3:
            ilines.append((f'+{len(inits) - 3}…', B["muted"], False))
        list_card(L, Inches(1.78), colw, Inches(1.22), rt(lang, "h_initiatives"), ilines, rt(lang, "no_initiative"))

        objs = det.get("objectives", [])
        olines = []
        for o in objs[:3]:
            rag = _status_rag(o["rag"])
            dl = f' · {rt(lang, "deadline")} {o["target_date"]}' if o.get("target_date") else ""
            olines.append((f'●  {o["title"]}   ({_status_label(o["rag"], lang)}{dl})', rgb(_RAG_BRAND[rag]), False))
        if len(objs) > 3:
            olines.append((f'+{len(objs) - 3}…', B["muted"], False))
        list_card(rcolx, Inches(1.78), colw, Inches(1.22),
                  f'{rt(lang, "h_otd_section")} {data["year"]}', olines, rt(lang, "no_obj"))

        # ----- Roadmap card with 4 quarter mini-cards + progress bars -----
        ry, rh = Inches(3.18), Inches(2.4)
        card(s, L, ry, FULLW, rh, f'{rt(lang, "h_roadmap")} {data["year"]}')
        qn, qgap = 4, Inches(0.16)
        inner_x = Emu(int(L) + int(Inches(0.18)))
        inner_w = Emu(int(FULLW) - int(Inches(0.36)))
        qw = Emu((int(inner_w) - (qn - 1) * int(qgap)) // qn)
        qy = Emu(int(ry) + int(Inches(0.52)))
        qh = Emu(int(rh) - int(Inches(0.68)))
        quarters = det.get("quarters", [])
        for i in range(qn):
            qd = quarters[i] if i < len(quarters) else {"q": i + 1, "pct": 0, "items": []}
            qx = Emu(int(inner_x) + i * (int(qw) + int(qgap)))
            qcard = rrect(s, qx, qy, qw, qh, rgb("#F8FAFC"), line=B["line"], radius=0.04)
            tx = Emu(int(qx) + int(Inches(0.1)))
            tw = Emu(int(qw) - int(Inches(0.2)))
            # Q label + percent live in the mini-card's own text frame (two runs).
            qtf = qcard.text_frame
            qtf.word_wrap = True; qtf.vertical_anchor = MSO_ANCHOR.TOP
            qtf.margin_left = Inches(0.1); qtf.margin_right = Inches(0.1)
            qtf.margin_top = Inches(0.07); qtf.margin_bottom = Inches(0.02)
            qp = qtf.paragraphs[0]
            r1 = qp.add_run(); r1.text = f'Q{qd["q"]}'
            r1.font.size = Pt(11); r1.font.bold = True; r1.font.color.rgb = B["navy"]
            r2 = qp.add_run(); r2.text = f'   {qd["pct"]}%'
            r2.font.size = Pt(10); r2.font.color.rgb = B["muted"]
            pbar(s, tx, Emu(int(qy) + int(Inches(0.42))), tw, qd["pct"], B["accent"])
            items = qd.get("items", [])
            lines = []
            for it in items[:5]:
                rag = _status_rag(it["status"])
                lines.append((f'●  {it["title"]}', rgb(_RAG_BRAND[rag]), False))
            if not items:
                lines.append((rt(lang, "no_jalon"), B["muted"], False))
            elif len(items) > 5:
                lines.append((f'+{len(items) - 5}…', B["muted"], False))
            bullets(s, tx, Emu(int(qy) + int(Inches(0.64))), tw, Emu(int(qh) - int(Inches(0.74))), lines, 8.5)

        # ----- Row 3: Key messages | Budget -----
        my, mh = Inches(5.68), Inches(1.5)
        kms = det.get("key_messages") or []
        klines = []
        for m in kms[:3]:
            rag = {"success": "green", "alert": "amber", "risk": "red"}.get(m["kind"], "grey")
            ts = f'   ({m["created_at"]})' if m.get("created_at") else ""
            klines.append((f'●  [{rt(lang, "km_" + m["kind"])}] {m["text"]}{ts}', rgb(_RAG_BRAND[rag]), False))
        if len(kms) > 3:
            klines.append((f'+{len(kms) - 3}…', B["muted"], False))
        list_card(L, my, colw, mh, rt(lang, "h_key_messages"), klines, rt(lang, "no_key_message"))

        # Budget card: title + figures are the card's own text; status stays a chip.
        bsh = rrect(s, rcolx, my, colw, mh, B["white"], line=B["line"], radius=0.05)
        place(bsh, [(rt(lang, "h_budget"), 12, B["navy"], True, PP_ALIGN.LEFT, 6)],
              anchor=MSO_ANCHOR.TOP, ml=0.18, mt=0.12, mr=0.18)
        bud = det.get("budget")
        btf = bsh.text_frame
        if bud is None:
            p = btf.add_paragraph(); rr = p.add_run(); rr.text = rt(lang, "no_budget")
            rr.font.size = Pt(9.5); rr.font.color.rgb = B["muted"]
        else:
            f = lambda v: "-" if v is None else f"{v:,.0f} €"
            st_color = {"on_track": "green", "at_risk": "amber", "over": "red"}[bud["status"]]
            st_lbl = rt(lang, {"on_track": "b_on_track", "at_risk": "b_at_risk", "over": "b_over"}[bud["status"]])
            rows = [
                (rt(lang, "b_total"), f(bud["total"])),
                (rt(lang, "b_spent"), f(bud["spent"]) + (f' · {bud["spent_pct"]}%' if bud.get("spent_pct") is not None else "")),
                (rt(lang, "b_forecast"), f(bud["forecast"]) + (f' · {bud["forecast_pct"]}%' if bud.get("forecast_pct") is not None else "")),
            ]
            for label, val in rows:
                p = btf.add_paragraph(); p.space_after = Pt(3)
                r1 = p.add_run(); r1.text = f'{label} : '
                r1.font.size = Pt(10); r1.font.color.rgb = B["muted"]
                r2 = p.add_run(); r2.text = val
                r2.font.size = Pt(10); r2.font.bold = True; r2.font.color.rgb = B["ink"]
            cw = Inches(0.26 + 0.082 * len(st_lbl))
            chip(s, Emu(int(rcolx) + int(colw) - int(cw) - int(Inches(0.16))), Emu(int(my) + int(Inches(0.12))),
                 st_lbl, rgb(_RAG_BRAND[st_color]))

    # --- Assemble the deck. A single-squad export mirrors the squad page (one
    # focused slide in page order); the multi-squad report keeps summary + grid.
    if data.get("squad_scoped"):
        for r in squads[:_MAX_DETAIL_SLIDES]:
            if r.get("detail"):
                squad_page_slide(r)
    else:
        summary_slide()
        for r in squads[:_MAX_DETAIL_SLIDES]:
            if r.get("detail"):
                detail_slide(r, with_objectives=True)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# Roadmap deck palette (mirrors the reference "Global Roadmap" deck).
_RM = {
    "dark": "#304957",   # quarter headers + swimlane labels
    "sub": "#97A3AA",    # month sub-headers
    "card": "#F2F2F2",   # milestone cards
    "card_ink": "#002060",
    "arrow": "#DCE3EA",
    "muted": "#6B7280",
    "white": "#FFFFFF",
}
_MONTHS = {
    "fr": ["Janv", "Févr", "Mars", "Avr", "Mai", "Juin", "Juil", "Août", "Sept", "Oct", "Nov", "Déc"],
    "en": ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"],
}


def render_roadmap_pptx(data: dict) -> bytes:
    """Roadmap swimlane deck (mirrors the reference layout): quarters in columns
    with month sub-headers and a timeline arrow, squads as swimlane rows, and one
    milestone card per (squad, quarter) with status-coloured bullets."""
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE = _pptx_toolkit()

    def rgb(hexstr: str) -> RGBColor:
        return RGBColor.from_string(hexstr.lstrip("#").upper())

    C = {k: rgb(v) for k, v in _RM.items()}
    STAGE = {k: rgb(v) for k, v in STAGE_COLOR.items()}  # EA=gold, GA=green
    lang = _lang(data.get("lang", "fr"))
    year = data["year"]
    gen = data["generated_at"]
    gen_str = gen.strftime("%d/%m/%Y %H:%M") if isinstance(gen, datetime) else str(gen)
    months = _MONTHS[lang]

    SLIDE_W, SLIDE_H = 13.333, 7.5
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # Column geometry (4 quarters).
    MARGIN, GAP = 0.5, 0.08
    COL_W = (SLIDE_W - 2 * MARGIN - GAP * 3) / 4
    def col_x(i): return MARGIN + i * (COL_W + GAP)
    Y_Q, H_Q = 0.92, 0.34          # quarter header
    Y_M, H_M = 1.30, 0.28          # month sub-headers
    Y_ARROW, H_ARROW = 1.66, 0.30  # timeline arrow
    Y_TOP, Y_BOTTOM = 2.10, 7.18   # swimlane content band

    squads = [r for blk in data["tribes"] for r in blk["squads"]]

    def shape(s, kind, x, y, w, h, fill, *, line=None, rot=0, round_adj=None):
        sh = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line; sh.line.width = Pt(0.75)
        sh.shadow.inherit = False
        if rot:
            sh.rotation = rot
        if round_adj is not None:
            try:
                sh.adjustments[0] = round_adj
            except Exception:
                pass
        return sh

    def set_text(holder, text, size, color, *, bold=False, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE):
        tf = holder.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.04)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

    def textbox(s, x, y, w, h, runs, size, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]; p.alignment = align
        for txt, color, bold in runs:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        return box

    def draw_header(s):
        textbox(s, MARGIN, 0.22, 8.6, 0.5,
                [(f'{rt(lang, "roadmap_report")} · {data["scope_name"]}', C["dark"], True)], 22)
        textbox(s, MARGIN, 0.66, 8.6, 0.25,
                [(f'{rt(lang, "year")} {year} · {rt(lang, "generated_full", d=gen_str)}', C["muted"], False)], 10.5)
        legend = [("EA  ", STAGE["EA"], True), (rt(lang, "stage_ea") + "      ", C["dark"], False),
                  ("GA  ", STAGE["GA"], True), (rt(lang, "stage_ga"), C["dark"], False)]
        textbox(s, SLIDE_W - 5.9, 0.34, 5.4, 0.3, legend, 10, align=PP_ALIGN.RIGHT)
        # timeline arrow spanning the columns
        shape(s, MSO_SHAPE.RIGHT_ARROW, MARGIN - 0.1, Y_ARROW,
              (col_x(3) + COL_W) - (MARGIN - 0.1) + 0.18, H_ARROW, C["arrow"], round_adj=None)
        for i, q in enumerate((1, 2, 3, 4)):
            x = col_x(i)
            set_text(shape(s, MSO_SHAPE.RECTANGLE, x, Y_Q, COL_W, H_Q, C["dark"]),
                     f'Q{q} {year}', 15, C["white"], bold=True)
            mw = (COL_W - 2 * 0.05) / 3
            for mi in range(3):
                mx = x + mi * (mw + 0.05)
                set_text(shape(s, MSO_SHAPE.RECTANGLE, mx, Y_M, mw, H_M, C["sub"]),
                         months[i * 3 + mi], 10, C["white"])

    def draw_card(s, x, y, w, h, items, fs, line_h):
        card = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, C["card"], round_adj=0.06)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.07); tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.03)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        if not items:
            return
        # Flatten into ordered paragraphs: a bold theme header then its milestones.
        specs: list[tuple[str, object]] = []
        for theme, group in group_by_theme(items):
            if theme:
                specs.append(("theme", theme))
            for it in group:
                specs.append(("item", it))
        item_fs = max(7, fs - 1.5)  # milestone lines a touch smaller than the theme header
        max_lines = max(1, int((h - 0.08) / line_h))
        shown = specs[:max_lines]
        for li, (kind, val) in enumerate(shown):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT  # themes + milestones are left-aligned, never centred
            p.space_after = Pt(0.5)
            if kind == "theme":
                r = p.add_run(); r.text = val
                r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = C["card_ink"]
                continue
            it = val
            r1 = p.add_run(); r1.text = it["title"]
            r1.font.size = Pt(item_fs); r1.font.color.rgb = C["card_ink"]
            stage = it.get("stage")
            if stage:
                ro = p.add_run(); ro.text = " ("
                ro.font.size = Pt(item_fs); ro.font.color.rgb = C["card_ink"]
                rs = p.add_run(); rs.text = stage
                rs.font.size = Pt(item_fs); rs.font.bold = True
                rs.font.color.rgb = STAGE.get(stage, C["card_ink"])
                rc = p.add_run(); rc.text = ")"
                rc.font.size = Pt(item_fs); rc.font.color.rgb = C["card_ink"]
        if len(specs) > max_lines:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = f'+{len(specs) - max_lines}…'
            r.font.size = Pt(max(6, item_fs - 0.5)); r.font.color.rgb = C["muted"]

    def draw_swimlanes(s, lanes):
        # Everything fits on ONE slide: band height + fonts scale with the count,
        # but with a readable floor (small decks get a comfortably large font).
        n = max(1, len(lanes))
        band_h = (Y_BOTTOM - Y_TOP) / n
        card_fs = (13 if band_h >= 1.05 else 12 if band_h >= 0.85 else 11 if band_h >= 0.68
                   else 10 if band_h >= 0.52 else 9 if band_h >= 0.40 else 8)
        line_h = card_fs * 0.020
        label_fs = max(8.5, min(13, band_h * 12))
        lbl_h = min(0.36, band_h * 0.6)
        for ri, sq in enumerate(lanes):
            by = Y_TOP + ri * band_h
            bcy = by + band_h / 2
            lbl_len = max(0.35, min(band_h - 0.08, 1.7))
            lbl = shape(s, MSO_SHAPE.RECTANGLE, 0.24 - lbl_len / 2, bcy - lbl_h / 2,
                        lbl_len, lbl_h, C["dark"], rot=270)
            set_text(lbl, sq["name"], label_fs, C["white"], bold=True)
            qmap = {qd["q"]: qd["items"] for qd in (sq.get("detail") or {}).get("quarters", [])}
            for i, q in enumerate((1, 2, 3, 4)):
                draw_card(s, col_x(i), by + 0.04, COL_W, band_h - 0.08, qmap.get(q, []), card_fs, line_h)

    s = prs.slides.add_slide(prs.slide_layouts[6])  # single page, always
    draw_header(s)
    draw_swimlanes(s, squads)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ----- Initiatives list (flat: initiative / owner / squad / deadline) -------------

_INIT_T = {
    "fr": {"title": "Initiatives", "h_init": "Initiative", "h_owner": "Owner",
           "h_squad": "Squad", "h_deadline": "Échéance", "none": "Aucune initiative", "year": "Année"},
    "en": {"title": "Initiatives", "h_init": "Initiative", "h_owner": "Owner",
           "h_squad": "Squad", "h_deadline": "Deadline", "none": "No initiative", "year": "Year"},
}


def _deadline_str(d) -> str:
    if not d:
        return "-"
    if isinstance(d, datetime):
        return d.date().isoformat()
    return str(d)[:10]


def build_initiative_list(db: Session, scope_tribe: int | None, year: int) -> dict:
    """Flat list of initiatives in scope: title, owner, squad, deadline."""
    from .models import Initiative, Tribe
    tribes = {t.id: t.name for t in db.scalars(select(Tribe)).all()}
    q = (select(Initiative).where(Initiative.year == year)
         .order_by(Initiative.display_order, Initiative.id))
    if scope_tribe is not None:
        q = q.where(Initiative.tribe_id == scope_tribe)
    items = [{"id": i.id, "title": i.title, "owner": i.owner,
              "squad_name": i.squad.name if i.squad else None,
              "deadline": _deadline_str(i.deadline)} for i in db.scalars(q).all()]
    scope_name = tribes.get(scope_tribe) if scope_tribe is not None else "Toutes les tribus"
    return {"year": year, "scope_name": scope_name or "-", "items": items}


def render_initiatives_html(data: dict, *, lang: str = "fr", standalone: bool = True) -> str:
    e = html.escape
    lang = _lang(lang)
    T = _INIT_T[lang]
    parts = [f'<div class="hdr"><h1>{e(T["title"])} - {e(data["scope_name"])}</h1>',
             f'<div class="sub">{e(T["year"])} {data["year"]}</div></div>']
    if not data["items"]:
        parts.append(f'<div class="muted small">{e(T["none"])}</div>')
    else:
        parts.append('<table><thead><tr>'
                     f'<th>{e(T["h_init"])}</th><th>{e(T["h_owner"])}</th>'
                     f'<th>{e(T["h_squad"])}</th><th>{e(T["h_deadline"])}</th></tr></thead><tbody>')
        for it in data["items"]:
            parts.append(f'<tr><td><strong>{e(it["title"])}</strong></td>'
                         f'<td>{e(it["owner"] or "-")}</td>'
                         f'<td>{e(it["squad_name"] or "-")}</td>'
                         f'<td>{e(it["deadline"])}</td></tr>')
        parts.append('</tbody></table>')
    body = "\n".join(parts)
    if not standalone:
        return f'<div class="tc-report">{_CSS}{body}</div>'
    return (f'<!doctype html><html lang="{e(lang)}"><head><meta charset="utf-8">'
            f'<title>{e(T["title"])}</title>{_CSS}</head>'
            f'<body><div class="tc-report">{body}</div></body></html>')


def render_initiatives_pptx(data: dict, *, lang: str = "fr") -> bytes:
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE = _pptx_toolkit()

    def rgb(h):
        return RGBColor.from_string(h.lstrip("#").upper())

    B = {k: rgb(v) for k, v in _BRAND.items()}
    lang = _lang(lang)
    T = _INIT_T[lang]
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    margin = Inches(0.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.92))
    head.fill.solid(); head.fill.fore_color.rgb = B["navy"]; head.line.fill.background(); head.shadow.inherit = False
    tf = head.text_frame; tf.margin_left = Inches(0.5); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; r = p.add_run(); r.text = f'{T["title"]} · {data["scope_name"]}'
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = B["white"]

    headers = [T["h_init"], T["h_owner"], T["h_squad"], T["h_deadline"]]
    wfrac = [0.40, 0.24, 0.22, 0.14]
    items = data["items"][:18]
    nrows = max(2, len(items) + 1)
    table_w = int(prs.slide_width - margin * 2)
    tbl = s.shapes.add_table(nrows, 4, margin, Inches(1.2), Emu(table_w),
                             Inches(0.34) + Inches(0.3) * (nrows - 1)).table
    for ci, f in enumerate(wfrac):
        tbl.columns[ci].width = Emu(int(table_w * f))

    def cell(c, text, size, color, *, bold=False, align=PP_ALIGN.LEFT, fill=None):
        c.margin_left = Inches(0.06); c.margin_right = Inches(0.04); c.margin_top = Emu(0); c.margin_bottom = Emu(0)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fill is not None:
            c.fill.solid(); c.fill.fore_color.rgb = fill
        else:
            c.fill.background()
        c.text = text or " "
        pp = c.text_frame.paragraphs[0]; pp.alignment = align
        if pp.runs:
            rr = pp.runs[0]; rr.font.size = Pt(size); rr.font.bold = bold; rr.font.color.rgb = color

    for ci, h in enumerate(headers):
        cell(tbl.cell(0, ci), h, 11, B["white"], bold=True, fill=B["navy"])
    if not items:
        cell(tbl.cell(1, 0), T["none"], 10, B["muted"])
        for ci in range(1, 4):
            cell(tbl.cell(1, ci), " ", 10, B["muted"])
    for ri, it in enumerate(items, start=1):
        zebra = B["zebra"] if ri % 2 == 0 else B["white"]
        cells = [(it["title"], B["ink"], True), (it["owner"] or "-", B["ink"], False),
                 (it["squad_name"] or "-", B["ink"], False), (it["deadline"], B["ink"], False)]
        for ci, (val, color, bold) in enumerate(cells):
            cell(tbl.cell(ri, ci), val, 10, color, bold=bold, fill=zebra)
    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()



# ----- Automatic weekly send ------------------------------------------------------

def _iso_week_key(dt: datetime) -> str:
    y, w, _ = dt.isocalendar()
    return f"{y}-W{w:02d}"


def send_due_weekly_reports(db: Session, now: datetime | None = None) -> int:
    """Send the weekly report if today/now matches the configured schedule.

    Idempotent within an ISO week (guarded by last_sent_week). Returns the number
    of emails sent. Safe to call repeatedly from the scheduler.
    """
    from .reportconfig import get_report, set_report
    from .smtpconfig import get_smtp
    from .mail import send_email
    from .models import User
    from .modulesconfig import get_modules, is_active

    now = now or utcnow()
    if not is_active(get_modules(db), "review", "weekly_report"):
        return 0
    cfg = get_report(db)
    if not cfg.get("enabled"):
        return 0
    smtp = get_smtp(db)
    if not smtp.get("enabled"):
        return 0
    if now.weekday() not in cfg["weekdays"] or now.hour < cfg["hour"]:
        return 0
    today = now.date().isoformat()
    if cfg.get("last_sent_day") == today:
        return 0

    since = cfg.get("since_days", 7)
    year = st.current_year_quarter(now)[0]
    lang = _lang(get_general(db).get("default_lang"))

    only_changes = bool(cfg.get("only_when_changes"))
    week = now.isocalendar()[1]
    sent = 0
    prepared: dict[str, dict] = {}

    def prepare(scope: int | None, scope_key: str, scope_label: str) -> dict:
        """Build the report for a scope once: data, changelog (vs baseline),
        HTML (with the "what's new" encart), PPTX and a prefixed subject."""
        if scope_key not in prepared:
            data = build_report_data(db, scope, year, since, now, lang=lang)
            sig = report_signature(data)
            changes = diff_report(get_baseline(db, scope_key), sig, lang)
            html_body = render_html(data, standalone=True, changes=changes)
            try:
                pptx_bytes = render_pptx(data)
            except Exception:
                pptx_bytes = b""
            subject = subject_prefix(changes, lang) + rt(lang, "subject", scope=scope_label, w=week)
            prepared[scope_key] = {"sig": sig, "changes": changes, "html": html_body,
                                   "pptx": pptx_bytes, "subject": subject}
        return prepared[scope_key]

    def attachment_of(p: dict):
        if p["pptx"] and cfg.get("attach_pptx", True):
            return (f"rapport_hebdo_{today}.pptx", p["pptx"],
                    "application", "vnd.openxmlformats-officedocument.presentationml.presentation")
        return None

    def worth_sending(p: dict) -> bool:
        c = p["changes"]
        # Always send the very first report (establishes the baseline); otherwise
        # honour the "only when changes" policy.
        return c.get("first") or not (only_changes and c["count"] == 0)

    # Fixed recipient list → global report.
    recipients = list(dict.fromkeys(a for a in (cfg.get("recipients") or []) if a))
    if recipients:
        p = prepare(None, "global", rt(lang, "all_tribes"))
        if worth_sending(p):
            att = attachment_of(p)
            for addr in recipients:
                if send_email(smtp, addr, p["subject"], p["html"], attachment=att, html=True):
                    sent += 1
        set_baseline(db, "global", p["sig"])

    # Optional: each tribe leader also receives their OWN tribe-scoped report,
    # with that tribe's squad leaders in CC.
    if cfg.get("tribe_leader_digest"):
        for tribe in db.scalars(select(Tribe).order_by(Tribe.display_order, Tribe.id)).all():
            leaders = [u for u in db.scalars(
                select(User).where(User.role == "tribe_leader", User.tribe_id == tribe.id)).all()
                if (u.email or "").strip() and u.status == "active"]
            if not leaders:
                continue
            to = ", ".join(dict.fromkeys(l.email.strip() for l in leaders))
            leader_emails = {l.email.strip().lower() for l in leaders}
            sq_leader_ids = [s.leader_user_id for s in db.scalars(
                select(Squad).where(Squad.tribe_id == tribe.id)).all() if s.leader_user_id]
            cc, seen_cc = [], set()
            if sq_leader_ids:
                for u in db.scalars(select(User).where(User.id.in_(sq_leader_ids))).all():
                    e = (u.email or "").strip()
                    el = e.lower()
                    if e and u.status == "active" and el not in leader_emails and el not in seen_cc:
                        seen_cc.add(el)
                        cc.append(e)
            scope_key = f"tribe:{tribe.id}"
            p = prepare(tribe.id, scope_key, tribe.name)
            if worth_sending(p):
                if send_email(smtp, to, p["subject"], p["html"],
                              attachment=attachment_of(p), html=True, cc=cc):
                    sent += 1
            set_baseline(db, scope_key, p["sig"])

    cfg["last_sent_day"] = today
    set_report(db, cfg)
    db.commit()
    return sent


def send_personal_subscriptions(db: Session, now: datetime | None = None) -> int:
    """Send the report to each subscription (global or per-squad) that is due.

    A global subscription (squad_id NULL) follows the user's visibility (admin →
    all tribes, others → their tribe); a per-squad subscription targets that squad.
    Returns the number of emails sent. Safe to call repeatedly from the scheduler.
    """
    from .smtpconfig import get_smtp
    from .mail import send_email
    from .models import ReportSubscription, Squad, Tribe, User
    from .modulesconfig import get_modules, is_active

    now = now or utcnow()
    if not is_active(get_modules(db), "review", "weekly_report"):
        return 0
    smtp = get_smtp(db)
    if not smtp.get("enabled"):
        return 0

    from .reportconfig import get_report
    year = st.current_year_quarter(now)[0]
    lang = _lang(get_general(db).get("default_lang"))
    only_changes = bool(get_report(db).get("only_when_changes"))
    # Cache report data + PPTX per (scope_tribe, squad_id, since); HTML is rendered
    # per subscription because the "what's new" encart is per-recipient baseline.
    rendered: dict[tuple, tuple[dict, bytes]] = {}

    def render(scope_tribe: int | None, squad_id: int | None, since: int) -> tuple[dict, bytes]:
        key = (scope_tribe, squad_id, since)
        if key not in rendered:
            data = build_report_data(db, scope_tribe, year, since, now, squad_id=squad_id, lang=lang)
            try:
                pptx_bytes = render_pptx(data)
            except Exception:
                pptx_bytes = b""
            rendered[key] = (data, pptx_bytes)
        return rendered[key]

    sent = 0
    for sub in db.scalars(select(ReportSubscription)).all():
        wd = sub.weekdays or []
        if not wd and sub.interval_days <= 0:
            continue  # inactive subscription
        user = db.get(User, sub.user_id)
        if user is None or not user.email:
            continue
        last = _aware(sub.last_sent_at)
        if wd:
            # Weekday schedule: fire on a chosen day, past the hour, once per day.
            if now.weekday() not in wd or now.hour < sub.hour:
                continue
            if last is not None and last.date() == now.date():
                continue
            since = 7
        else:
            # Legacy "every N days" cadence.
            if last is not None and (now - last) < timedelta(days=sub.interval_days):
                continue
            since = max(sub.interval_days, 7)
        if sub.squad_id is not None:
            data, pptx_bytes = render(None, sub.squad_id, since)
        else:
            scope_tribe = None if user.role == "admin" else user.tribe_id
            data, pptx_bytes = render(scope_tribe, None, since)

        scope_key = f"sub:{sub.id}"
        sig = report_signature(data)
        changes = diff_report(get_baseline(db, scope_key), sig, lang)

        def _mark_done():
            set_baseline(db, scope_key, sig)
            sub.last_sent_at = now
            if sub.squad_id is None:
                user.report_last_sent_at = now

        # "Only when changes": skip the email but still advance the cadence/baseline.
        if only_changes and not changes.get("first") and changes["count"] == 0:
            _mark_done()
            continue

        html_body = render_html(data, standalone=True, changes=changes)
        attachment = None
        if pptx_bytes:
            attachment = (f"rapport_{now.date().isoformat()}.pptx", pptx_bytes,
                          "application", "vnd.openxmlformats-officedocument.presentationml.presentation")
        if sub.squad_id is not None:
            sq = db.get(Squad, sub.squad_id)
            scope_lbl = sq.name if sq else rt(lang, "h_squad")
        elif user.role == "admin":
            scope_lbl = rt(lang, "all_tribes")
        else:
            tr = db.get(Tribe, user.tribe_id) if user.tribe_id else None
            scope_lbl = tr.name if tr else rt(lang, "all_tribes")
        subject = subject_prefix(changes, lang) + rt(lang, "subject_personal", scope=scope_lbl, n=sub.interval_days)
        if send_email(smtp, user.email, subject, html_body, attachment=attachment, html=True):
            _mark_done()
            sent += 1
    db.commit()  # persist baselines / last_sent even when only skips occurred
    return sent
