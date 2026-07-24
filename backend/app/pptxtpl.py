"""Optional PPTX export template (Admin > Settings).

An admin uploads a ``.pptx`` once; from then on **every** PPTX export is built on top
of it, so the decks inherit the organisation's master slides, theme, colours, fonts and
branding (e.g. your corporate template). Off by default -> python-pptx's plain white deck.

Storage: the file (a few MB) is kept in ``app_settings`` as base64 under the ``pptx_template``
key, alongside light metadata (filename, size, who/when). The **active** template for the
current request is carried in a ``ContextVar`` so the pure renderers in ``report.py`` /
``routers/steerco.py`` / ``orgrender.py`` stay database-free and testable: an export
endpoint calls :func:`use` once with the stored bytes, the renderers call
:func:`new_presentation` and :func:`blank_layout`.
"""
from __future__ import annotations

import base64
import io
import json
from contextvars import ContextVar
from datetime import datetime, timezone

_KEY = "pptx_template"
# The template bytes to build the current request's decks on (None = default deck).
_current: ContextVar[bytes | None] = ContextVar("pptx_template_current", default=None)


# --------------------------------------------------------------------------
# Persistence (app_settings)
# --------------------------------------------------------------------------
def _row(db):
    from .models import AppSetting
    return db.get(AppSetting, _KEY)


def meta(db) -> dict:
    """Public status of the stored template (never the bytes): presence + light info."""
    row = _row(db)
    if not row or not row.value:
        return {"present": False}
    try:
        d = json.loads(row.value)
    except (ValueError, TypeError):
        return {"present": False}
    return {"present": True, "filename": d.get("filename"), "size": d.get("size"),
            "uploaded_at": d.get("uploaded_at"), "uploaded_by": d.get("uploaded_by")}


def get(db) -> bytes | None:
    """The stored template bytes, or None when none is configured."""
    row = _row(db)
    if not row or not row.value:
        return None
    try:
        return base64.b64decode(json.loads(row.value)["data_b64"])
    except (ValueError, TypeError, KeyError):
        return None


def validate(data: bytes) -> None:
    """Raise ValueError unless ``data`` is a real, openable ``.pptx``."""
    from pptx import Presentation
    try:
        Presentation(io.BytesIO(data))
    except Exception:
        raise ValueError("Fichier invalide : ce n'est pas un .pptx exploitable.")


def save(db, filename: str, data: bytes, uploaded_by: str | None = None) -> dict:
    """Validate then store the template (replacing any previous one). Returns its meta."""
    from .models import AppSetting
    validate(data)
    payload = {
        "filename": filename, "data_b64": base64.b64encode(data).decode("ascii"),
        "size": len(data), "uploaded_at": datetime.now(timezone.utc).isoformat(),
        "uploaded_by": uploaded_by,
    }
    row = _row(db)
    if row is None:
        db.add(AppSetting(key=_KEY, value=json.dumps(payload)))
    else:
        row.value = json.dumps(payload)
    db.flush()
    return {k: payload[k] for k in ("filename", "size", "uploaded_at", "uploaded_by")}


def clear(db) -> None:
    """Remove the template; exports fall back to the default deck."""
    row = _row(db)
    if row is not None:
        db.delete(row)
        db.flush()


# --------------------------------------------------------------------------
# Render-time API (used by the pure PPTX renderers)
# --------------------------------------------------------------------------
def use(template: bytes | None) -> None:
    """Set the template the current request's decks are built on (None = default)."""
    _current.set(template)


def new_presentation():
    """A fresh ``Presentation`` on the active template - its masters, layouts, theme and
    branding kept, its sample slides stripped so we start empty - or python-pptx's default
    deck when no template is configured."""
    from pptx import Presentation
    template = _current.get()
    if not template:
        return Presentation()
    prs = Presentation(io.BytesIO(template))
    # Drop the template's own example slides; keep everything else (masters/layouts/theme).
    # Removing the sldId reference alone leaves the slide PART orphaned in the package, and
    # it then collides with the first new slide ("Duplicate name slide1.xml") on save - so
    # drop the relationship too, which makes the orphaned part unreachable and unsaved.
    _RID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    sldIdLst = prs.slides._sldIdLst
    prs_part = prs.part
    for sldId in list(sldIdLst):
        rId = sldId.get(_RID)
        if rId:
            try:
                prs_part.drop_rel(rId)
            except KeyError:  # pragma: no cover - already gone
                pass
        sldIdLst.remove(sldId)
    return prs


def _shows_master(layout) -> bool:
    """False when the layout hides the master's own shapes (``showMasterSp="0"``). Many
    templates set this on their "blank" layout, which then suppresses the master footer
    band and logo - exactly the branding an export is supposed to show."""
    return layout.element.get("showMasterSp") != "0"


def blank_layout(prs):
    """The best slide layout to draw a full custom deck on: one that still **shows the
    master's shapes** (so the template's footer band and logo appear), then a
    ``blank``-named one, then the fewest own decorations and placeholders. Picking the
    absolute emptiest layout is wrong when that layout is the one with ``showMasterSp=0``
    (it would hide all the branding). Placeholders are stripped per-slide in
    :func:`add_slide`, so they barely weigh here."""
    layouts = list(prs.slide_layouts)
    if not layouts:  # pragma: no cover - every real deck has layouts
        return prs.slide_masters[0].slide_layouts[0]
    branding = any(not sh.is_placeholder for sh in prs.slide_masters[0].shapes)

    def score(lay):
        hides_branding = 1 if (branding and not _shows_master(lay)) else 0
        not_blank = 0 if "blank" in (lay.name or "").lower() else 1
        own_decorations = sum(1 for sh in lay.shapes if not sh.is_placeholder)
        return (hides_branding, not_blank, own_decorations, len(list(lay.placeholders)))

    return min(layouts, key=score)


def add_slide(prs):
    """Add a slide on :func:`blank_layout` and remove the layout's inherited placeholders
    so no "click to add" prompts show. The template master's own branding (footer, logo)
    keeps rendering - it comes from the master, not from these placeholders."""
    slide = prs.slides.add_slide(blank_layout(prs))
    for ph in list(slide.placeholders):
        el = ph._element
        el.getparent().remove(el)
    return slide
