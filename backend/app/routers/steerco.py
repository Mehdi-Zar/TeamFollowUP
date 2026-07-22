"""Steerco (steering committee) - monthly squad snapshots + consolidated one-pager.

A squad opts in (``Squad.steerco_enabled``, self-service) and each month reports a
small **snapshot** for the current month (KPI counts, the month's SLA per COTS, the
month's incident count, plus events). Snapshots accumulate one per (squad, period),
so the 12-month charts and the "last 12 months" SLA row are **computed** from the
last 12 snapshots - no need to re-enter history every time. For the first report a
backfill endpoint seeds past months in one shot (grid / paste-from-Excel in the UI).

Leadership reads/export a KPI one-pager per squad (HTML / PPTX), rendered in the
requested language (default English). Gated by the optional ``steerco`` module.

Stored snapshot shape (``SteercoEntry.data``), see frontend/src/steerco.ts:
    {"kpis":[{label,value}], "sla":{"services":[...],"cells":[{v}]},
     "incidents": <number>, "last_events":[...], "next_events":[...]}

Only raw values are entered. The KPI variation vs M-1 (``trend`` / ``delta``) and the
SLA colour (``s``) are recomputed at render time from the numbers themselves, see
``_kpi_change`` and ``_sla_status``.
"""
import io
import math
from html import escape

from fastapi import APIRouter, Body, Depends, HTTPException, Query
from fastapi.responses import HTMLResponse, Response
from sqlalchemy.orm import Session

from ..database import get_db
from ..deps import (assert_can_edit_squad, get_current_user, record_audit,
                    require_module, require_tribe_or_admin, require_writer,
                    visible_tribe_id)
from ..models import SteercoEntry, Squad, User

router = APIRouter(prefix="/api/steerco", tags=["steerco"],
                   dependencies=[Depends(require_module("steerco"))])

NAVY = "#0B2545"
SLA_CLASS = {"ok": "b-ok", "warn": "b-warn", "ko": "b-ko"}
# Event severity (entered on each event, in the wizard or the Excel): tints the type
# chip so a critical event reads at a glance. Same palette as the SLA statuses.
SEV_BG = {"red": "#F7E0E0", "amber": "#FBF0D9", "green": "#E4F3EA", "ice": "#E8F0FE"}
SEV_FG = {"red": "#B42318", "amber": "#9C7212", "green": "#027A48", "ice": "#13315C"}
TREND_CLASS = {"up": "up", "down": "down", "flat": "flat"}
TREND_ARROW = {"up": "▲", "down": "▼", "flat": "▬"}
SERIES_COLORS = ["#0B2545", "#2E9E5B", "#E0A526", "#8DA9C4", "#13315C", "#D24545"]

# Fixed labels of the one-pager, per language (default English). Data (KPI labels,
# service names, event text) is user-entered and rendered as-is.
I18N = {
    "en": {
        "months": ["January", "February", "March", "April", "May", "June", "July",
                   "August", "September", "October", "November", "December"],
        "key_figures": "KPI",
        "sla": "SLA", "sla_sub": "incidents &amp; SwF (by COTS)",
        "kpi_chart": "KPI trend", "kpi_chart_sub": "12 months (base 100)",
        "inc_chart": "Incidents", "inc_chart_sub": "over 12 months",
        "period": "Period", "current": "Current month", "trailing": "Last 12 months",
        "last_events": "Last events", "next_events": "Next events",
        "no_kpi": "No KPI yet.", "no_sla": "No SLA data.", "no_data": "No data.",
        "no_event": "No event.", "no_squads": "No Steerco-enabled squad in your scope.",
    },
    "fr": {
        "months": ["janvier", "février", "mars", "avril", "mai", "juin", "juillet",
                   "août", "septembre", "octobre", "novembre", "décembre"],
        "key_figures": "KPI",
        "sla": "SLA", "sla_sub": "incidents &amp; SwF (par COTS)",
        "kpi_chart": "Évolution KPI", "kpi_chart_sub": "12 mois (base 100)",
        "inc_chart": "Incidents", "inc_chart_sub": "sur 12 mois",
        "period": "Période", "current": "Mois en cours", "trailing": "12 derniers mois",
        "last_events": "Derniers évènements", "next_events": "Prochains évènements",
        "no_kpi": "Aucun KPI renseigné.", "no_sla": "Aucune donnée SLA.", "no_data": "Aucune donnée.",
        "no_event": "Aucun évènement.", "no_squads": "Aucune squad Steerco activée dans votre périmètre.",
    },
}


def _lang(v: str | None) -> str:
    return "fr" if (v or "").lower().startswith("fr") else "en"


# --------------------------------------------------------------------------
# Opt-in flag (self-service by the squad leader)
# --------------------------------------------------------------------------

@router.put("/squad/{squad_id}/enabled")
def set_squad_enabled(squad_id: int, enabled: bool = Body(..., embed=True),
                      db: Session = Depends(get_db), user: User = Depends(require_writer)):
    """Turn Steerco reporting on/off for a squad (the squad leader's own toggle)."""
    squad = db.get(Squad, squad_id)
    if squad is None:
        raise HTTPException(status_code=404, detail="Squad introuvable")
    assert_can_edit_squad(db, user, squad_id)
    squad.steerco_enabled = bool(enabled)
    record_audit(db, user.id, "steerco.enabled", entity="squad", entity_id=squad_id,
                 detail={"enabled": squad.steerco_enabled})
    db.commit()
    return {"squad_id": squad_id, "steerco_enabled": squad.steerco_enabled}


# --------------------------------------------------------------------------
# Data access (one monthly snapshot per squad+period)
# --------------------------------------------------------------------------

def _squads_in_scope(db: Session, user: User) -> list[Squad]:
    q = db.query(Squad)
    tid = visible_tribe_id(user)
    if tid is not None:
        q = q.filter(Squad.tribe_id == tid)
    return q.order_by(Squad.name).all()


def _clamp_pct(v):
    """SLA cells are percentages: keep them inside 0 to 100 (a typo like 994 -> 100).
    Non-numeric text is returned untouched."""
    n = _num(v)
    if n is None or 0 <= n <= 100:
        return v
    c = min(max(n, 0.0), 100.0)
    return f"{c:.0f}%" if str(v).strip().endswith("%") else f"{c:.0f}"


def _sanitized(payload: dict) -> dict:
    """A snapshot as it is stored: SLA percentages capped at 100."""
    sla = (payload or {}).get("sla")
    if not isinstance(sla, dict) or not isinstance(sla.get("cells"), list):
        return payload or {}
    cells = [({**c, "v": _clamp_pct(c.get("v"))} if isinstance(c, dict) else c)
             for c in sla["cells"]]
    return {**payload, "sla": {**sla, "cells": cells}}


def _entry(db: Session, squad_id: int, period: str) -> SteercoEntry | None:
    return (db.query(SteercoEntry)
            .filter(SteercoEntry.squad_id == squad_id, SteercoEntry.period == period)
            .one_or_none())


@router.get("/entries")
def list_entries(period: str = Query(...), db: Session = Depends(get_db),
                 user: User = Depends(require_tribe_or_admin)):
    """Every steerco-enabled in-scope squad's snapshot for a period (drives the tab).

    Leadership-only, like the one-pager endpoints it feeds: it returns every squad's
    figures at once, which a squad member has no business reading."""
    squads = [s for s in _squads_in_scope(db, user) if s.steerco_enabled]
    by_squad = {e.squad_id: e for e in
                db.query(SteercoEntry).filter(SteercoEntry.period == period).all()}
    return [{
        "squad_id": s.id, "squad_name": s.name, "tribe_id": s.tribe_id,
        "data": (by_squad[s.id].data if s.id in by_squad else {}),
        "filled": s.id in by_squad,
        "updated_at": (by_squad[s.id].updated_at if s.id in by_squad else None),
    } for s in squads]


@router.get("/squad/{squad_id}")
def get_squad_entry(squad_id: int, period: str = Query(...), db: Session = Depends(get_db),
                    user: User = Depends(get_current_user)):
    """One squad's monthly snapshot for a period ({} when not yet filled).

    Also returns whether this month is already filled and when/who last updated it,
    so the reporting screen can show the monthly-cadence status at a glance."""
    if db.get(Squad, squad_id) is None:
        raise HTTPException(status_code=404, detail="Squad introuvable")
    e = _entry(db, squad_id, period)
    return {
        "squad_id": squad_id, "period": period,
        "data": (e.data if e else {}),
        "filled": bool(e and e.data),
        "updated_at": e.updated_at.isoformat() if (e and e.updated_at) else None,
        "updated_by": (e.updated_by.display_name if (e and e.updated_by) else None),
    }


@router.put("/squad/{squad_id}")
def upsert_squad_entry(squad_id: int, period: str = Query(...),
                       data: dict = Body(default=None), db: Session = Depends(get_db),
                       user: User = Depends(require_writer)):
    """Create/update a squad's monthly snapshot for a period (writer + edit rights)."""
    squad = db.get(Squad, squad_id)
    if squad is None:
        raise HTTPException(status_code=404, detail="Squad introuvable")
    assert_can_edit_squad(db, user, squad_id)
    payload = _sanitized(data or {})
    e = _entry(db, squad_id, period)
    if e is None:
        e = SteercoEntry(squad_id=squad_id, period=period, data=payload, updated_by_user_id=user.id)
        db.add(e)
    else:
        e.data = payload
        e.updated_by_user_id = user.id
    db.flush()
    record_audit(db, user.id, "steerco.upsert", entity="steerco", entity_id=e.id,
                 detail={"squad_id": squad_id, "period": period})
    db.commit()
    return {"squad_id": squad_id, "period": period, "data": e.data}


# --------------------------------------------------------------------------
# 12-month history (backfill + read), for the auto-built charts
# --------------------------------------------------------------------------

def month_keys(period: str, n: int = 12) -> list[str]:
    """The n month keys ("YYYY-MM") ending at ``period`` (oldest first)."""
    y, m = (int(x) for x in period.split("-")[:2])
    keys = []
    for i in range(n - 1, -1, -1):
        mm, yy = m - i, y
        while mm <= 0:
            mm += 12
            yy -= 1
        keys.append(f"{yy:04d}-{mm:02d}")
    return keys


def _month_short(key: str) -> str:
    y, m = key.split("-")[:2]
    return f"{m}/{y[2:]}"


def _period_long(period: str, L: dict) -> str:
    """"2026-07" -> "July 2026" / "juillet 2026" (spelled-out month, per language)."""
    try:
        y, m = period.split("-")[:2]
        return f"{L['months'][int(m) - 1]} {y}"
    except (ValueError, IndexError, KeyError):
        return period


@router.get("/squad/{squad_id}/history")
def get_history(squad_id: int, period: str = Query(...), db: Session = Depends(get_db),
                user: User = Depends(get_current_user)):
    """The 12 monthly snapshots ending at ``period`` (for the backfill grid)."""
    if db.get(Squad, squad_id) is None:
        raise HTTPException(status_code=404, detail="Squad introuvable")
    keys = month_keys(period, 12)
    by_period = {e.period: (e.data or {}) for e in db.query(SteercoEntry)
                 .filter(SteercoEntry.squad_id == squad_id, SteercoEntry.period.in_(keys)).all()}
    return {"squad_id": squad_id, "period": period,
            "months": [{"period": k, "data": by_period.get(k, {})} for k in keys]}


@router.put("/squad/{squad_id}/history")
def upsert_history(squad_id: int, months: dict = Body(..., embed=True),
                   db: Session = Depends(get_db), user: User = Depends(require_writer)):
    """Backfill several months at once. Body: ``{"months": {"2026-06": {...}, ...}}``.

    Each value is merged onto that month's snapshot (existing events are kept; the
    numeric metrics that feed the charts are updated). Writer + edit rights."""
    squad = db.get(Squad, squad_id)
    if squad is None:
        raise HTTPException(status_code=404, detail="Squad introuvable")
    assert_can_edit_squad(db, user, squad_id)
    for period, snap in (months or {}).items():
        snap = _sanitized(snap or {})
        e = _entry(db, squad_id, period)
        if e is None:
            db.add(SteercoEntry(squad_id=squad_id, period=period, data=snap,
                                updated_by_user_id=user.id))
        else:
            e.data = {**(e.data or {}), **snap}
            e.updated_by_user_id = user.id
    record_audit(db, user.id, "steerco.history", entity="steerco", entity_id=squad_id,
                 detail={"squad_id": squad_id, "months": list((months or {}).keys())})
    db.commit()
    return {"squad_id": squad_id, "count": len(months or {})}


# --------------------------------------------------------------------------
# Aggregation: build the render-data (charts + rolling SLA) from 12 snapshots
# --------------------------------------------------------------------------

def _num(v) -> float | None:
    try:
        return float(str(v).replace("%", "").replace(",", ".").strip())
    except (TypeError, ValueError):
        return None


def _round_up(x: float) -> int:
    x = max(x, 1.0)
    mag = 10 ** (len(str(int(x))) - 1)
    return int(math.ceil(x / mag) * mag)


def _round_down(x: float) -> int:
    x = max(x, 0.0)
    if x < 1:
        return 0
    mag = 10 ** (len(str(int(x))) - 1)
    return int(math.floor(x / mag) * mag)


# --- Auto-computed indicators (never entered by the squad leader) -------------
# SLA colour: above 90% green, 80 to 90% amber, below 80% red.
SLA_GREEN, SLA_AMBER = 90.0, 80.0


def _sla_status(v) -> str | None:
    """RAG status of an SLA value, derived from the number itself (None when empty)."""
    n = _num(v)
    if n is None:
        return None
    return "ok" if n > SLA_GREEN else "warn" if n >= SLA_AMBER else "ko"


def _fmt_delta(d: float) -> str:
    r = round(d, 1)
    if r == 0:
        return "0"
    return (f"{r:+.10g}").replace(".", ",")


def _kpi_change(cur, prev) -> tuple[str, str]:
    """A KPI's (trend, delta) vs the previous month, computed from both values."""
    c, p = _num(cur), _num(prev)
    if c is None or p is None:
        return "flat", ""
    d = c - p
    return ("up" if d > 0 else "down" if d < 0 else "flat"), _fmt_delta(d)


def _kpi_value(snap: dict, label: str):
    """That month's value for a KPI label (case-insensitive), None when absent."""
    target = (label or "").strip().lower()
    for k in (snap.get("kpis") or []):
        if (k.get("label") or "").strip().lower() == target:
            return k.get("value")
    return None


def _aggregate(db: Session, squad_id: int, period: str, override: dict | None = None) -> dict:
    """Assemble the one-pager render-data for a squad+month from the last 12 snapshots:
    KPI cards + events from the current month; SLA table (current row + 12-month average
    row); KPI and incident charts as the 12-month series.

    ``override`` (optional) replaces the current month's snapshot in memory only (never
    persisted) so the wizard can preview unsaved edits before submitting."""
    keys = month_keys(period, 12)
    by_period = {e.period: (e.data or {}) for e in db.query(SteercoEntry)
                 .filter(SteercoEntry.squad_id == squad_id, SteercoEntry.period.in_(keys)).all()}
    if override is not None:
        by_period[period] = _sanitized(override)
    cur = by_period.get(period, {}) or {}
    prev = by_period.get(keys[-2], {}) or {} if len(keys) > 1 else {}

    # KPI change vs M-1: computed here from the two snapshots, so the squad leader
    # never types a variation and it can never go stale.
    kpis = []
    for k in (cur.get("kpis") or []):
        trend, delta = _kpi_change(k.get("value"), _kpi_value(prev, k.get("label") or ""))
        kpis.append({**k, "trend": trend, "delta": delta})

    sla_cur = cur.get("sla") or {}
    services = sla_cur.get("services") or []
    # SLA colour: computed from each value (see _sla_status), never chosen by hand.
    cur_cells = [{**(c or {}), "v": _clamp_pct((c or {}).get("v")),
                  "s": _sla_status(_clamp_pct((c or {}).get("v")))}
                 for c in (sla_cur.get("cells") or [])]
    avg_cells = []
    for i in range(len(services)):
        vals = []
        for k in keys:
            cells = (by_period.get(k, {}).get("sla") or {}).get("cells") or []
            if i < len(cells):
                pv = _num(cells[i].get("v"))
                if pv is not None:
                    vals.append(min(max(pv, 0.0), 100.0))
        if vals:
            a = sum(vals) / len(vals)
            avg_cells.append({"v": f"{a:.1f}".replace(".", ",") + "%", "s": _sla_status(a)})
        else:
            avg_cells.append({"v": "-", "s": None})
    sla = {"services": services, "rows": [
        {"period": "__current__", "cells": cur_cells},
        {"period": "__trailing__", "cells": avg_cells},
    ]}

    labels = [_month_short(k) for k in keys]
    kpi_series = []
    for idx, k in enumerate(kpis):
        label = k.get("label")
        # Same (case-insensitive) label matching as the vs-M-1 delta, so a KPI typed
        # "K8aaS" one month and "K8AAS" the next stays one single series.
        raw = [_num(_kpi_value(by_period.get(key, {}), label or "")) for key in keys]
        # Index each KPI to base 100 at its first value, so metrics of very different
        # magnitudes (Users ~250 vs Cluster K8s ~8) share the axis and read as trends.
        base = next((v for v in raw if v not in (None, 0)), None)
        data = [round(v / base * 100, 1) if (v is not None and base) else v for v in raw]
        kpi_series.append({"name": label, "color": SERIES_COLORS[idx % len(SERIES_COLORS)], "data": data})
    all_idx = [v for s in kpi_series for v in s["data"] if v is not None]
    kpi_ymax = _round_up(max(all_idx + [110]))
    kpi_ymin = _round_down(min(all_idx + [90]))

    inc_data = [_num(by_period.get(k, {}).get("incidents")) for k in keys]
    inc_ymax = _round_up(max([v for v in inc_data if v is not None] + [10]))

    return {
        "kpis": kpis,
        "sla": sla,
        "kpi_chart": {"labels": labels, "series": kpi_series, "y_max": kpi_ymax, "y_min": kpi_ymin, "indexed": True},
        "incidents_chart": {"labels": labels, "y_max": inc_ymax, "y_min": 0,
                            "series": [{"name": "Incidents", "color": "#D24545", "data": inc_data}]},
        "last_events": cur.get("last_events") or [],
        "next_events": cur.get("next_events") or [],
    }


# --------------------------------------------------------------------------
# One-pager rendering (HTML + PPTX), language-aware (default English)
# --------------------------------------------------------------------------

_PAGE_CSS = """
/* Palette aligned on the app theme (theme.css) so the in-app view feels native.
   The PPTX export keeps its own colours (rendered in _render_pptx). */
:root{--navy:#1E2761;--navy2:#141B47;--ice:#CADCFC;--ice-light:#E8F0FE;--ice-line:#CADCFC;
--bg:#F5F7FA;--card:#FFFFFF;--line:#E2E8F0;--txt:#1E293B;--muted:#64748B;--green:#027A48;--amber:#B54708;--red:#B42318;}
*{box-sizing:border-box;margin:0;padding:0;}
body{font-family:Calibri,"Segoe UI",system-ui,sans-serif;background:var(--bg);color:var(--txt);padding:22px;font-size:14px;}
.page{max-width:1280px;margin:0 auto;}
.page + .page{margin-top:26px;}
/* Compact, discreet header (no heavy band). */
.hdr{display:flex;align-items:baseline;justify-content:space-between;margin:2px 2px 12px;}
.hdr h1{font-size:19px;font-weight:700;letter-spacing:.3px;color:var(--navy);}
.hdr .date{color:var(--muted);font-size:12px;font-weight:600;}
/* 2x2 grid + events row: identical 16px gutters so every column lines up. */
.grid{display:grid;grid-template-columns:1fr 1fr;gap:16px;margin-bottom:16px;align-items:stretch;}
.events-grid{display:grid;grid-template-columns:1fr 1fr;gap:16px;}
/* Every block is the same panel: tinted header strip + body. */
.panel{display:flex;flex-direction:column;background:var(--card);border:1px solid var(--line);border-radius:12px;overflow:hidden;box-shadow:0 1px 2px rgba(16,37,66,.05);}
.panel .hd{display:flex;align-items:center;gap:8px;background:var(--ice-light);color:var(--navy);font-size:13.5px;font-weight:700;text-transform:uppercase;letter-spacing:.5px;padding:10px 14px;border-bottom:1px solid var(--ice-line);}
.panel .hd::before{content:"";width:8px;height:8px;border-radius:2px;background:var(--navy);flex:0 0 auto;}
.panel .hd .sub{color:var(--muted);text-transform:none;letter-spacing:0;font-weight:400;font-size:11px;}
.panel .bd{padding:12px 14px;flex:1;display:flex;flex-direction:column;min-height:0;}
/* Users alone and wider on top; the 4 others in a row below. */
.kpi-wrap{display:flex;flex-direction:column;gap:10px;height:100%;}
.kpi-hero-row{display:flex;justify-content:center;}
.kpi-row{flex:1;display:grid;grid-template-columns:repeat(4,1fr);gap:10px;}
/* 3 lines: number, name (single line), trend at the bottom-left. */
.kpi{background:#fff;border:1px solid #AFC0D6;border-radius:10px;padding:10px 12px;display:flex;flex-direction:column;}
.kpi.hero{width:44%;min-width:190px;padding:14px;}
.kpi .num{font-size:21px;font-weight:800;color:var(--navy);line-height:1.05;}
.kpi.hero .num{font-size:32px;}
.kpi .num .unit{font-size:12px;font-weight:400;color:var(--muted);margin-left:1px;}
.kpi .lbl{font-size:10px;font-weight:700;color:var(--navy2);text-transform:uppercase;letter-spacing:.3px;white-space:nowrap;overflow:hidden;text-overflow:ellipsis;margin-top:3px;}
.kpi.hero .lbl{font-size:14px;}
.kpi .t{font-size:11px;font-weight:600;margin-top:auto;padding-top:8px;}
.kpi-sub{display:flex;flex-wrap:wrap;gap:2px 8px;font-size:9px;color:var(--muted);margin-top:6px;line-height:1.3;}
.kpi-sub b{color:var(--navy);font-weight:700;}
.up{color:var(--green);}.down{color:var(--red);}.flat{color:var(--muted);}
/* SLA table fills the whole panel body; rows share the height evenly. */
table{width:100%;border-collapse:collapse;font-size:13px;height:100%;}
.bd > table{flex:1;}
th,td{padding:8px 10px;text-align:center;border:1px solid var(--line);vertical-align:middle;}
thead th{background:var(--navy);color:#fff;font-weight:600;}
tbody th{background:var(--ice-light);color:var(--navy);font-weight:700;text-align:left;}
/* Whole SLA cell is filled with the status colour, computed from the value (matches the PPTX). */
td.b-ok,td.b-warn,td.b-ko{font-weight:700;}
.b-ok{background:#E4F3EA;color:var(--green);}
.b-warn{background:#FBF0D9;color:#9c7212;}
.b-ko{background:#F7E0E0;color:var(--red);}
.b-none{color:var(--muted);}
.chart{flex:1;min-height:0;display:flex;}
.chart svg{width:100%;height:100%;display:block;}
.legend{display:flex;flex-wrap:wrap;gap:12px;margin-top:8px;font-size:11px;color:var(--muted);}
.legend span{display:flex;align-items:center;gap:5px;}
.legend i{width:14px;height:3px;border-radius:2px;display:inline-block;}
/* Uniform event row: date / type (colour chip) / description. */
.ev-list{list-style:none;}
.ev-list li{display:flex;align-items:center;gap:10px;padding:9px 2px;border-bottom:1px solid var(--line);}
.ev-list li:last-child{border-bottom:none;}
.ev-date{color:var(--muted);font-size:12px;font-weight:700;width:52px;flex:0 0 auto;}
.ev-type{font-size:11.5px;font-weight:700;color:var(--navy2);background:var(--ice-light);border:1px solid var(--ice-line);padding:2px 9px;border-radius:5px;flex:0 0 auto;min-width:52px;text-align:center;white-space:nowrap;}
.ev-desc{flex:1;font-size:13px;color:var(--txt);}
.muted{color:var(--muted);}
.empty{color:var(--muted);font-size:12px;font-style:italic;padding:6px 2px;}
"""


def _svg_line_chart(chart: dict, empty: str) -> str:
    """Server-render a multi-series line chart as SVG. Missing points (None) are gaps."""
    series = [s for s in (chart.get("series") or []) if any(v is not None for v in (s.get("data") or []))]
    labels = chart.get("labels") or []
    y_max = float(chart.get("y_max") or 100) or 100
    y_min = float(chart.get("y_min") or 0)
    span = (y_max - y_min) or 1
    if not series:
        return f"<div class='empty'>{escape(empty)}</div>"
    W, H, pl, pr, pt, pb = 560, 190, 34, 14, 12, 24
    iw, ih = W - pl - pr, H - pt - pb
    n = max(len(labels), max((len(s["data"]) for s in series), default=0), 2)

    def x(i):
        return pl + (iw * i / (n - 1))

    def y(v):
        return pt + ih - (ih * (float(v) - y_min) / span)

    out = []
    for g in range(5):
        yy = pt + ih * g / 4
        out.append(f'<line x1="{pl}" y1="{yy:.1f}" x2="{W-pr}" y2="{yy:.1f}" stroke="#E7ECF2" stroke-width="1"/>')
        out.append(f'<text x="{pl-6}" y="{yy+3:.1f}" font-size="9" fill="#6B7C90" text-anchor="end">{round(y_max - span*g/4)}</text>')
    out.append(f'<line x1="{pl}" y1="{pt}" x2="{pl}" y2="{pt+ih}" stroke="#B7C2CF" stroke-width="1"/>')
    out.append(f'<line x1="{pl}" y1="{pt+ih}" x2="{W-pr}" y2="{pt+ih}" stroke="#B7C2CF" stroke-width="1"/>')
    for i, m in enumerate(labels):
        out.append(f'<text x="{x(i):.1f}" y="{H-8}" font-size="9" fill="#6B7C90" text-anchor="middle">{escape(str(m))}</text>')
    for s in series:
        color = escape(str(s.get("color") or NAVY))
        pts = [(i, v) for i, v in enumerate(s["data"]) if v is not None]
        if pts:
            # Line only (no point markers).
            d = " ".join(("M" if j == 0 else "L") + f"{x(i):.1f} {y(v):.1f}" for j, (i, v) in enumerate(pts))
            out.append(f'<path d="{d}" fill="none" stroke="{color}" stroke-width="2.2" stroke-linejoin="round" stroke-linecap="round"/>')
    return f'<svg viewBox="0 0 {W} {H}" preserveAspectRatio="none" width="100%">{"".join(out)}</svg>'


def _legend(series: list[dict]) -> str:
    return "".join(
        f'<span><i style="background:{escape(str(s.get("color") or NAVY))}"></i>{escape(str(s.get("name") or ""))}</span>'
        for s in (series or []) if any(v is not None for v in (s.get("data") or []))
    )


def _kpi_card_html(k: dict, hero: bool = False) -> str:
    """3 lines: (1) the number, (2) the name on a single line, (3) the trend at the
    bottom-left. Software Factory additionally shows its sub-metrics small."""
    trend = k.get("trend") or "flat"
    arrow = TREND_ARROW.get(trend, "▬")
    unit = f"<span class='unit'>{escape(str(k.get('unit')))}</span>" if k.get("unit") else ""
    delta = f"{arrow} {escape(str(k.get('delta')))}" if k.get("delta") else ""
    label = escape(str(k.get("label") or "-").upper())
    val = escape(str(k.get("value") or "-"))
    trend_html = f"<div class='t {TREND_CLASS.get(trend, 'flat')}'>{delta}</div>" if delta else ""
    sub = k.get("sub") or []
    sub_html = ""
    if sub:
        parts = "".join(
            f"<span>{escape(str(s.get('label') or ''))} <b>{escape(str(s.get('value') or '-'))}</b></span>"
            for s in sub
        )
        sub_html = f"<div class='kpi-sub'>{parts}</div>"
    return (f"<div class='kpi{' hero' if hero else ''}'>"
            f"<div class='num'>{val}{unit}</div>"
            f"<div class='lbl'>{label}</div>"
            f"{sub_html}{trend_html}</div>")


def _kpi_cards(kpis: list[dict], L: dict) -> str:
    """First KPI (Users) on its own, centred and wider at the top; the rest in a row
    below (Landing Zone / K8aaS / DBaaS / Software Factory, the last with sub-metrics)."""
    if not kpis:
        return f"<div class='empty'>{escape(L['no_kpi'])}</div>"
    hero = _kpi_card_html(kpis[0], hero=True)
    rest = "".join(_kpi_card_html(k) for k in kpis[1:])
    return (f"<div class='kpi-wrap'>"
            f"<div class='kpi-hero-row'>{hero}</div>"
            f"<div class='kpi-row'>{rest}</div></div>")


def _row_label(period: str, L: dict) -> str:
    return {"__current__": L["current"], "__trailing__": L["trailing"]}.get(period, period)


def _sla_table(sla: dict, L: dict) -> str:
    services = (sla or {}).get("services") or []
    rows = (sla or {}).get("rows") or []
    if not services or not rows:
        return f"<div class='empty'>{escape(L['no_sla'])}</div>"
    head = "".join(f"<th>{escape(str(s))}</th>" for s in services)
    body = []
    for r in rows:
        cells = r.get("cells") or []
        tds = []
        for i in range(len(services)):
            c = cells[i] if i < len(cells) else {}
            v = escape(str((c or {}).get("v") or "-"))
            cls = SLA_CLASS.get((c or {}).get("s"), "b-none")
            tds.append(f'<td class="{cls}">{v}</td>')
        body.append(f"<tr><th>{escape(_row_label(str(r.get('period') or '-'), L))}</th>{''.join(tds)}</tr>")
    return (f"<table><thead><tr><th>{escape(L['period'])}</th>{head}</tr></thead>"
            f"<tbody>{''.join(body)}</tbody></table>")


def _events_list(events: list[dict], L: dict) -> str:
    """Uniform event row: date / type (colour chip) / description. The chip is tinted
    with the event's severity (the "Gravité" entered in the wizard / the Excel)."""
    if not events:
        return f"<div class='empty'>{escape(L['no_event'])}</div>"
    items = []
    for e in events:
        typ = escape(str(e.get("tag") or "-"))
        sev = e.get("sev") if e.get("sev") in SEV_FG else None
        chip = (f' style="background:{SEV_BG[sev]};color:{SEV_FG[sev]};border-color:{SEV_FG[sev]}"'
                if sev else "")
        items.append(
            f'<li><span class="ev-date">{escape(str(e.get("date") or ""))}</span>'
            f'<span class="ev-type"{chip}>{typ}</span>'
            f'<span class="ev-desc">{escape(str(e.get("text") or ""))}</span></li>'
        )
    return f'<ul class="ev-list">{"".join(items)}</ul>'


def _panel(title: str, sub: str, body: str) -> str:
    subhtml = f' <span class="sub">{sub}</span>' if sub else ""
    return f'<div class="panel"><div class="hd">{title}{subhtml}</div><div class="bd">{body}</div></div>'


def _chart_body(chart: dict, L: dict) -> str:
    return (f'<div class="chart">{_svg_line_chart(chart, L["no_data"])}</div>'
            f'<div class="legend">{_legend(chart.get("series") or [])}</div>')


def _onepager(squad_name: str, period: str, data: dict, L: dict) -> str:
    """One squad's KPI one-pager: header band, then a 2x2 panel grid (row 1 = KPIs |
    KPI chart, row 2 = SLA | incidents chart) with identical panels/gutters, then the
    last/next events row - every column and row edge lines up."""
    d = data or {}
    return (
        '<div class="page">'
        f'<div class="hdr"><h1>{escape(squad_name)}</h1><div class="date">{escape(_period_long(period, L))}</div></div>'
        '<div class="grid">'
        + _panel(escape(L["key_figures"]), "", _kpi_cards(d.get("kpis") or [], L))
        + _panel(escape(L["kpi_chart"]), escape(L["kpi_chart_sub"]), _chart_body(d.get("kpi_chart") or {}, L))
        + _panel(L["sla"], "", _sla_table(d.get("sla") or {}, L))
        + _panel(escape(L["inc_chart"]), escape(L["inc_chart_sub"]), _chart_body(d.get("incidents_chart") or {}, L))
        + '</div>'
        '<div class="events-grid">'
        + _panel(escape(L["last_events"]), "", _events_list(d.get("last_events") or [], L))
        + _panel(escape(L["next_events"]), "", _events_list(d.get("next_events") or [], L))
        + '</div>'
        '</div>'
    )


def _document(title: str, body: str, lang: str) -> str:
    return (f"<!doctype html><html lang='{lang}'><head><meta charset='utf-8'>"
            f"<meta name='viewport' content='width=device-width, initial-scale=1.0'>"
            f"<title>{escape(title)}</title><style>{_PAGE_CSS}</style></head><body>{body}</body></html>")


def _render_pptx(squads: list[dict], period: str, L: dict) -> bytes:
    """Native PPTX: one 16:9 slide per squad, same layout as the HTML (left KPIs+SLA,
    right KPI chart over incidents chart, bottom events)."""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    def rgb(h):
        h = (h or "#000000").lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def add_text(slide, x, y, w, h, text, size=12, bold=False, color="#1B2A3D"):
        tf = slide.shapes.add_textbox(x, y, w, h).text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
        return tf

    def no_shadow(s):
        s.shadow.inherit = False

    def panel(slide, x, y, w, h):
        p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        p.fill.solid(); p.fill.fore_color.rgb = rgb("#FFFFFF")
        p.line.color.rgb = rgb("#E1E7EF"); p.line.width = Pt(0.75); no_shadow(p)
        return p

    def titled_panel(slide, x, y, w, h, title, sub=""):
        """Panel + tinted header strip (navy title). x/y/w/h are inches; returns the
        inner body box (bx, by, bw, bh) in inches."""
        panel(slide, Inches(x), Inches(y), Inches(w), Inches(h))
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
        strip.fill.solid(); strip.fill.fore_color.rgb = rgb("#EDF2F8")
        strip.line.fill.background(); no_shadow(strip)
        tf = strip.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.14); tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.02)
        r = tf.paragraphs[0].add_run(); r.text = title.upper()
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = rgb(NAVY)
        if sub:
            rs = tf.paragraphs[0].add_run(); rs.text = "  " + sub.replace("&amp;", "&")
            rs.font.size = Pt(10); rs.font.color.rgb = rgb("#6B7C90")
        return x + 0.14, y + 0.44, w - 0.28, h - 0.56

    def kpi_card(slide, x, y, w, h, k, hero=False):
        """3 lines: (1) number, (2) name on one line, (3) trend at the bottom-left.
        Software Factory also shows its sub-metrics small. x/y/w/h in EMU."""
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.fill.solid(); card.fill.fore_color.rgb = rgb("#FFFFFF")
        card.line.color.rgb = rgb("#AFC0D6"); card.line.width = Pt(1.0); no_shadow(card)
        num_sz = 24 if hero else 19
        lbl_sz = 13 if hero else 9
        pad = Inches(0.1)
        iw = w - Inches(0.16)
        unit = str(k.get("unit")) if k.get("unit") else ""
        trend = k.get("trend") if k.get("trend") in ("up", "down", "flat") else "flat"
        tcolor = {"up": "#2E9E5B", "down": "#D24545", "flat": "#6B7C90"}[trend]

        def box(yy, hh, wrap=False):
            tf = slide.shapes.add_textbox(x + pad, yy, iw, hh).text_frame
            tf.word_wrap = wrap
            tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = Emu(9000); tf.margin_bottom = Emu(9000)
            return tf

        yy = y + Inches(0.07)
        num_h = Inches(0.36 if hero else 0.28)
        nbox = box(yy, num_h)
        rn = nbox.paragraphs[0].add_run(); rn.text = f"{k.get('value') or '-'}{unit}"
        rn.font.size = Pt(num_sz); rn.font.bold = True; rn.font.color.rgb = rgb(NAVY)
        yy += num_h
        # Name: dedicated box, no wrap -> always a single line.
        lbox = box(yy, Inches(0.2))
        rl = lbox.paragraphs[0].add_run(); rl.text = str(k.get("label") or "-").upper()
        rl.font.size = Pt(lbl_sz); rl.font.bold = True; rl.font.color.rgb = rgb("#141B47")
        yy += Inches(0.2)
        sub = k.get("sub") or []
        if sub:
            sbox = box(yy, Inches(0.3), wrap=True)
            p = sbox.paragraphs[0]
            for s in sub:
                rl2 = p.add_run(); rl2.text = f"{s.get('label') or ''} "
                rl2.font.size = Pt(8); rl2.font.color.rgb = rgb("#64748B")
                rv2 = p.add_run(); rv2.text = f"{s.get('value') or '-'}   "
                rv2.font.size = Pt(8); rv2.font.bold = True; rv2.font.color.rgb = rgb(NAVY)
        # Trend: inline for the hero (right under the name), bottom-left otherwise.
        if k.get("delta"):
            ty = yy if hero else (y + h - Inches(0.26))
            tbox = box(ty, Inches(0.22))
            rt = tbox.paragraphs[0].add_run(); rt.text = f"{TREND_ARROW.get(trend, '▬')} {k.get('delta')}".strip()
            rt.font.size = Pt(10 if hero else 9); rt.font.bold = True; rt.font.color.rgb = rgb(tcolor)

    def set_cell(cell, fill, color, size=10, bold=False):
        cell.fill.solid(); cell.fill.fore_color.rgb = rgb(fill)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Emu(14000); cell.margin_bottom = Emu(14000)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)

    def sla_table(slide, x, y, w, h, sla):
        services = sla.get("services") or []
        rows = sla.get("rows") or []
        if not services or not rows:
            return
        nrows, ncols = len(rows) + 1, len(services) + 1
        tbl = slide.shapes.add_table(nrows, ncols, x, y, w, Inches(h)).table
        # Header a touch shorter, data rows share the rest so the table fills the box.
        head_h = min(0.42, h / nrows)
        data_h = (h - head_h) / (nrows - 1) if nrows > 1 else head_h
        tbl.rows[0].height = Inches(head_h)
        for i in range(1, nrows):
            tbl.rows[i].height = Inches(data_h)
        for j, txt in enumerate([L["period"]] + [str(s) for s in services]):
            tbl.cell(0, j).text = txt
            set_cell(tbl.cell(0, j), NAVY, "#FFFFFF", bold=True)
        fills = {"ok": "#E4F3EA", "warn": "#FBF0D9", "ko": "#F7E0E0"}
        texts = {"ok": "#2E9E5B", "warn": "#9C7212", "ko": "#D24545"}
        for i, r in enumerate(rows):
            tbl.cell(i + 1, 0).text = _row_label(str(r.get("period", "-")), L)
            set_cell(tbl.cell(i + 1, 0), "#EEF4F8", NAVY, bold=True)
            cells = r.get("cells") or []
            for j in range(len(services)):
                cell = cells[j] if j < len(cells) else {}
                s = (cell or {}).get("s")
                tbl.cell(i + 1, j + 1).text = str((cell or {}).get("v", "-"))
                set_cell(tbl.cell(i + 1, j + 1), fills.get(s, "#FFFFFF"), texts.get(s, "#6B7C90"), bold=True)

    def line_chart(slide, x, y, w, h, chart):
        series = [s for s in (chart.get("series") or []) if any(v is not None for v in (s.get("data") or []))]
        if not series:
            add_text(slide, x, y, w, Inches(0.3), L["no_data"], size=10, color="#98A2B3")
            return
        cd = CategoryChartData()
        cd.categories = chart.get("labels") or [str(i + 1) for i in range(max(len(s["data"]) for s in series))]
        for s in series:
            cd.add_series(str(s.get("name") or ""), [None if v is None else float(v) for v in s["data"]])
        gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, w, h, cd)
        ch = gf.chart
        ch.has_title = False
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(8)
        for i, ps in enumerate(ch.series):
            ps.format.line.color.rgb = rgb(series[i].get("color") or NAVY)
            ps.format.line.width = Pt(2)
        try:
            ch.value_axis.maximum_scale = float(chart.get("y_max") or 100)
            ch.value_axis.minimum_scale = float(chart.get("y_min") or 0)
            ch.value_axis.tick_labels.font.size = Pt(8)
            ch.category_axis.tick_labels.font.size = Pt(8)
        except Exception:  # pragma: no cover
            pass

    def events(slide, x, y, w, h, evs):
        tf = slide.shapes.add_textbox(x, y, w, h).text_frame
        tf.word_wrap = True
        first = True
        if not evs:
            re = tf.paragraphs[0].add_run(); re.text = L["no_event"]
            re.font.size = Pt(12); re.font.italic = True; re.font.color.rgb = rgb("#6B7C90")
        for e in evs:
            # Uniform row: date / type / description.
            pe = tf.paragraphs[0] if first else tf.add_paragraph()
            pe.space_after = Pt(5)
            first = False
            rd = pe.add_run(); rd.text = f"{e.get('date', '')}    "
            rd.font.size = Pt(12); rd.font.bold = True; rd.font.color.rgb = rgb("#6B7C90")
            # Type coloured by the event's severity, like the HTML chip.
            rtp = pe.add_run(); rtp.text = f"{str(e.get('tag') or '-')}    "
            rtp.font.size = Pt(12); rtp.font.bold = True
            rtp.font.color.rgb = rgb(SEV_FG.get(e.get("sev"), "#13315C"))
            rt = pe.add_run(); rt.text = str(e.get("text", ""))
            rt.font.size = Pt(12); rt.font.color.rgb = rgb("#1B2A3D")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Symmetric grid geometry (inches): equal columns, equal gutters, aligned rows.
    LM = 0.45
    GUT = 0.22
    COLW = (13.333 - 2 * LM - GUT) / 2          # left/right column width
    LX, RX = LM, LM + COLW + GUT
    for s in squads:
        slide = prs.slides.add_slide(blank)
        d = s["data"] or {}
        # Compact, discreet header: navy name (left) + muted period (right).
        add_text(slide, Inches(LM), Inches(0.26), Inches(9), Inches(0.4), s["squad_name"], size=19, bold=True, color=NAVY)
        pmeta = slide.shapes.add_textbox(Inches(13.333 - LM - 4), Inches(0.32), Inches(4), Inches(0.3)).text_frame
        pp = pmeta.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
        rp = pp.add_run(); rp.text = _period_long(period, L); rp.font.size = Pt(12); rp.font.bold = True; rp.font.color.rgb = rgb("#6B7C90")

        # 2x2 grid: row 1 = KPIs | KPI chart, row 2 = SLA | incidents chart.
        gy, gb, vg = 0.92, 5.9, 0.18
        row_h = (gb - gy - vg) / 2
        r1y, r2y = gy, gy + row_h + vg

        # Row 1 left: KPI panel. Users (hero) centred on top, the rest in a row below.
        bx, by, bw, bh = titled_panel(slide, LX, r1y, COLW, row_h, L["key_figures"])
        kpis = d.get("kpis") or []
        if kpis:
            cg = 0.1
            hero_h = min(0.92, bh * 0.48)
            hero_w = bw * 0.46
            kpi_card(slide, Inches(bx + (bw - hero_w) / 2), Inches(by), Inches(hero_w), Inches(hero_h), kpis[0], hero=True)
            rest = kpis[1:5]
            if rest:
                ry = by + hero_h + cg
                rh = bh - hero_h - cg
                cw = (bw - cg * (len(rest) - 1)) / len(rest)
                for i, k in enumerate(rest):
                    kpi_card(slide, Inches(bx + i * (cw + cg)), Inches(ry), Inches(cw), Inches(rh), k)
        # Row 1 right: KPI chart.
        bx, by, bw, bh = titled_panel(slide, RX, r1y, COLW, row_h, L["kpi_chart"], L["kpi_chart_sub"])
        line_chart(slide, Inches(bx), Inches(by), Inches(bw), Inches(bh), d.get("kpi_chart") or {})

        # Row 2 left: SLA table.
        bx, by, bw, bh = titled_panel(slide, LX, r2y, COLW, row_h, L["sla"])
        sla_table(slide, Inches(bx), Inches(by), Inches(bw), bh, d.get("sla") or {})
        # Row 2 right: incidents chart.
        bx, by, bw, bh = titled_panel(slide, RX, r2y, COLW, row_h, L["inc_chart"], L["inc_chart_sub"])
        line_chart(slide, Inches(bx), Inches(by), Inches(bw), Inches(bh), d.get("incidents_chart") or {})

        # Bottom: events row, same two columns.
        ey, eh = gb + 0.16, 1.28
        bx, by, bw, bh = titled_panel(slide, LX, ey, COLW, eh, L["last_events"])
        events(slide, Inches(bx), Inches(by), Inches(bw), Inches(bh), d.get("last_events") or [])
        bx, by, bw, bh = titled_panel(slide, RX, ey, COLW, eh, L["next_events"])
        events(slide, Inches(bx), Inches(by), Inches(bw), Inches(bh), d.get("next_events") or [])

    if not squads:
        slide = prs.slides.add_slide(blank)
        add_text(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.6), f"Steerco {period}", size=24, bold=True, color=NAVY)
        add_text(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5), L["no_squads"], size=12, color="#6B7C90")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --------------------------------------------------------------------------
# Document endpoints
# --------------------------------------------------------------------------

def _squad_in_scope(db: Session, user: User, squad_id: int) -> Squad:
    squad = db.get(Squad, squad_id)
    tid = visible_tribe_id(user)
    if squad is None or (tid is not None and squad.tribe_id != tid):
        raise HTTPException(status_code=404, detail="Squad introuvable")
    return squad


def _enabled_squads(db: Session, user: User) -> list[Squad]:
    return [s for s in _squads_in_scope(db, user) if s.steerco_enabled]


@router.get("/onepager.html", response_class=HTMLResponse)
def onepager_html(squad_id: int = Query(...), period: str = Query(...), lang: str | None = Query(None),
                  db: Session = Depends(get_db), user: User = Depends(require_tribe_or_admin)):
    """One squad's KPI one-pager (auto-built from the last 12 monthly snapshots)."""
    squad = _squad_in_scope(db, user, squad_id)
    L = I18N[_lang(lang)]
    body = _onepager(squad.name, period, _aggregate(db, squad_id, period), L)
    return HTMLResponse(_document(f"Steerco {squad.name} {period}", body, _lang(lang)))


@router.post("/squad/{squad_id}/preview.html", response_class=HTMLResponse)
def preview_html(squad_id: int, period: str = Query(...), data: dict = Body(...),
                 lang: str | None = Query(None), db: Session = Depends(get_db),
                 user: User = Depends(require_writer)):
    """Live preview of the one-pager for the wizard, using the still-unsaved snapshot
    (``data``) as the current month. Nothing is persisted. Squad-leader accessible for
    their own squad (they are the one filling the report)."""
    squad = db.get(Squad, squad_id)
    if squad is None:
        raise HTTPException(status_code=404, detail="Squad introuvable")
    assert_can_edit_squad(db, user, squad_id)
    L = I18N[_lang(lang)]
    body = _onepager(squad.name, period, _aggregate(db, squad_id, period, override=data), L)
    return HTMLResponse(_document(f"Steerco {squad.name} {period}", body, _lang(lang)))


@router.get("/document.html", response_class=HTMLResponse)
def document_html(period: str = Query(...), lang: str | None = Query(None), db: Session = Depends(get_db),
                  user: User = Depends(require_tribe_or_admin)):
    """Consolidated steerco (all steerco-enabled squads) as HTML one-pagers."""
    L = I18N[_lang(lang)]
    squads = _enabled_squads(db, user)
    pages = "".join(_onepager(s.name, period, _aggregate(db, s.id, period), L) for s in squads)
    if not pages:
        pages = f"<div class='page'><p class='empty'>{escape(L['no_squads'])}</p></div>"
    return HTMLResponse(_document(f"Steerco {period}", pages, _lang(lang)))


@router.get("/document.pptx")
def document_pptx(period: str = Query(...), lang: str | None = Query(None), db: Session = Depends(get_db),
                  user: User = Depends(require_tribe_or_admin)):
    """Consolidated steerco as PPTX. 501 when python-pptx is unavailable."""
    L = I18N[_lang(lang)]
    squads = [{"squad_name": s.name, "data": _aggregate(db, s.id, period)} for s in _enabled_squads(db, user)]
    try:
        payload = _render_pptx(squads, period, L)
    except ModuleNotFoundError:
        raise HTTPException(status_code=501, detail="Generation PPTX indisponible (python-pptx non installe)")
    return Response(
        content=payload,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="steerco_{period}.pptx"'},
    )
