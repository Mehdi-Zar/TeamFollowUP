"""PowerPoint rendering: the four decks the product exports.

Dashboard, roadmap swimlane, initiatives and cross-team dependencies. Kept apart
from ``report`` because it is 900 lines that share almost nothing with the HTML
path - eight names, all of which now live in ``reportcommon`` - and because the
python-pptx machinery reads very differently from string templating.

The public names stay importable from ``report`` as before: it re-exports them,
so nothing that used ``from .report import render_pptx`` had to change.
"""
from __future__ import annotations

import html
import io
import os
from datetime import datetime, timedelta, timezone

from sqlalchemy import select
from sqlalchemy.orm import Session, selectinload

from . import pptxtpl
from . import status as st
from .generalconfig import get_general
from .models import Squad, Tribe, utcnow
from .serializers import annual_progress, budget_out, dependency_label
from .reportcommon import (STAGE_COLOR, _DEP_T, _INIT_T, _MONTHS, _lang, _status_label,
                           _status_rag, group_by_theme, rt)


# ----- PPTX rendering -------------------------------------------------------------

# Brand palette (mirrors the app theme): navy / accent + RAG.
_BRAND = {
    "navy": "#1E2761", "navy_deep": "#141B47", "accent": "#175CD3",
    "green": "#027A48", "orange": "#B54708", "red": "#B42318",
    "ink": "#1F2937", "muted": "#6B7280", "card": "#F1F5F9",
    "line": "#E2E8F0", "white": "#FFFFFF", "zebra": "#F8FAFC",
}


_RAG_BRAND = {"red": "#B42318", "amber": "#B54708", "green": "#027A48", "grey": "#6B7280"}


# Runaway guard on per-squad detail slides. Set well above any realistic squad
# count so an explicit export never silently drops the squads the user picked;
# if it is ever exceeded, render_pptx adds a visible "N more squads" notice slide
# rather than dropping them without a trace.
_MAX_DETAIL_SLIDES = 300


def _pptx_toolkit():
    """Import python-pptx lazily and return the bits used to build a deck."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    return Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE


def render_pptx(data: dict) -> bytes:
    """Render the weekly report as a branded deck (requires python-pptx):
    a summary one-pager + one full detail slide per squad (annual objectives +
    roadmap/milestones by quarter + advancement). The roadmap-only swimlane deck
    is produced separately by render_roadmap_pptx."""
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE = _pptx_toolkit()

    def rgb(hexstr: str) -> RGBColor:
        return RGBColor.from_string(hexstr.lstrip("#").upper())

    B = {k: rgb(v) for k, v in _BRAND.items()}
    lang = data.get("lang", "fr")

    prs = pptxtpl.new_presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def new_slide():
        return pptxtpl.add_slide(prs)

    def textbox(s, left, top, width, height, text, size, *, bold=False, color=None,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        box = s.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color if color is not None else B["ink"]
        return box

    def rect(s, left, top, width, height, fill, line=None):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line
        sh.shadow.inherit = False
        return sh

    def place(sh, lines, *, anchor=MSO_ANCHOR.TOP, ml=0.1, mt=0.06, mr=0.1, mb=0.06):
        """Write paragraphs INTO a shape's own text frame, so the text is part of
        the shape (not a separate textbox floating on top). Each line is
        (text, size_pt, color[, bold, align, space_after_pt])."""
        tf = sh.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = Inches(ml); tf.margin_right = Inches(mr)
        tf.margin_top = Inches(mt); tf.margin_bottom = Inches(mb)
        for i, ln in enumerate(lines):
            txt, size, color = ln[0], ln[1], ln[2]
            bold = ln[3] if len(ln) > 3 else False
            align = ln[4] if len(ln) > 4 else PP_ALIGN.LEFT
            sa = ln[5] if len(ln) > 5 else 2
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(sa)
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        return sh

    def bullets(s, left, top, width, height, lines, size):
        """A text box with one paragraph per (text, color, bold) line."""
        box = s.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        for i, (txt, color, bold) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(2)
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        return box

    def style_cell(cell, text, size, color, *, bold=False, align=PP_ALIGN.LEFT, fill=None):
        cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
        cell.margin_top = Emu(0); cell.margin_bottom = Emu(0)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fill is not None:
            cell.fill.solid(); cell.fill.fore_color.rgb = fill
        else:
            cell.fill.background()
        cell.text = text or " "
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        if p.runs:
            r = p.runs[0]
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

    # ----- app-styled primitives (rounded cards, chips, progress bars) -----------
    def rrect(s, left, top, width, height, fill, *, line=None, radius=0.08):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line; sh.line.width = Pt(1)
        sh.shadow.inherit = False
        return sh

    def card(s, left, top, width, height, title=None):
        sh = rrect(s, left, top, width, height, B["white"], line=B["line"], radius=0.05)
        if title:
            # Title lives in the card's own text frame (top-anchored), not as an
            # overlay; body content is added by the caller as child shapes.
            place(sh, [(title, 12, B["navy"], True)], anchor=MSO_ANCHOR.TOP, ml=0.18, mt=0.12, mr=0.18)
        return sh

    def chip(s, left, top, text, fill, *, color=None, size=10):
        w = Inches(0.26 + 0.082 * len(text))
        sh = rrect(s, left, top, w, Inches(0.3), fill, radius=0.5)
        tf = sh.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = Inches(0.07)
        tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = color if color is not None else B["white"]
        return Emu(int(left) + int(w))

    def chips_row(s, left, top, items, *, gap=0.1):
        x = int(left)
        for text, fill in items:
            x = int(chip(s, Emu(x), top, text, fill))
            x += int(Inches(gap))

    def pbar(s, left, top, width, pct, color):
        rrect(s, left, top, width, Inches(0.12), rgb("#E2E8F0"), radius=0.5)
        p = max(0, min(100, int(pct or 0)))
        if p > 0:
            rrect(s, left, top, Emu(int(int(width) * p / 100)), Inches(0.12), color, radius=0.5)

    gen = data["generated_at"]
    gen_str = gen.strftime("%d/%m/%Y %H:%M") if isinstance(gen, datetime) else str(gen)
    sm = data["summary"]
    squads = [r for blk in data["tribes"] for r in blk["squads"]]
    margin = Inches(0.5)

    # ---------------- Summary one-pager (report kind) -----------------------------
    def summary_slide():
        s = new_slide()
        band = rect(s, Inches(0), Inches(0), prs.slide_width, Inches(1.12), B["navy"])
        # Title + subtitle are the band's own text; the right-corner meta stays a
        # corner label (left title and right meta can't share one text frame).
        place(band, [
            (f'{data["app_name"]} - {rt(lang, "report")}', 26, B["white"], True, PP_ALIGN.LEFT, 5),
            (f'{data["scope_name"]} · {rt(lang, "year")} {data["year"]}', 13, rgb("#C7D2FE"), False, PP_ALIGN.LEFT, 0),
        ], anchor=MSO_ANCHOR.TOP, ml=0.55, mt=0.18, mr=4.2)
        textbox(s, Inches(9.3), Inches(0.3), Inches(3.5), Inches(0.6),
                f'{rt(lang, "generated_full", d=gen_str)}\n{rt(lang, "window_full", n=data["since_days"])}', 11,
                color=rgb("#C7D2FE"), align=PP_ALIGN.RIGHT)

        kpis = [
            (rt(lang, "k_squads"), str(sm["squads_total"]), B["navy"]),
            (rt(lang, "k_progress"), f'{sm["avg_progress"]}%', B["accent"]),
            (rt(lang, "k_blocked"), str(sm["blocked"]), B["red"] if sm["blocked"] else B["ink"]),
            (rt(lang, "k_atrisk"), str(sm["at_risk"]), B["orange"] if sm["at_risk"] else B["ink"]),
            (rt(lang, "k_obj_red"), str(sm["objectives_red"]), B["red"] if sm["objectives_red"] else B["ink"]),
            (rt(lang, "k_stale"), str(sm["stale"]), B["orange"] if sm["stale"] else B["ink"]),
        ]
        gap = Inches(0.14)
        n = len(kpis)
        total_w = prs.slide_width - margin * 2
        card_w = Emu(int((total_w - gap * (n - 1)) / n))
        ky, kh = Inches(1.34), Inches(1.05)
        for i, (label, val, color) in enumerate(kpis):
            left = Emu(int(margin) + i * (int(card_w) + int(gap)))
            kp = rect(s, left, ky, card_w, kh, B["card"], line=B["line"])
            # Value + label are the card's own text, vertically centered.
            place(kp, [(val, 26, color, True, PP_ALIGN.CENTER, 4),
                       (label, 10, B["muted"], False, PP_ALIGN.CENTER, 0)],
                  anchor=MSO_ANCHOR.MIDDLE, ml=0.05, mr=0.05)

        headers = [rt(lang, "h_squad"), rt(lang, "h_leader"), rt(lang, "h_status"),
                   rt(lang, "h_progress"), rt(lang, "h_delta"), rt(lang, "h_blocked"), rt(lang, "h_atrisk")]
        wfrac = [0.275, 0.21, 0.12, 0.105, 0.082, 0.103, 0.105]
        table_w = int(prs.slide_width - margin * 2)
        widths = [Emu(int(table_w * f)) for f in wfrac]

        MAX = 18
        shown = squads[:MAX]
        overflow = len(squads) - len(shown)
        nrows = len(shown) + 1 + (1 if overflow > 0 else 0)
        top = Inches(2.62)
        tbl = s.shapes.add_table(max(nrows, 2), len(headers), margin, top, Emu(table_w),
                                 Inches(0.34) + Inches(0.255) * (nrows - 1)).table
        for ci, w in enumerate(widths):
            tbl.columns[ci].width = w
        for ci, h in enumerate(headers):
            align = PP_ALIGN.LEFT if ci < 2 else PP_ALIGN.CENTER
            style_cell(tbl.cell(0, ci), h, 10, B["white"], bold=True, align=align, fill=B["navy"])
        for ri, r in enumerate(shown, start=1):
            zebra = B["zebra"] if ri % 2 == 0 else B["white"]
            delta = r["delta"]
            cells = [
                (r["name"], B["ink"], True, PP_ALIGN.LEFT),
                (r["leader"] or "-", B["muted"], False, PP_ALIGN.LEFT),
                (_status_label(r["status"], lang), rgb(_RAG_BRAND[r["status_rag"]]), True, PP_ALIGN.CENTER),
                (f'{r["annual_pct"]}%', B["ink"], False, PP_ALIGN.CENTER),
                ((f'+{delta}' if delta > 0 else str(delta)),
                 (B["green"] if delta > 0 else B["red"]) if delta else B["muted"], False, PP_ALIGN.CENTER),
                (str(r["blocked"] or "-"), B["red"] if r["blocked"] else B["muted"], r["blocked"] > 0, PP_ALIGN.CENTER),
                (str(r["at_risk"] or "-"), B["orange"] if r["at_risk"] else B["muted"], r["at_risk"] > 0, PP_ALIGN.CENTER),
            ]
            for ci, (val, color, bold, align) in enumerate(cells):
                style_cell(tbl.cell(ri, ci), val, 9.5, color, bold=bold, align=align, fill=zebra)
        if overflow > 0:
            last = len(shown) + 1
            style_cell(tbl.cell(last, 0), rt(lang, "more_squads", n=overflow), 9, B["muted"], align=PP_ALIGN.LEFT)
            for ci in range(1, len(headers)):
                style_cell(tbl.cell(last, ci), " ", 9, B["muted"])

    # ---------------- Per-squad detail slide --------------------------------------
    def detail_slide(r, *, with_objectives):
        det = r.get("detail") or {}
        s = new_slide()
        band = rect(s, Inches(0), Inches(0), prs.slide_width, Inches(0.92), B["navy"])
        place(band, [
            (r["name"], 22, B["white"], True, PP_ALIGN.LEFT, 4),
            (f'{rt(lang, "year")} {data["year"]} · {rt(lang, "h_progress_long")} {r["annual_pct"]}%',
             12, rgb("#C7D2FE"), False, PP_ALIGN.LEFT, 0),
        ], anchor=MSO_ANCHOR.TOP, ml=0.5, mt=0.12, mr=0.5)

        # Budget line (status + figures), just under the header band.
        bud = det.get("budget")
        if bud is not None:
            f = lambda v: "-" if v is None else f"{v:,.0f} €"
            st_color = {"on_track": "green", "at_risk": "amber", "over": "red"}[bud["status"]]
            st_lbl = rt(lang, {"on_track": "b_on_track", "at_risk": "b_at_risk", "over": "b_over"}[bud["status"]])
            bline = f'{rt(lang, "h_budget")}: {st_lbl} · {rt(lang, "b_total")} {f(bud["total"])}'
            if bud["spent"] is not None:
                pc = f' ({bud["spent_pct"]}%)' if bud.get("spent_pct") is not None else ""
                bline += f' · {rt(lang, "b_spent")} {f(bud["spent"])}{pc}'
            if bud["forecast"] is not None:
                pc = f' ({bud["forecast_pct"]}%)' if bud.get("forecast_pct") is not None else ""
                bline += f' · {rt(lang, "b_forecast")} {f(bud["forecast"])}{pc}'
            if bud["status"] == "over":
                bline += f' · +{f(bud["overrun"])} ({bud["overrun_pct"]}%)'
            textbox(s, margin, Inches(0.96), Inches(12.3), Inches(0.22), bline, 10, bold=True,
                    color=rgb(_RAG_BRAND[st_color]))

        top = Inches(1.22) if bud is not None else Inches(1.15)
        if with_objectives:
            textbox(s, margin, top, Inches(12.3), Inches(0.3),
                    rt(lang, "h_otd_section"), 13, bold=True, color=B["navy"])
            objs = det.get("objectives", [])
            if objs:
                lines = []
                for o in objs[:6]:
                    rag = _status_rag(o["rag"])
                    dl = f' · {rt(lang, "deadline")} {o["target_date"]}' if o.get("target_date") else ""
                    lines.append((f'•  {o["title"]}   ({_status_label(o["rag"], lang)}{dl})',
                                  rgb(_RAG_BRAND[rag]), False))
                if len(objs) > 6:
                    lines.append((f'+{len(objs) - 6}…', B["muted"], False))
                bullets(s, margin, Inches(1.5), Inches(12.3), Inches(1.4), lines, 11)
            else:
                textbox(s, margin, Inches(1.5), Inches(12.3), Inches(0.3),
                        rt(lang, "no_obj"), 11, color=B["muted"])
            top = Inches(3.05)

        textbox(s, margin, top, Inches(12.3), Inches(0.3),
                rt(lang, "h_roadmap"), 13, bold=True, color=B["navy"])
        qtop = Emu(int(top) + int(Inches(0.42)))
        gap = Inches(0.2)
        n = 4
        total_w = prs.slide_width - margin * 2
        col_w = Emu(int((total_w - gap * (n - 1)) / n))
        # Reserve a band at the bottom for key messages when there are any.
        kms = det.get("key_messages") or []
        roadmap_bottom = int(Inches(6.35)) if kms else int(Inches(7.5))
        qh = Emu(roadmap_bottom - int(qtop) - int(Inches(0.3)))
        for i, qd in enumerate(det.get("quarters", [])):
            left = Emu(int(margin) + i * (int(col_w) + int(gap)))
            qc = rect(s, left, qtop, col_w, qh, B["card"], line=B["line"])
            items = qd.get("items", [])
            lines = []
            for it in items[:10]:
                rag = _status_rag(it["status"])
                dep = f'   · {rt(lang, "dep")} {it["dependency"]}' if it.get("dependency") else ""
                lines.append((f'•  {it["title"]}{dep}', rgb(_RAG_BRAND[rag]), False))
            if not items:
                lines.append((rt(lang, "no_jalon"), B["muted"], False))
            elif len(items) > 10:
                lines.append((f'+{len(items) - 10}…', B["muted"], False))
            # Header + milestones are the card's own text.
            paras = [(f'Q{qd["q"]} - {qd["pct"]}%', 12, B["ink"], True, PP_ALIGN.LEFT, 6)]
            paras += [(txt, 9.5, color, bold, PP_ALIGN.LEFT, 2) for (txt, color, bold) in lines]
            place(qc, paras, anchor=MSO_ANCHOR.TOP, ml=0.1, mt=0.08, mr=0.1)

        # Key messages band (success / alert / risk) along the bottom.
        if kms:
            ktop = Emu(roadmap_bottom + int(Inches(0.05)))
            textbox(s, margin, ktop, Inches(12.3), Inches(0.3),
                    rt(lang, "h_key_messages"), 12, bold=True, color=B["navy"])
            km_rag = {"success": "green", "alert": "amber", "risk": "red"}
            klines = []
            for m in kms[:4]:
                rag = km_rag.get(m["kind"], "grey")
                ts = f'   ({m["created_at"]})' if m.get("created_at") else ""
                klines.append((f'•  [{rt(lang, "km_" + m["kind"])}] {m["text"]}{ts}',
                               rgb(_RAG_BRAND[rag]), False))
            if len(kms) > 4:
                klines.append((f'+{len(kms) - 4}…', B["muted"], False))
            bullets(s, margin, Emu(int(ktop) + int(Inches(0.32))), Inches(12.3), Inches(0.7), klines, 9.5)

    # ---------------- Single-squad page slide (mirrors the squad page order) ------
    def squad_page_slide(r):
        det = r.get("detail") or {}
        s = new_slide()
        SW = prs.slide_width
        rect(s, Inches(0), Inches(0), SW, prs.slide_height, rgb("#F5F7FA"))  # app background
        L = Inches(0.4)
        FULLW = Emu(int(SW) - 2 * int(L))
        colw = Emu((int(FULLW) - int(Inches(0.2))) // 2)
        rcolx = Emu(int(L) + int(colw) + int(Inches(0.2)))

        # ----- header card (navy) -----
        hdr = rrect(s, L, Inches(0.32), FULLW, Inches(0.84), B["navy"], radius=0.08)
        place(hdr, [(r["name"], 20, B["white"], True)], anchor=MSO_ANCHOR.MIDDLE, ml=0.3, mr=4.6)
        textbox(s, Emu(int(SW) - int(Inches(4.6))), Inches(0.5), Inches(4.0), Inches(0.3),
                f'{rt(lang, "h_leader")} : {r["leader"] or "-"}', 11, color=rgb("#C7D2FE"), align=PP_ALIGN.RIGHT)

        # ----- badges row (mirrors the page header) -----
        items = [(f'{rt(lang, "h_progress_long")} {r["annual_pct"]}%', B["accent"])]
        if r["blocked"]:
            items.append((f'{r["blocked"]} {rt(lang, "h_blocked")}', B["red"]))
        if r["at_risk"]:
            items.append((f'{r["at_risk"]} {rt(lang, "h_atrisk")}', B["orange"]))
        items.append((rt(lang, "stale") if r["is_stale"] else rt(lang, "h_freshness_ok"),
                      B["orange"] if r["is_stale"] else B["green"]))
        chips_row(s, L, Inches(1.3), items)

        def list_card(x, y, w, h, title, lines, empty):
            # Title + lines are all the card's own text (one shape, no overlay).
            sh = rrect(s, x, y, w, h, B["white"], line=B["line"], radius=0.05)
            paras = [(title, 12, B["navy"], True, PP_ALIGN.LEFT, 5)]
            if lines:
                paras += [(txt, 9.5, color, bold, PP_ALIGN.LEFT, 2) for (txt, color, bold) in lines]
            else:
                paras.append((empty, 9.5, B["muted"], False, PP_ALIGN.LEFT, 0))
            place(sh, paras, anchor=MSO_ANCHOR.TOP, ml=0.18, mt=0.12, mr=0.18, mb=0.1)

        # ----- Row 1: Initiatives | OTD -----
        inits = det.get("initiatives") or []
        ilines = []
        for ini in inits[:3]:
            meta = [x for x in (ini.get("owner"),
                                (f'{rt(lang, "deadline")} {ini["deadline"]}' if ini.get("deadline") else None)) if x]
            tail = f'   ({" · ".join(meta)})' if meta else ""
            ilines.append((f'{ini["title"]}{tail}', B["ink"], False))
        if len(inits) > 3:
            ilines.append((f'+{len(inits) - 3}…', B["muted"], False))
        list_card(L, Inches(1.78), colw, Inches(1.22), rt(lang, "h_initiatives"), ilines, rt(lang, "no_initiative"))

        objs = det.get("objectives", [])
        olines = []
        for o in objs[:3]:
            rag = _status_rag(o["rag"])
            dl = f' · {rt(lang, "deadline")} {o["target_date"]}' if o.get("target_date") else ""
            olines.append((f'●  {o["title"]}   ({_status_label(o["rag"], lang)}{dl})', rgb(_RAG_BRAND[rag]), False))
        if len(objs) > 3:
            olines.append((f'+{len(objs) - 3}…', B["muted"], False))
        list_card(rcolx, Inches(1.78), colw, Inches(1.22),
                  f'{rt(lang, "h_otd_section")} {data["year"]}', olines, rt(lang, "no_obj"))

        # ----- Roadmap card with 4 quarter mini-cards + progress bars -----
        ry, rh = Inches(3.18), Inches(2.4)
        card(s, L, ry, FULLW, rh, f'{rt(lang, "h_roadmap")} {data["year"]}')
        qn, qgap = 4, Inches(0.16)
        inner_x = Emu(int(L) + int(Inches(0.18)))
        inner_w = Emu(int(FULLW) - int(Inches(0.36)))
        qw = Emu((int(inner_w) - (qn - 1) * int(qgap)) // qn)
        qy = Emu(int(ry) + int(Inches(0.52)))
        qh = Emu(int(rh) - int(Inches(0.68)))
        quarters = det.get("quarters", [])
        for i in range(qn):
            qd = quarters[i] if i < len(quarters) else {"q": i + 1, "pct": 0, "items": []}
            qx = Emu(int(inner_x) + i * (int(qw) + int(qgap)))
            qcard = rrect(s, qx, qy, qw, qh, rgb("#F8FAFC"), line=B["line"], radius=0.04)
            tx = Emu(int(qx) + int(Inches(0.1)))
            tw = Emu(int(qw) - int(Inches(0.2)))
            # Q label + percent live in the mini-card's own text frame (two runs).
            qtf = qcard.text_frame
            qtf.word_wrap = True; qtf.vertical_anchor = MSO_ANCHOR.TOP
            qtf.margin_left = Inches(0.1); qtf.margin_right = Inches(0.1)
            qtf.margin_top = Inches(0.07); qtf.margin_bottom = Inches(0.02)
            qp = qtf.paragraphs[0]
            r1 = qp.add_run(); r1.text = f'Q{qd["q"]}'
            r1.font.size = Pt(11); r1.font.bold = True; r1.font.color.rgb = B["navy"]
            r2 = qp.add_run(); r2.text = f'   {qd["pct"]}%'
            r2.font.size = Pt(10); r2.font.color.rgb = B["muted"]
            pbar(s, tx, Emu(int(qy) + int(Inches(0.42))), tw, qd["pct"], B["accent"])
            items = qd.get("items", [])
            lines = []
            for it in items[:5]:
                rag = _status_rag(it["status"])
                lines.append((f'●  {it["title"]}', rgb(_RAG_BRAND[rag]), False))
            if not items:
                lines.append((rt(lang, "no_jalon"), B["muted"], False))
            elif len(items) > 5:
                lines.append((f'+{len(items) - 5}…', B["muted"], False))
            bullets(s, tx, Emu(int(qy) + int(Inches(0.64))), tw, Emu(int(qh) - int(Inches(0.74))), lines, 8.5)

        # ----- Row 3: Key messages | Budget -----
        my, mh = Inches(5.68), Inches(1.5)
        kms = det.get("key_messages") or []
        klines = []
        for m in kms[:3]:
            rag = {"success": "green", "alert": "amber", "risk": "red"}.get(m["kind"], "grey")
            ts = f'   ({m["created_at"]})' if m.get("created_at") else ""
            klines.append((f'●  [{rt(lang, "km_" + m["kind"])}] {m["text"]}{ts}', rgb(_RAG_BRAND[rag]), False))
        if len(kms) > 3:
            klines.append((f'+{len(kms) - 3}…', B["muted"], False))
        list_card(L, my, colw, mh, rt(lang, "h_key_messages"), klines, rt(lang, "no_key_message"))

        # Budget card: title + figures are the card's own text; status stays a chip.
        bsh = rrect(s, rcolx, my, colw, mh, B["white"], line=B["line"], radius=0.05)
        place(bsh, [(rt(lang, "h_budget"), 12, B["navy"], True, PP_ALIGN.LEFT, 6)],
              anchor=MSO_ANCHOR.TOP, ml=0.18, mt=0.12, mr=0.18)
        bud = det.get("budget")
        btf = bsh.text_frame
        if bud is None:
            p = btf.add_paragraph(); rr = p.add_run(); rr.text = rt(lang, "no_budget")
            rr.font.size = Pt(9.5); rr.font.color.rgb = B["muted"]
        else:
            f = lambda v: "-" if v is None else f"{v:,.0f} €"
            st_color = {"on_track": "green", "at_risk": "amber", "over": "red"}[bud["status"]]
            st_lbl = rt(lang, {"on_track": "b_on_track", "at_risk": "b_at_risk", "over": "b_over"}[bud["status"]])
            rows = [
                (rt(lang, "b_total"), f(bud["total"])),
                (rt(lang, "b_spent"), f(bud["spent"]) + (f' · {bud["spent_pct"]}%' if bud.get("spent_pct") is not None else "")),
                (rt(lang, "b_forecast"), f(bud["forecast"]) + (f' · {bud["forecast_pct"]}%' if bud.get("forecast_pct") is not None else "")),
            ]
            for label, val in rows:
                p = btf.add_paragraph(); p.space_after = Pt(3)
                r1 = p.add_run(); r1.text = f'{label} : '
                r1.font.size = Pt(10); r1.font.color.rgb = B["muted"]
                r2 = p.add_run(); r2.text = val
                r2.font.size = Pt(10); r2.font.bold = True; r2.font.color.rgb = B["ink"]
            cw = Inches(0.26 + 0.082 * len(st_lbl))
            chip(s, Emu(int(rcolx) + int(colw) - int(cw) - int(Inches(0.16))), Emu(int(my) + int(Inches(0.12))),
                 st_lbl, rgb(_RAG_BRAND[st_color]))

    # --- Assemble the deck. A single-squad export mirrors the squad page (one
    # focused slide in page order); the multi-squad report keeps summary + grid.
    if data.get("squad_scoped"):
        for r in squads[:_MAX_DETAIL_SLIDES]:
            if r.get("detail"):
                squad_page_slide(r)
    else:
        summary_slide()
        for r in squads[:_MAX_DETAIL_SLIDES]:
            if r.get("detail"):
                detail_slide(r, with_objectives=True)
        # Never silently drop squads the user explicitly selected: if the runaway
        # guard is ever hit, say how many were omitted instead of losing them.
        omitted = len(squads) - _MAX_DETAIL_SLIDES
        if omitted > 0:
            s = new_slide()
            rect(s, Inches(0), Inches(0), prs.slide_width, Inches(0.92), B["navy"])
            textbox(s, margin, Inches(3.2), prs.slide_width - margin * 2, Inches(1),
                    rt(lang, "more_squads", n=omitted), 24, bold=True,
                    color=B["navy"], align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# Roadmap deck palette (mirrors the reference "Global Roadmap" deck).
_RM = {
    "dark": "#304957",   # quarter headers + swimlane labels
    "sub": "#97A3AA",    # month sub-headers
    "card": "#F2F2F2",   # milestone cards
    "card_ink": "#002060",
    "arrow": "#DCE3EA",
    "muted": "#6B7280",
    "white": "#FFFFFF",
}




def render_roadmap_pptx(data: dict) -> bytes:
    """Roadmap swimlane deck (mirrors the reference layout): quarters in columns
    with month sub-headers and a timeline arrow, squads as swimlane rows, and one
    milestone card per (squad, quarter) with status-coloured bullets."""
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE = _pptx_toolkit()

    def rgb(hexstr: str) -> RGBColor:
        return RGBColor.from_string(hexstr.lstrip("#").upper())

    C = {k: rgb(v) for k, v in _RM.items()}
    STAGE = {k: rgb(v) for k, v in STAGE_COLOR.items()}  # EA=gold, GA=green
    lang = _lang(data.get("lang", "fr"))
    year = data["year"]
    gen = data["generated_at"]
    gen_str = gen.strftime("%d/%m/%Y %H:%M") if isinstance(gen, datetime) else str(gen)
    months = _MONTHS[lang]

    SLIDE_W, SLIDE_H = 13.333, 7.5
    prs = pptxtpl.new_presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # Column geometry (4 quarters).
    MARGIN, GAP = 0.5, 0.08
    COL_W = (SLIDE_W - 2 * MARGIN - GAP * 3) / 4
    def col_x(i): return MARGIN + i * (COL_W + GAP)
    Y_Q, H_Q = 0.92, 0.34          # quarter header
    Y_M, H_M = 1.30, 0.28          # month sub-headers
    Y_ARROW, H_ARROW = 1.66, 0.30  # timeline arrow
    Y_TOP, Y_BOTTOM = 2.10, 7.18   # swimlane content band

    squads = [r for blk in data["tribes"] for r in blk["squads"]]

    def shape(s, kind, x, y, w, h, fill, *, line=None, rot=0, round_adj=None):
        sh = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line; sh.line.width = Pt(0.75)
        sh.shadow.inherit = False
        if rot:
            sh.rotation = rot
        if round_adj is not None:
            try:
                sh.adjustments[0] = round_adj
            except Exception:
                pass
        return sh

    def set_text(holder, text, size, color, *, bold=False, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE):
        tf = holder.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.04)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

    def textbox(s, x, y, w, h, runs, size, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]; p.alignment = align
        for txt, color, bold in runs:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        return box

    def draw_header(s):
        textbox(s, MARGIN, 0.22, 8.6, 0.5,
                [(f'{rt(lang, "roadmap_report")} · {data["scope_name"]}', C["dark"], True)], 22)
        textbox(s, MARGIN, 0.66, 8.6, 0.25,
                [(f'{rt(lang, "year")} {year} · {rt(lang, "generated_full", d=gen_str)}', C["muted"], False)], 10.5)
        legend = [("EA  ", STAGE["EA"], True), (rt(lang, "stage_ea") + "      ", C["dark"], False),
                  ("GA  ", STAGE["GA"], True), (rt(lang, "stage_ga"), C["dark"], False)]
        textbox(s, SLIDE_W - 5.9, 0.34, 5.4, 0.3, legend, 10, align=PP_ALIGN.RIGHT)
        # timeline arrow spanning the columns
        shape(s, MSO_SHAPE.RIGHT_ARROW, MARGIN - 0.1, Y_ARROW,
              (col_x(3) + COL_W) - (MARGIN - 0.1) + 0.18, H_ARROW, C["arrow"], round_adj=None)
        for i, q in enumerate((1, 2, 3, 4)):
            x = col_x(i)
            set_text(shape(s, MSO_SHAPE.RECTANGLE, x, Y_Q, COL_W, H_Q, C["dark"]),
                     f'Q{q} {year}', 15, C["white"], bold=True)
            mw = (COL_W - 2 * 0.05) / 3
            for mi in range(3):
                mx = x + mi * (mw + 0.05)
                set_text(shape(s, MSO_SHAPE.RECTANGLE, mx, Y_M, mw, H_M, C["sub"]),
                         months[i * 3 + mi], 10, C["white"])

    def draw_card(s, x, y, w, h, items, fs, line_h):
        card = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, C["card"], round_adj=0.06)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.07); tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.03)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        if not items:
            return
        # Flatten into ordered paragraphs: a bold theme header then its milestones.
        specs: list[tuple[str, object]] = []
        for theme, group in group_by_theme(items):
            if theme:
                specs.append(("theme", theme))
            for it in group:
                specs.append(("item", it))
        item_fs = max(7, fs - 1.5)  # milestone lines a touch smaller than the theme header
        max_lines = max(1, int((h - 0.08) / line_h))
        shown = specs[:max_lines]
        for li, (kind, val) in enumerate(shown):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT  # themes + milestones are left-aligned, never centred
            p.space_after = Pt(0.5)
            if kind == "theme":
                r = p.add_run(); r.text = val
                r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = C["card_ink"]
                continue
            it = val
            r1 = p.add_run(); r1.text = it["title"]
            r1.font.size = Pt(item_fs); r1.font.color.rgb = C["card_ink"]
            stage = it.get("stage")
            if stage:
                ro = p.add_run(); ro.text = " ("
                ro.font.size = Pt(item_fs); ro.font.color.rgb = C["card_ink"]
                rs = p.add_run(); rs.text = stage
                rs.font.size = Pt(item_fs); rs.font.bold = True
                rs.font.color.rgb = STAGE.get(stage, C["card_ink"])
                rc = p.add_run(); rc.text = ")"
                rc.font.size = Pt(item_fs); rc.font.color.rgb = C["card_ink"]
        if len(specs) > max_lines:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = f'+{len(specs) - max_lines}…'
            r.font.size = Pt(max(6, item_fs - 0.5)); r.font.color.rgb = C["muted"]

    def draw_swimlanes(s, lanes):
        # Everything fits on ONE slide: band height + fonts scale with the count,
        # but with a readable floor (small decks get a comfortably large font).
        n = max(1, len(lanes))
        band_h = (Y_BOTTOM - Y_TOP) / n
        card_fs = (13 if band_h >= 1.05 else 12 if band_h >= 0.85 else 11 if band_h >= 0.68
                   else 10 if band_h >= 0.52 else 9 if band_h >= 0.40 else 8)
        line_h = card_fs * 0.020
        label_fs = max(8.5, min(13, band_h * 12))
        lbl_h = min(0.36, band_h * 0.6)
        for ri, sq in enumerate(lanes):
            by = Y_TOP + ri * band_h
            bcy = by + band_h / 2
            lbl_len = max(0.35, min(band_h - 0.08, 1.7))
            lbl = shape(s, MSO_SHAPE.RECTANGLE, 0.24 - lbl_len / 2, bcy - lbl_h / 2,
                        lbl_len, lbl_h, C["dark"], rot=270)
            set_text(lbl, sq["name"], label_fs, C["white"], bold=True)
            qmap = {qd["q"]: qd["items"] for qd in (sq.get("detail") or {}).get("quarters", [])}
            for i, q in enumerate((1, 2, 3, 4)):
                draw_card(s, col_x(i), by + 0.04, COL_W, band_h - 0.08, qmap.get(q, []), card_fs, line_h)

    s = pptxtpl.add_slide(prs)  # single page, always
    draw_header(s)
    draw_swimlanes(s, squads)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_initiatives_pptx(data: dict, *, lang: str = "fr") -> bytes:
    """Render the initiatives list as a single branded table slide (max 18 rows)."""
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE = _pptx_toolkit()

    def rgb(h):
        return RGBColor.from_string(h.lstrip("#").upper())

    B = {k: rgb(v) for k, v in _BRAND.items()}
    lang = _lang(lang)
    T = _INIT_T[lang]
    prs = pptxtpl.new_presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    margin = Inches(0.5)
    s = pptxtpl.add_slide(prs)
    head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.92))
    head.fill.solid(); head.fill.fore_color.rgb = B["navy"]; head.line.fill.background(); head.shadow.inherit = False
    tf = head.text_frame; tf.margin_left = Inches(0.5); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; r = p.add_run(); r.text = f'{T["title"]} · {data["scope_name"]}'
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = B["white"]

    headers = [T["h_init"], T["h_owner"], T["h_squad"], T["h_deadline"]]
    wfrac = [0.40, 0.24, 0.22, 0.14]
    items = data["items"][:18]
    nrows = max(2, len(items) + 1)
    table_w = int(prs.slide_width - margin * 2)
    tbl = s.shapes.add_table(nrows, 4, margin, Inches(1.2), Emu(table_w),
                             Inches(0.34) + Inches(0.3) * (nrows - 1)).table
    for ci, f in enumerate(wfrac):
        tbl.columns[ci].width = Emu(int(table_w * f))

    def cell(c, text, size, color, *, bold=False, align=PP_ALIGN.LEFT, fill=None):
        c.margin_left = Inches(0.06); c.margin_right = Inches(0.04); c.margin_top = Emu(0); c.margin_bottom = Emu(0)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fill is not None:
            c.fill.solid(); c.fill.fore_color.rgb = fill
        else:
            c.fill.background()
        c.text = text or " "
        pp = c.text_frame.paragraphs[0]; pp.alignment = align
        if pp.runs:
            rr = pp.runs[0]; rr.font.size = Pt(size); rr.font.bold = bold; rr.font.color.rgb = color

    for ci, h in enumerate(headers):
        cell(tbl.cell(0, ci), h, 11, B["white"], bold=True, fill=B["navy"])
    if not items:
        cell(tbl.cell(1, 0), T["none"], 10, B["muted"])
        for ci in range(1, 4):
            cell(tbl.cell(1, ci), " ", 10, B["muted"])
    for ri, it in enumerate(items, start=1):
        zebra = B["zebra"] if ri % 2 == 0 else B["white"]
        cells = [(it["title"], B["ink"], True), (it["owner"] or "-", B["ink"], False),
                 (it["squad_name"] or "-", B["ink"], False), (it["deadline"], B["ink"], False)]
        for ci, (val, color, bold) in enumerate(cells):
            cell(tbl.cell(ri, ci), val, 10, color, bold=bold, fill=zebra)
    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()


def render_dependencies_pptx(data: dict) -> bytes:
    """Paginated table deck of milestone dependencies, grouped by the entity waited
    on. Rows flow across slides so no dependency is ever dropped."""
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE = _pptx_toolkit()

    def rgb(hexstr: str) -> RGBColor:
        return RGBColor.from_string(hexstr.lstrip("#").upper())

    B = {k: rgb(v) for k, v in _BRAND.items()}
    lang = _lang(data.get("lang", "fr"))
    T = _DEP_T[lang]

    prs = pptxtpl.new_presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW = prs.slide_width
    margin = int(Inches(0.5))
    content_w = int(SW) - 2 * margin
    fracs = [0.405, 0.225, 0.075, 0.18, 0.115]
    aligns = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER]
    colw = [int(content_w * f) for f in fracs]
    colx, acc = [], margin
    for w in colw:
        colx.append(acc); acc += w
    ROW_H, HDR_H, GRP_H = int(Inches(0.30)), int(Inches(0.32)), int(Inches(0.46))
    TOP0, BOTTOM = int(Inches(1.25)), int(Inches(7.22))

    gen = data["generated_at"]
    gen_str = gen.strftime("%d/%m/%Y %H:%M") if isinstance(gen, datetime) else str(gen)

    def new_slide():
        return pptxtpl.add_slide(prs)

    def textbox(s, left, top, width, height, text, size, *, bold=False, color=None,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
        box = s.shapes.add_textbox(Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height)))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color if color is not None else B["ink"]
        return box

    def rect(s, left, top, width, height, fill):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height)))
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        sh.line.fill.background(); sh.shadow.inherit = False
        return sh

    def band(s, cont):
        rect(s, 0, 0, SW, Inches(1.05), B["navy"])
        textbox(s, margin, Inches(0.14), Inches(9.2), Inches(0.5),
                T["title"] + (T["suite"] if cont else ""), 22, bold=True, color=B["white"], anchor=MSO_ANCHOR.TOP)
        textbox(s, margin, Inches(0.64), Inches(9.2), Inches(0.3),
                f'{data["scope_name"]} · {rt(lang, "year")} {data["year"]} · {T["total"].format(n=data["total"])}',
                12, color=rgb("#C7D2FE"), anchor=MSO_ANCHOR.TOP)
        textbox(s, int(SW) - int(Inches(4.3)), Inches(0.32), Inches(3.8), Inches(0.4),
                rt(lang, "generated_full", d=gen_str), 10, color=rgb("#C7D2FE"), align=PP_ALIGN.RIGHT)

    def col_headers(s, y):
        rect(s, margin, y, content_w, HDR_H, B["navy"])
        for i, key in enumerate(("c_jalon", "c_squad", "c_trim", "c_owner", "c_status")):
            textbox(s, colx[i], y, colw[i], HDR_H, T[key], 9.5, bold=True, color=B["white"], align=aligns[i])
        return y + HDR_H

    def group_header(s, y, g, cont):
        rect(s, margin, y, content_w, GRP_H, rgb("#EEF2FF"))
        tlbl = {"tribe": T["t_tribe"], "squad": T["t_squad"], "text": T["t_text"]}.get(g["target_type"], "")
        if g["target_type"] == "squad" and g.get("target_tribe"):
            tlbl = f'{tlbl} · {g["target_tribe"]}'
        txt = f'▶  {g["target_label"]}   ({tlbl})' + (T["suite"] if cont else "")
        textbox(s, colx[0], y, content_w - int(Inches(2.4)), GRP_H, txt, 13, bold=True, color=B["navy"])
        textbox(s, margin + content_w - int(Inches(2.4)), y, Inches(2.3), GRP_H,
                T["gcount"].format(n=len(g["items"])), 10, color=B["muted"], align=PP_ALIGN.RIGHT)
        return y + GRP_H

    def trunc(x, n):
        x = x or ""
        return x if len(x) <= n else x[:n - 1] + "…"

    def data_row(s, y, it, zebra):
        if zebra:
            rect(s, margin, y, content_w, ROW_H, rgb("#F8FAFC"))
        vals = [trunc(it["jalon"], 62),
                trunc(f'{it["squad_name"]} · {it["tribe_name"]}', 34),
                f'Q{it["quarter"]} {str(it["year"])[2:]}',
                trunc(it["owner"] or "-", 24),
                _status_label(it["status"], lang)]
        colors = [B["ink"], B["muted"], B["ink"], B["ink"], rgb(_RAG_BRAND[_status_rag(it["status"])])]
        bolds = [True, False, False, False, True]
        for i, v in enumerate(vals):
            textbox(s, colx[i], y, colw[i], ROW_H, v, 9, bold=bolds[i], color=colors[i], align=aligns[i])
        return y + ROW_H

    if data["total"] == 0:
        s = new_slide(); band(s, False)
        textbox(s, margin, Inches(3.2), content_w, Inches(0.6), T["none"], 18, bold=True,
                color=B["muted"], align=PP_ALIGN.CENTER)
        buf = io.BytesIO(); prs.save(buf); return buf.getvalue()

    state = {"s": None, "y": 0}

    def open_slide(cont):
        s = new_slide(); band(s, cont)
        state["s"] = s; state["y"] = TOP0

    open_slide(False)
    for g in data["groups"]:
        if state["y"] + GRP_H + HDR_H + ROW_H > BOTTOM:
            open_slide(True)
        state["y"] = group_header(state["s"], state["y"], g, False)
        state["y"] = col_headers(state["s"], state["y"])
        for idx, it in enumerate(g["items"]):
            if state["y"] + ROW_H > BOTTOM:
                open_slide(True)
                state["y"] = group_header(state["s"], state["y"], g, True)
                state["y"] = col_headers(state["s"], state["y"])
            state["y"] = data_row(state["s"], state["y"], it, idx % 2 == 1)

    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()
